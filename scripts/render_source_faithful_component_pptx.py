#!/usr/bin/env python3
"""Render source-faithful component crops into a review PPTX.

The component SVGs retain the editable source vectors. This companion deck is
for visual review: it crops PowerPoint's rendering of the source PPTX at the
same planned component bounds, so font availability in an SVG renderer cannot
change what the reviewer sees.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
KIT_ROOT = ROOT / "templates" / "components" / "source_templates" / "nsfc_defense_distilled_kit"
DEFAULT_CATALOG = KIT_ROOT / "source_faithful_component_catalog.json"
DEFAULT_SOURCE_RENDERS = ROOT / "projects" / "nsfc_source_component_extraction_20260728" / "powerpoint_source_png"
DEFAULT_OUTPUT = ROOT / "projects" / "nsfc_source_component_extraction_20260728" / "nsfc_defense_source_faithful_components.pptx"
DEFAULT_CROP_DIR = ROOT / "projects" / "nsfc_source_component_extraction_20260728" / "powerpoint_component_crops"


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8-sig"))
    if not isinstance(payload, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return payload


def _slide_index(value: str) -> int:
    match = re.search(r"(\d+)$", value)
    if not match:
        raise ValueError(f"could not determine source slide index from {value}")
    return int(match.group(1))


def _source_pngs(render_dir: Path) -> dict[int, Path]:
    found: dict[int, Path] = {}
    for path in render_dir.glob("*.PNG"):
        match = re.search(r"(\d+)$", path.stem)
        if match:
            found[int(match.group(1))] = path
    if not found:
        raise FileNotFoundError(f"no PowerPoint PNG exports found in {render_dir}")
    return found


def _crop_source_component(*, source: Path, bounds: dict[str, Any], destination: Path) -> tuple[int, int]:
    with Image.open(source) as image:
        width, height = image.size
        scale_x = width / 1280.0
        scale_y = height / 720.0
        left = round(float(bounds["x"]) * scale_x)
        top = round(float(bounds["y"]) * scale_y)
        right = round((float(bounds["x"]) + float(bounds["width"])) * scale_x)
        bottom = round((float(bounds["y"]) + float(bounds["height"])) * scale_y)
        if left < 0 or top < 0 or right > width or bottom > height or right <= left or bottom <= top:
            raise ValueError(f"invalid crop bounds {bounds} for {source.name} ({width} x {height})")
        destination.parent.mkdir(parents=True, exist_ok=True)
        image.crop((left, top, right, bottom)).save(destination, "PNG")
        return right - left, bottom - top


def _set_text(shape: Any, text: str, *, size: float, color: str, bold: bool = False, align: Any = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _add_text(slide: Any, *, x: float, y: float, width: float, height: float, text: str, size: float, color: str, bold: bool = False, align: Any = None) -> None:
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    _set_text(shape, text, size=size, color=color, bold=bold, align=align)


def _add_rect(slide: Any, *, x: float, y: float, width: float, height: float, fill: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.fill.background()
    shape.line.width = Pt(0)


def _fit(width: int, height: int, *, max_width: float, max_height: float) -> tuple[float, float]:
    ratio = width / height
    if max_width / max_height > ratio:
        fitted_height = max_height
        return fitted_height * ratio, fitted_height
    fitted_width = max_width
    return fitted_width, fitted_width / ratio


def build_review_pptx(
    *,
    catalog_path: Path = DEFAULT_CATALOG,
    source_render_dir: Path = DEFAULT_SOURCE_RENDERS,
    crop_dir: Path = DEFAULT_CROP_DIR,
    output_path: Path = DEFAULT_OUTPUT,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    catalog = _read_json(catalog_path)
    components = [entry for entry in catalog.get("components", []) if isinstance(entry, dict)]
    if not components:
        raise ValueError("source-faithful component catalog has no components")
    renders = _source_pngs(source_render_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    _add_rect(cover, x=0, y=0, width=13.333333, height=7.5, fill="FFFFFF")
    _add_rect(cover, x=0, y=0, width=13.333333, height=0.86, fill="751497")
    _add_text(cover, x=0.55, y=1.45, width=12.2, height=0.65, text="国自然答辩：原样内容组件", size=30, color="1C1B22", bold=True)
    _add_text(cover, x=0.57, y=2.18, width=11.9, height=0.40, text=f"{len(components)} 个组件 · 按源坐标原样裁切", size=14, color="5F5864")

    records: list[dict[str, Any]] = []
    for component_number, component in enumerate(components, 1):
        component_id = str(component["component_id"])
        source_index = _slide_index(str(component["source_slide"]))
        source_png = renders.get(source_index)
        if source_png is None:
            raise FileNotFoundError(f"missing PowerPoint export for source slide {source_index}")
        crop_path = crop_dir / f"{component_number:02d}_{component_id}.png"
        crop_width, crop_height = _crop_source_component(
            source=source_png,
            bounds=dict(component["bounds"]),
            destination=crop_path,
        )

        slide = prs.slides.add_slide(blank)
        _add_rect(slide, x=0, y=0, width=13.333333, height=7.5, fill="FFFFFF")
        _add_rect(slide, x=0, y=0, width=13.333333, height=0.62, fill="751497")
        _add_text(slide, x=0.42, y=0.11, width=9.7, height=0.30, text=str(component["display_name"]), size=18, color="FFFFFF", bold=True)
        _add_text(slide, x=10.2, y=0.14, width=2.7, height=0.22, text=f"源第 {source_index} 页", size=10, color="FFFFFF", align=PP_ALIGN.RIGHT)
        _add_text(slide, x=0.52, y=6.78, width=10.8, height=0.25, text="原样裁切", size=9.5, color="645B69")
        _add_text(slide, x=11.52, y=6.78, width=1.25, height=0.25, text=f"{component_number}/{len(components)}", size=9.5, color="751497", align=PP_ALIGN.RIGHT)
        fitted_width, fitted_height = _fit(crop_width, crop_height, max_width=12.0, max_height=5.72)
        x = (13.333333 - fitted_width) / 2
        y = 0.87 + (5.72 - fitted_height) / 2
        slide.shapes.add_picture(str(crop_path), Inches(x), Inches(y), width=Inches(fitted_width), height=Inches(fitted_height))
        records.append(
            {
                "component_id": component_id,
                "source_slide": source_index,
                "crop_path": crop_path.relative_to(ROOT).as_posix(),
                "crop_pixels": {"width": crop_width, "height": crop_height},
            }
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    report = {
        "schema_version": "easyslides.source_faithful_component_pptx_review.v1",
        "status": "pass",
        "component_count": len(records),
        "slide_count": len(prs.slides),
        "source_renderer": "Microsoft PowerPoint Export",
        "pptx": output_path.resolve().relative_to(ROOT).as_posix(),
        "components": records,
    }
    report_path = output_path.with_suffix(".json")
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Render source-faithful components into a PowerPoint review deck.")
    parser.add_argument("--catalog", type=Path, default=DEFAULT_CATALOG)
    parser.add_argument("--source-render-dir", type=Path, default=DEFAULT_SOURCE_RENDERS)
    parser.add_argument("--crop-dir", type=Path, default=DEFAULT_CROP_DIR)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    report = build_review_pptx(
        catalog_path=args.catalog,
        source_render_dir=args.source_render_dir,
        crop_dir=args.crop_dir,
        output_path=args.out,
    )
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"Source-faithful review PPTX: pass ({report['component_count']} components)")
        print(args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
