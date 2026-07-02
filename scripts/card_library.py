"""Card library registry, capacity validation, and preview export."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import sys
from typing import Any


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_LIBRARY = REPO_ROOT / "templates" / "cards" / "card_library.json"
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from scripts.text_capacity import SlotCapacity, fit_text_to_capacity


def load_card_library(path: str | Path | None = None) -> dict[str, Any]:
    library_path = Path(path) if path else DEFAULT_LIBRARY
    with library_path.open("r", encoding="utf-8") as handle:
        library = json.load(handle)
    styles = library.get("styles")
    if not isinstance(styles, list) or not styles:
        raise ValueError("card library must define a non-empty styles list")
    return library


def card_styles(library: dict[str, Any] | None = None) -> list[dict[str, Any]]:
    return list((library or load_card_library())["styles"])


def count_card_styles(library: dict[str, Any] | None = None) -> int:
    return len(card_styles(library))


def get_card_style(card_id: str, library: dict[str, Any] | None = None) -> dict[str, Any]:
    for style in card_styles(library):
        if style.get("card_id") == card_id:
            return style
    raise KeyError(f"unknown card_id: {card_id}")


def select_cards(
    content_shape: str | None = None,
    item_count: int | None = None,
    density: str | None = None,
    library: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    """Return matching card styles, highest score first."""
    matches: list[tuple[int, str, dict[str, Any]]] = []
    for style in card_styles(library):
        selection = style.get("selection", {})
        score = 0
        shapes = set(selection.get("content_shapes") or [])
        if content_shape:
            if content_shape in shapes:
                score += 6
            else:
                continue
        if item_count is not None:
            minimum = int(selection.get("item_count_min", 1))
            maximum = int(selection.get("item_count_max", minimum))
            if minimum <= item_count <= maximum:
                score += 4
                if minimum == maximum == item_count:
                    score += 1
            else:
                continue
        if density:
            if selection.get("density") == density:
                score += 2
            else:
                score -= 1
        matches.append((score, str(style["card_id"]), style))
    matches.sort(key=lambda item: (-item[0], item[1]))
    return [style for _, _, style in matches]


def _slot_capacity(slot: dict[str, Any]) -> SlotCapacity:
    max_lines = int(slot["max_lines"])
    chars = int(slot["max_chars_per_line_zh"])
    return SlotCapacity(
        slot_id=str(slot["slot_id"]),
        role=str(slot.get("role") or "body"),
        font_size_px=float(slot["font_size_px"]),
        min_font_size_px=float(slot["min_font_size_px"]),
        line_height=float(slot["line_height"]),
        max_chars_per_line_zh=chars,
        max_lines=max_lines,
        capacity_chars=max_lines * chars,
        overflow_action=str(slot["overflow_action"]),
    )


def _check_text(slot: dict[str, Any], value: Any, location: str) -> dict[str, Any] | None:
    text = "" if value is None else str(value).strip()
    result = fit_text_to_capacity(text, _slot_capacity(slot))
    if result.input_over_capacity or result.output_overflow:
        return {
            "location": location,
            "slot_id": slot["slot_id"],
            "input_chars": result.input_chars,
            "capacity_chars": result.rendered_chars if result.output_overflow else _slot_capacity(slot).capacity_chars,
            "max_lines": slot["max_lines"],
            "max_chars_per_line_zh": slot["max_chars_per_line_zh"],
            "overflow_action": slot["overflow_action"],
        }
    return None


def validate_card_payload(
    card_id: str,
    payload: dict[str, Any],
    library: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Validate a card payload against slot capacities and item-count limits."""
    style = get_card_style(card_id, library)
    selection = style.get("selection", {})
    items = payload.get("items")
    item_count = len(items) if isinstance(items, list) else 1
    minimum = int(selection.get("item_count_min", 1))
    maximum = int(selection.get("item_count_max", minimum))
    violations: list[dict[str, Any]] = []

    if not (minimum <= item_count <= maximum):
        violations.append(
            {
                "location": "items",
                "slot_id": "item_count",
                "input_count": item_count,
                "allowed_min": minimum,
                "allowed_max": maximum,
                "overflow_action": "choose_matching_card_or_split",
            }
        )

    checked = 0
    for slot in style.get("slots", []):
        slot_id = str(slot["slot_id"])
        required = bool(slot.get("required", False))
        checked_here = False

        if slot_id in payload:
            checked += 1
            checked_here = True
            violation = _check_text(slot, payload[slot_id], slot_id)
            if violation:
                violations.append(violation)

        if isinstance(items, list):
            for index, item in enumerate(items, start=1):
                if not isinstance(item, dict):
                    violations.append(
                        {
                            "location": f"items[{index}]",
                            "slot_id": slot_id,
                            "overflow_action": "use_object_items",
                        }
                    )
                    continue
                if slot_id in item:
                    checked += 1
                    checked_here = True
                    violation = _check_text(slot, item[slot_id], f"items[{index}].{slot_id}")
                    if violation:
                        violations.append(violation)
                elif required and slot_id not in payload:
                    violations.append(
                        {
                            "location": f"items[{index}]",
                            "slot_id": slot_id,
                            "missing": True,
                            "overflow_action": "fill_required_slot",
                        }
                    )

        if required and not checked_here:
            violations.append(
                {
                    "location": "payload",
                    "slot_id": slot_id,
                    "missing": True,
                    "overflow_action": "fill_required_slot",
                }
            )

    return {
        "passed": not violations,
        "card_id": card_id,
        "checked_slots": checked,
        "violations": violations,
    }


def _px(value: float) -> float:
    return value / 96.0


def export_preview_pptx(
    output: str | Path,
    library: dict[str, Any] | None = None,
) -> Path:
    """Export a polished PPTX catalog for all card styles."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    library = library or load_card_library()
    output_path = Path(output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    colors = {
        "metric": (0, 85, 135),
        "parallel": (0, 118, 168),
        "comparison": (220, 38, 38),
        "process": (245, 158, 11),
        "evidence": (13, 148, 136),
        "method": (12, 39, 68),
        "literature": (71, 85, 105),
        "callout": (26, 26, 46),
    }
    accent_cycle = [
        (220, 38, 38),
        (245, 158, 11),
        (0, 118, 168),
        (74, 144, 164),
        (100, 116, 139),
    ]
    family_labels = {
        "metric": "METRIC",
        "parallel": "PARALLEL",
        "comparison": "COMPARE",
        "process": "PROCESS",
        "evidence": "EVIDENCE",
        "method": "METHOD",
        "literature": "LIT NOTE",
        "callout": "CALLOUT",
    }
    family_codes = {
        "metric": "ME",
        "parallel": "PA",
        "comparison": "CO",
        "process": "PR",
        "evidence": "EV",
        "method": "MD",
        "literature": "LT",
        "callout": "CA",
    }

    def add_textbox(slide, x, y, w, h, text, size=16, bold=False, color=(35, 35, 35), align=PP_ALIGN.LEFT):
        if len(str(text)) > 24 and h < 36:
            h = 36
        box = slide.shapes.add_textbox(Inches(_px(x)), Inches(_px(y)), Inches(_px(w)), Inches(_px(h)))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        para = frame.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return box

    def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(_px(x)),
            Inches(_px(y)),
            Inches(_px(w)),
            Inches(_px(h)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        if line is None:
            shape.line.color.rgb = RGBColor(*fill)
        else:
            shape.line.color.rgb = RGBColor(*line)
        shape.line.width = Pt(1)
        return shape

    def add_line(slide, x1, y1, x2, y2, color=(226, 232, 240), width=1):
        line = slide.shapes.add_connector(
            1,
            Inches(_px(x1)),
            Inches(_px(y1)),
            Inches(_px(x2)),
            Inches(_px(y2)),
        )
        line.line.color.rgb = RGBColor(*color)
        line.line.width = Pt(width)
        return line

    def add_circle(slide, cx, cy, r, fill, line=None):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(_px(cx - r)),
            Inches(_px(cy - r)),
            Inches(_px(r * 2)),
            Inches(_px(r * 2)),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*fill)
        circle.line.color.rgb = RGBColor(*(line or fill))
        circle.line.width = Pt(1)
        return circle

    def _payload_text(payload, *slot_ids):
        for slot_id in slot_ids:
            value = payload.get(slot_id)
            if value:
                return str(value)
        return ""

    def _slot_lookup(slots):
        return {slot["slot_id"]: slot for slot in slots}

    def _slot_line(slot):
        if not slot:
            return ""
        return f"{slot['max_lines']} lines x {slot['max_chars_per_line_zh']} chars"

    def draw_card(slide, box, family, payload, slots, ordinal=None):
        accent = accent_cycle[(ordinal - 1) % len(accent_cycle)] if ordinal else colors.get(family, (55, 55, 55))
        x, y, w, h = box["x"], box["y"], box["width"], box["height"]

        add_rect(slide, x + 6, y + 8, w, h, (219, 225, 232), (219, 225, 232))
        add_rect(slide, x, y, w, h, (255, 255, 255), (226, 232, 240))
        add_rect(slide, x, y, w, 8, accent, accent, radius=False)

        compact = h < 250 or w < 280
        slot_map = _slot_lookup(slots)
        label = family_labels.get(family, "CARD")
        if ordinal is not None:
            add_textbox(slide, x + 16, y + 22, 48, 30, f"{ordinal:02d}", size=16, bold=True, color=accent)
        else:
            add_textbox(slide, x + 16, y + 22, 48, 30, family_codes.get(family, "CD"), size=13, bold=True, color=accent)

        if not compact:
            add_textbox(slide, x + w - 112, y + 24, 90, 22, label, size=8, bold=True, color=accent, align=PP_ALIGN.RIGHT)
            add_circle(slide, x + w / 2, y + 72, 26, (242, 247, 250), (226, 232, 240))
            add_textbox(slide, x + w / 2 - 14, y + 62, 28, 22, label[:1], size=13, bold=True, color=accent, align=PP_ALIGN.CENTER)

        metric = _payload_text(payload, "metric")
        title = _payload_text(payload, "title", "label", "takeaway", "claim", "statement", "citation", "date")
        body = _payload_text(payload, "body", "note", "bullets", "evidence", "problem", "method", "result", "limitation")
        caption = _payload_text(payload, "source", "synthesis")

        left = x + 22
        content_w = w - 44
        if metric:
            metric_size = 30 if compact else 40
            add_textbox(slide, left, y + 48, content_w, 50, metric, size=metric_size, bold=True, color=accent, align=PP_ALIGN.CENTER)
            add_textbox(slide, left, y + 98, content_w, 28, title, size=14 if compact else 16, bold=True, color=(30, 41, 59), align=PP_ALIGN.CENTER)
            if body:
                add_line(slide, left, y + 132, x + w - 22, y + 132)
                add_textbox(slide, left, y + 144, content_w, max(34, h - 170), body, size=11 if compact else 13, color=(71, 85, 105), align=PP_ALIGN.CENTER)
        else:
            title_y = y + (52 if compact else 120)
            title_h = 34 if compact else 58
            title_size = 14 if compact else min(20, max(14, w / 20))
            add_textbox(slide, left, title_y, content_w, title_h, title, size=title_size, bold=True, color=(30, 41, 59), align=PP_ALIGN.CENTER if compact else PP_ALIGN.LEFT)
            divider_y = title_y + title_h + 8
            add_line(slide, left, divider_y, x + w - 22, divider_y)
            body_y = divider_y + 14
            bottom_guard = 56 if h > 300 else 32
            add_textbox(slide, left, body_y, content_w, max(30, y + h - body_y - bottom_guard), body, size=11 if compact else 14, color=(71, 85, 105))
            if caption and not compact:
                add_textbox(slide, left, y + h - 84, content_w, 32, caption, size=11, color=(100, 116, 139))

        capacity_bits = []
        for slot_id in ("body", "bullets", "evidence", "note", "statement"):
            line = _slot_line(slot_map.get(slot_id))
            if line:
                capacity_bits.append(line)
                break
        pill_text = capacity_bits[0] if capacity_bits else "capacity locked"
        if h >= 220:
            pill_w = min(w - 44, 156)
            add_rect(slide, left, y + h - 38, pill_w, 22, (241, 245, 249), (226, 232, 240))
            add_textbox(slide, left + 10, y + h - 34, pill_w - 20, 16, pill_text, size=8, bold=True, color=(100, 116, 139), align=PP_ALIGN.CENTER)

    def add_header(slide, title, subtitle):
        add_rect(slide, 0, 0, 1280, 8, (0, 85, 135), (0, 85, 135), radius=False)
        add_textbox(slide, 42, 30, 1010, 34, title, size=19, bold=True, color=(15, 23, 42))
        add_rect(slide, 42, 70, 1196, 42, (24, 52, 94), (24, 52, 94))
        add_rect(slide, 42, 70, 5, 42, (56, 189, 248), (56, 189, 248), radius=False)
        add_textbox(slide, 58, 82, 1128, 18, subtitle, size=10, bold=True, color=(248, 250, 252))
        add_textbox(slide, 42, 126, 450, 18, "FIXED GEOMETRY  |  DECLARED CAPACITY  |  PPT-MASTER STYLE SKIN", size=8, color=(148, 163, 184))

    title_slide = prs.slides.add_slide(blank)
    add_rect(title_slide, 0, 0, 1280, 8, (0, 85, 135), (0, 85, 135), radius=False)
    add_textbox(title_slide, 72, 72, 880, 48, "EasySlides Card Library", size=28, bold=True, color=(15, 23, 42))
    add_textbox(
        title_slide,
        74,
        128,
        980,
        54,
        f"{count_card_styles(library)} polished card styles with fixed geometry, slot capacity, and overflow rules.",
        size=16,
        color=(71, 85, 105),
    )
    add_rect(title_slide, 74, 198, 1130, 46, (24, 52, 94), (24, 52, 94))
    add_rect(title_slide, 74, 198, 6, 46, (56, 189, 248), (56, 189, 248), radius=False)
    add_textbox(title_slide, 94, 212, 950, 18, "A native PowerPoint component catalog: choose by content shape, then validate every slot before rendering.", size=11, bold=True, color=(248, 250, 252))
    overview_y = 286
    for index, style in enumerate(card_styles(library), start=1):
        col = (index - 1) % 4
        row = (index - 1) // 4
        x = 74 + col * 292
        y = overview_y + row * 96
        family = style.get("family", "parallel")
        accent = colors.get(family, (55, 55, 55))
        add_rect(title_slide, x + 4, y + 5, 248, 68, (219, 225, 232), (219, 225, 232))
        add_rect(title_slide, x, y, 248, 68, (255, 255, 255), (226, 232, 240))
        add_rect(title_slide, x, y, 248, 6, accent, accent, radius=False)
        add_textbox(title_slide, x + 16, y + 16, 34, 22, f"{index:02d}", size=12, bold=True, color=accent)
        add_textbox(title_slide, x + 58, y + 15, 174, 20, style["name_zh"], size=11, bold=True, color=(30, 41, 59))
        add_textbox(title_slide, x + 58, y + 38, 174, 16, style["card_id"], size=8, color=(100, 116, 139))

    for index, style in enumerate(card_styles(library), start=1):
        slide = prs.slides.add_slide(blank)
        selection = style.get("selection", {})
        add_header(
            slide,
            f"{index:02d}. {style['name_zh']} / {style['card_id']}",
            f"family={style.get('family')} | density={selection.get('density')} | items={selection.get('item_count_min')}-{selection.get('item_count_max')}",
        )
        payload = style.get("preview_payload", {})
        boxes = style.get("layout", {}).get("boxes", [])
        slots = style.get("slots", [])
        items = payload.get("items") if isinstance(payload, dict) else None
        if isinstance(items, list):
            for item_index, box in enumerate(boxes[: len(items)], start=1):
                item = items[item_index - 1]
                if isinstance(item, dict):
                    draw_card(slide, box, style.get("family", ""), item, slots, item_index)
        else:
            for box in boxes:
                if box.get("id") == "figure":
                    figure = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(_px(box["x"])),
                        Inches(_px(box["y"])),
                        Inches(_px(box["width"])),
                        Inches(_px(box["height"])),
                    )
                    figure.fill.solid()
                    figure.fill.fore_color.rgb = RGBColor(235, 239, 243)
                    figure.line.color.rgb = RGBColor(190, 198, 207)
                    figure.text = "FIGURE AREA"
                    figure.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    figure.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                else:
                    draw_card(slide, box, style.get("family", ""), payload, slots)
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"Best for: {selection.get('best_for', '')}\nAvoid when: {selection.get('avoid_when', '')}"

    prs.save(output_path)
    return output_path


def _json_print(value: Any) -> None:
    print(json.dumps(value, ensure_ascii=False, indent=2))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Query, validate, and preview EasySlides card styles.")
    parser.add_argument("--library", default=str(DEFAULT_LIBRARY), help="Path to card_library.json.")
    subparsers = parser.add_subparsers(dest="command", required=True)

    subparsers.add_parser("count", help="Print the number of card styles.")

    list_parser = subparsers.add_parser("list", help="List card styles.")
    list_parser.add_argument("--json", action="store_true", help="Emit JSON.")

    query_parser = subparsers.add_parser("query", help="Find suitable card styles.")
    query_parser.add_argument("--content-shape", default=None)
    query_parser.add_argument("--item-count", type=int, default=None)
    query_parser.add_argument("--density", default=None)
    query_parser.add_argument("--json", action="store_true", help="Emit JSON.")

    validate_parser = subparsers.add_parser("validate", help="Validate a card payload.")
    validate_parser.add_argument("--card-id", required=True)
    payload_group = validate_parser.add_mutually_exclusive_group(required=True)
    payload_group.add_argument("--payload-json")
    payload_group.add_argument("--payload-file")

    preview_parser = subparsers.add_parser("preview", help="Export a PPTX preview catalog.")
    preview_parser.add_argument("--output", required=True)

    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    library = load_card_library(args.library)

    if args.command == "count":
        print(count_card_styles(library))
        return 0

    if args.command == "list":
        rows = [
            {
                "card_id": style["card_id"],
                "name_zh": style["name_zh"],
                "family": style["family"],
                "density": style["selection"]["density"],
            }
            for style in card_styles(library)
        ]
        if args.json:
            _json_print(rows)
        else:
            for row in rows:
                print(f"{row['card_id']}\t{row['name_zh']}\t{row['family']}\t{row['density']}")
        return 0

    if args.command == "query":
        matches = select_cards(args.content_shape, args.item_count, args.density, library)
        rows = [
            {
                "card_id": style["card_id"],
                "name_zh": style["name_zh"],
                "family": style["family"],
                "best_for": style["selection"]["best_for"],
            }
            for style in matches
        ]
        if args.json:
            _json_print(rows)
        else:
            for row in rows:
                print(f"{row['card_id']}\t{row['name_zh']}\t{row['best_for']}")
        return 0 if rows else 1

    if args.command == "validate":
        if args.payload_file:
            with Path(args.payload_file).open("r", encoding="utf-8") as handle:
                payload = json.load(handle)
        else:
            payload = json.loads(args.payload_json)
        result = validate_card_payload(args.card_id, payload, library)
        _json_print(result)
        return 0 if result["passed"] else 1

    if args.command == "preview":
        output = export_preview_pptx(args.output, library)
        print(output)
        return 0

    raise AssertionError(f"unhandled command {args.command!r}")


if __name__ == "__main__":
    raise SystemExit(main())
