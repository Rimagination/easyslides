#!/usr/bin/env python3
"""Validate exported PPTX text layout before delivery.

This is an exit-gate validator for issues that are easy to miss in SVG or
template-level checks: text overflow, stacked text boxes, off-slide text, and
font-size collapse. It intentionally uses conservative estimates so a warning
can be inspected before the deck ships.
"""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation

try:
    from scripts import layout_metrics
except ImportError:  # pragma: no cover - direct script execution
    import layout_metrics

try:  # Pillow is listed in requirements, but keep the CLI useful without it.
    from PIL import ImageFont

    HAS_PIL = True
except Exception:  # pragma: no cover - depends on local environment
    ImageFont = None
    HAS_PIL = False


EMU_PER_INCH = layout_metrics.EMU_PER_INCH
PT_PER_INCH = layout_metrics.PT_PER_INCH
SCHEMA_VERSION = "easyslides.pptx_text_layout_report.v1"
MIN_READABLE_FONT_PT = 8.0
OVERFLOW_BLOCK_RATIO = 1.10
OVERLAP_BLOCK_RATIO = 0.65
OVERLAP_MIN_AREA_IN2 = 0.08
OFF_SLIDE_TOLERANCE_IN = 0.08
AffineMatrix = tuple[float, float, float, float, float, float]
IDENTITY_MATRIX: AffineMatrix = (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)
PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


@dataclass(frozen=True)
class TextBox:
    slide_number: int
    shape_index: int
    name: str
    text: str
    x: float
    y: float
    w: float
    h: float
    usable_w: float
    usable_h: float
    font_name: str
    font_size_pt: float
    line_spacing: float
    wrap: str = "square"

    @property
    def right(self) -> float:
        return self.x + self.w

    @property
    def bottom(self) -> float:
        return self.y + self.h

    @property
    def area(self) -> float:
        return max(0.0, self.w) * max(0.0, self.h)


def _inches(value: int | float | None) -> float:
    return layout_metrics.emu_to_in(value)


def _matrix_multiply(left: AffineMatrix, right: AffineMatrix) -> AffineMatrix:
    return layout_metrics.matrix_multiply(left, right)


def _rotate_matrix(angle_deg: float, cx: float, cy: float) -> AffineMatrix:
    return layout_metrics.rotate_matrix(angle_deg, cx, cy)


def _transform_point(matrix: AffineMatrix, x: float, y: float) -> tuple[float, float]:
    return layout_metrics.transform_point(matrix, x, y)


def _transform_rect(
    matrix: AffineMatrix,
    x: float,
    y: float,
    w: float,
    h: float,
) -> tuple[float, float, float, float]:
    box = layout_metrics.transform_box(layout_metrics.Box(x, y, w, h), matrix)
    return box.x, box.y, box.width, box.height


def _xfrm_rect(xfrm: ET.Element | None) -> tuple[float, float, float, float] | None:
    off = xfrm.find("a:off", PPTX_NS) if xfrm is not None else None
    ext = xfrm.find("a:ext", PPTX_NS) if xfrm is not None else None
    if off is None or ext is None:
        return None
    return (
        _inches(off.attrib.get("x")),
        _inches(off.attrib.get("y")),
        _inches(ext.attrib.get("cx")),
        _inches(ext.attrib.get("cy")),
    )


def _xfrm_rotation_matrix(xfrm: ET.Element | None, rect: tuple[float, float, float, float]) -> AffineMatrix:
    if xfrm is None:
        return IDENTITY_MATRIX
    raw = xfrm.attrib.get("rot")
    if not raw:
        return IDENTITY_MATRIX
    try:
        angle_deg = float(raw) / 60000.0
    except ValueError:
        return IDENTITY_MATRIX
    if not angle_deg:
        return IDENTITY_MATRIX
    x, y, w, h = rect
    return _rotate_matrix(angle_deg, x + w / 2, y + h / 2)


def _group_child_matrix(grp: ET.Element) -> AffineMatrix:
    xfrm = grp.find("p:grpSpPr/a:xfrm", PPTX_NS)
    outer = _xfrm_rect(xfrm)
    ch_off = xfrm.find("a:chOff", PPTX_NS) if xfrm is not None else None
    ch_ext = xfrm.find("a:chExt", PPTX_NS) if xfrm is not None else None
    if outer is None or ch_off is None or ch_ext is None:
        return IDENTITY_MATRIX
    x, y, w, h = outer
    child_x = _inches(ch_off.attrib.get("x"))
    child_y = _inches(ch_off.attrib.get("y"))
    child_w = _inches(ch_ext.attrib.get("cx"))
    child_h = _inches(ch_ext.attrib.get("cy"))
    sx = w / child_w if child_w else 1.0
    sy = h / child_h if child_h else 1.0
    scale_translate: AffineMatrix = (sx, 0.0, 0.0, sy, x - child_x * sx, y - child_y * sy)
    return _matrix_multiply(_xfrm_rotation_matrix(xfrm, outer), scale_translate)


def _first_run_font(shape) -> tuple[str, float]:
    font_name = "Arial"
    font_size = 12.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                font_name = run.font.name
            if run.font.size is not None:
                font_size = float(run.font.size.pt)
            if run.text:
                return font_name, font_size
    return font_name, font_size


def _line_spacing(shape, font_size: float) -> float:
    for para in shape.text_frame.paragraphs:
        value = para.line_spacing
        if isinstance(value, float):
            return max(value, 1.0)
        if value is not None:
            try:
                return max(float(value.pt) / max(font_size, 1.0), 1.0)
            except Exception:
                pass
    return 1.2


def _usable_dimensions(shape) -> tuple[float, float]:
    tf = shape.text_frame
    margin_left = _inches(tf.margin_left) if tf.margin_left is not None else 0.05
    margin_right = _inches(tf.margin_right) if tf.margin_right is not None else 0.05
    margin_top = _inches(tf.margin_top) if tf.margin_top is not None else 0.05
    margin_bottom = _inches(tf.margin_bottom) if tf.margin_bottom is not None else 0.05
    usable_w = max(0.01, _inches(shape.width) - margin_left - margin_right)
    usable_h = max(0.01, _inches(shape.height) - margin_top - margin_bottom)
    return usable_w, usable_h


def _iter_text_boxes(prs: Presentation) -> list[TextBox]:
    boxes: list[TextBox] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            font_name, font_size = _first_run_font(shape)
            usable_w, usable_h = _usable_dimensions(shape)
            boxes.append(
                TextBox(
                    slide_number=slide_number,
                    shape_index=shape_index,
                    name=getattr(shape, "name", ""),
                    text=text,
                    x=_inches(shape.left),
                    y=_inches(shape.top),
                    w=_inches(shape.width),
                    h=_inches(shape.height),
                    usable_w=usable_w,
                    usable_h=usable_h,
                    font_name=font_name,
                    font_size_pt=font_size,
                    line_spacing=_line_spacing(shape, font_size),
                    wrap="none" if shape.text_frame.word_wrap is False else "square",
                )
            )
    return boxes


def _pptx_slide_order(pptx_path: Path) -> list[str]:
    with zipfile.ZipFile(pptx_path) as zf:
        pres = ET.fromstring(zf.read("ppt/presentation.xml"))
        rels = ET.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
        relmap = {
            rel.attrib["Id"]: rel.attrib["Target"]
            for rel in rels.findall("rel:Relationship", PPTX_NS)
            if rel.attrib.get("Id") and rel.attrib.get("Target")
        }
        ordered: list[str] = []
        for slide_id in pres.findall(".//p:sldId", PPTX_NS):
            rid = slide_id.attrib.get(f"{{{PPTX_NS['r']}}}id")
            target = relmap.get(rid or "")
            if target:
                ordered.append("ppt/" + target.lstrip("../"))
        return ordered


def _raw_text(sp: ET.Element) -> str:
    paragraphs: list[str] = []
    for paragraph in sp.findall(".//a:p", PPTX_NS):
        chunks = [node.text or "" for node in paragraph.findall(".//a:t", PPTX_NS)]
        text = "".join(chunks).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs).strip()


def _raw_font(sp: ET.Element) -> tuple[str, float]:
    font_name = "Arial"
    font_size = 12.0
    rpr = sp.find(".//a:rPr", PPTX_NS)
    if rpr is not None:
        sz = rpr.attrib.get("sz")
        if sz:
            try:
                font_size = float(sz) / 100.0
            except ValueError:
                pass
        latin = rpr.find("a:latin", PPTX_NS)
        if latin is not None and latin.attrib.get("typeface"):
            font_name = latin.attrib["typeface"]
    return font_name, font_size


def _raw_body_margins(sp: ET.Element) -> tuple[float, float, float, float]:
    body_pr = sp.find(".//a:bodyPr", PPTX_NS)
    if body_pr is None:
        return 0.05, 0.05, 0.05, 0.05

    def margin(name: str) -> float:
        value = body_pr.attrib.get(name)
        if value is None:
            return 0.05
        try:
            return max(0.0, layout_metrics.emu_to_in(float(value)))
        except (TypeError, ValueError):
            return 0.05

    return margin("lIns"), margin("rIns"), margin("tIns"), margin("bIns")


def _raw_alignment(sp: ET.Element) -> str:
    ppr = sp.find(".//a:pPr", PPTX_NS)
    return ppr.attrib.get("algn", "l") if ppr is not None else "l"


def _raw_body_anchor(sp: ET.Element) -> str:
    body_pr = sp.find(".//a:bodyPr", PPTX_NS)
    return str(body_pr.attrib.get("anchor") or "t").lower() if body_pr is not None else "t"


def _raw_body_wrap(sp: ET.Element) -> str:
    body_pr = sp.find(".//a:bodyPr", PPTX_NS)
    if body_pr is None:
        return "square"
    return str(body_pr.attrib.get("wrap") or "square").lower()


def _measure_text_width_pt(text: str, font_name: str, size_pt: float) -> float:
    font = _load_font(font_name, size_pt)
    return max(
        (
            sum(_char_width(ch, font, size_pt) for ch in line)
            for line in (text.splitlines() or [text])
        ),
        default=0.0,
    )


def _raw_visual_box(
    sp: ET.Element,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_name: str,
    font_size: float,
) -> tuple[float, float, float, float]:
    visual_w = max(0.01, _measure_text_width_pt(text, font_name, font_size) / PT_PER_INCH)
    line_count = max(1, len(text.splitlines()))
    visual_h = max(0.01, line_count * font_size * 1.2 / PT_PER_INCH)
    alignment = _raw_alignment(sp)
    if alignment == "ctr":
        visual_x = x + (w - visual_w) / 2
    elif alignment == "r":
        visual_x = x + w - visual_w
    else:
        visual_x = x
    anchor = _raw_body_anchor(sp)
    if anchor in {"ctr", "middle", "center"}:
        visual_y = y + (h - min(visual_h, h)) / 2
    elif anchor in {"b", "bottom"}:
        visual_y = y + h - min(visual_h, h)
    else:
        visual_y = y
    return visual_x, visual_y, visual_w, min(max(visual_h, 0.01), max(h, visual_h))


def _iter_text_boxes_from_xml(pptx_path: str | Path) -> list[TextBox]:
    pptx_path = Path(pptx_path)
    boxes: list[TextBox] = []
    with zipfile.ZipFile(pptx_path) as zf:
        for slide_number, slide_name in enumerate(_pptx_slide_order(pptx_path), start=1):
            root = ET.fromstring(zf.read(slide_name))
            shape_index = 0

            def walk(parent: ET.Element, matrix: AffineMatrix) -> None:
                nonlocal shape_index
                for child in parent:
                    tag = child.tag.rsplit("}", 1)[-1] if "}" in child.tag else child.tag
                    if tag == "grpSp":
                        walk(child, _matrix_multiply(matrix, _group_child_matrix(child)))
                        continue
                    if tag != "sp":
                        walk(child, matrix)
                        continue
                    shape_index += 1
                    sp = child
                    text = _raw_text(sp)
                    if not text:
                        continue
                    xfrm = sp.find("p:spPr/a:xfrm", PPTX_NS)
                    rect = _xfrm_rect(xfrm)
                    if rect is None:
                        continue
                    x, y, w, h = rect
                    margin_left, margin_right, margin_top, margin_bottom = _raw_body_margins(sp)
                    font_name, font_size = _raw_font(sp)
                    visual_x, visual_y, visual_w, visual_h = _raw_visual_box(
                        sp, text, x, y, w, h, font_name, font_size
                    )
                    visual_x, visual_y, visual_w, visual_h = _transform_rect(
                        _xfrm_rotation_matrix(xfrm, rect),
                        visual_x,
                        visual_y,
                        visual_w,
                        visual_h,
                    )
                    visual_x, visual_y, visual_w, visual_h = _transform_rect(
                        matrix,
                        visual_x,
                        visual_y,
                        visual_w,
                        visual_h,
                    )
                    name = ""
                    c_nv_pr = sp.find("p:nvSpPr/p:cNvPr", PPTX_NS)
                    if c_nv_pr is not None:
                        name = c_nv_pr.attrib.get("name", "")
                    boxes.append(
                        TextBox(
                            slide_number=slide_number,
                            shape_index=shape_index,
                            name=name or f"raw_xml_shape_{shape_index}",
                            text=text,
                            x=visual_x,
                            y=visual_y,
                            w=visual_w,
                            h=visual_h,
                            usable_w=max(0.01, w - margin_left - margin_right),
                            usable_h=max(0.01, h - margin_top - margin_bottom),
                            font_name=font_name,
                            font_size_pt=font_size,
                            line_spacing=1.2,
                            wrap=_raw_body_wrap(sp),
                        )
                    )

            walk(root, IDENTITY_MATRIX)
    return boxes


def _shape_name(shape: ET.Element, fallback: str) -> str:
    c_nv_pr = shape.find(".//p:cNvPr", PPTX_NS)
    if c_nv_pr is not None and c_nv_pr.attrib.get("name"):
        return c_nv_pr.attrib["name"]
    return fallback


def _negative_extent_issues_from_xml(pptx_path: str | Path) -> list[dict[str, Any]]:
    pptx_path = Path(pptx_path)
    issues: list[dict[str, Any]] = []
    shape_tags = {"sp", "pic", "grpSp", "graphicFrame", "cxnSp"}
    with zipfile.ZipFile(pptx_path) as zf:
        for slide_number, slide_name in enumerate(_pptx_slide_order(pptx_path), start=1):
            root = ET.fromstring(zf.read(slide_name))
            shape_index = 0

            def walk(parent: ET.Element) -> None:
                nonlocal shape_index
                for child in parent:
                    tag = child.tag.rsplit("}", 1)[-1] if "}" in child.tag else child.tag
                    if tag in shape_tags:
                        shape_index += 1
                        xfrm = child.find(".//a:xfrm", PPTX_NS)
                        ext = xfrm.find("a:ext", PPTX_NS) if xfrm is not None else None
                        if ext is not None:
                            try:
                                cx = int(float(ext.attrib.get("cx", "0")))
                                cy = int(float(ext.attrib.get("cy", "0")))
                            except ValueError:
                                cx = cy = 0
                            if cx < 0 or cy < 0:
                                issues.append(
                                    {
                                        "code": "PPTX-INVALID-NEGATIVE-EXTENT",
                                        "severity": "blocking",
                                        "slide_number": slide_number,
                                        "shape_index": shape_index,
                                        "shape_name": _shape_name(child, f"raw_xml_shape_{shape_index}"),
                                        "message": "DrawingML shape has a negative a:ext cx/cy; PowerPoint may repair or corrupt the slide.",
                                        "details": {
                                            "slide_part": slide_name,
                                            "cx_emu": cx,
                                            "cy_emu": cy,
                                            "cx_in": round(_inches(cx), 3),
                                            "cy_in": round(_inches(cy), 3),
                                        },
                                    }
                                )
                    walk(child)

            walk(root)
    return issues


def _font_candidates(font_name: str) -> list[Path]:
    normalized = font_name.lower().replace(" ", "")
    names = [
        f"{normalized}.ttf",
        f"{normalized}.otf",
        "arial.ttf",
        "calibri.ttf",
        "aptos.ttf",
        "segoeui.ttf",
    ]
    roots = [
        Path("C:/Windows/Fonts"),
        Path("/System/Library/Fonts"),
        Path("/Library/Fonts"),
        Path("/usr/share/fonts/truetype"),
        Path("/usr/share/fonts"),
    ]
    return [root / name for root in roots for name in names]


def _load_font(font_name: str, size_pt: float):
    if not HAS_PIL:
        return None
    for path in _font_candidates(font_name):
        if path.exists():
            try:
                return ImageFont.truetype(str(path), max(1, int(round(size_pt))))
            except Exception:
                continue
    try:  # pragma: no cover - only used when no local font path is found
        return ImageFont.load_default()
    except Exception:
        return None


def _char_width(ch: str, font, size_pt: float) -> float:
    if font is not None:
        if hasattr(font, "getlength"):
            return float(font.getlength(ch))
        bbox = font.getbbox(ch)
        return float(bbox[2] - bbox[0]) if bbox else 0.0
    size_px = layout_metrics.pt_to_px(size_pt)
    return layout_metrics.px_to_pt(layout_metrics.estimate_text_width_px(ch, size_px))


def _measure_lines(text: str, font_name: str, size_pt: float, width_pt: float) -> int:
    if width_pt <= 0:
        return math.inf
    font = _load_font(font_name, size_pt)
    lines = 0
    for paragraph in text.splitlines() or [""]:
        current = 0.0
        saw_char = False
        for ch in paragraph:
            char_w = _char_width(ch, font, size_pt)
            if saw_char and current + char_w > width_pt:
                lines += 1
                current = 0.0
            current += char_w
            saw_char = True
        lines += 1
    return max(lines, 1)


def _explicit_line_count(text: str) -> int:
    return max(1, len(text.splitlines()) or 1)


def _measure_lines_for_box(box: TextBox, width_pt: float) -> int:
    if box.wrap == "none":
        return _explicit_line_count(box.text)
    return _measure_lines(box.text, box.font_name, box.font_size_pt, width_pt)


def _issue(code: str, message: str, box: TextBox, severity: str = "blocking", **details: Any) -> dict[str, Any]:
    return {
        "code": code,
        "severity": severity,
        "slide_number": box.slide_number,
        "shape_index": box.shape_index,
        "shape_name": box.name,
        "message": message,
        "details": {
            "text_preview": box.text[:120],
            "x": round(box.x, 3),
            "y": round(box.y, 3),
            "w": round(box.w, 3),
            "h": round(box.h, 3),
            "font_name": box.font_name,
            "font_size_pt": round(box.font_size_pt, 2),
            "wrap": box.wrap,
            **details,
        },
    }


def _overlap_area(a: TextBox, b: TextBox) -> float:
    dx = min(a.right, b.right) - max(a.x, b.x)
    dy = min(a.bottom, b.bottom) - max(a.y, b.y)
    if dx <= 0 or dy <= 0:
        return 0.0
    return dx * dy


def _is_page_number_like(box: TextBox) -> bool:
    text = box.text.strip()
    return len(text) <= 4 and text.replace("/", "").replace("-", "").isdigit()


def _is_short_label_like(box: TextBox) -> bool:
    text = box.text.strip().replace("\n", "")
    if not text or len(text) > 6:
        return False
    if any(ch.isspace() for ch in text):
        return False
    return True


def _is_vertical_stack_label(box: TextBox) -> bool:
    lines = [line.strip() for line in box.text.splitlines() if line.strip()]
    return 2 <= len(lines) <= 6 and all(len(line) == 1 for line in lines)


def _is_template_placeholder_text(text: str) -> bool:
    stripped = re.sub(r"\s+", "", text)
    if not stripped:
        return False
    if "请添加" in stripped:
        return True
    return "添加" in stripped and any(
        token in stripped for token in ("标题", "描述", "正文", "章节", "图片名称")
    )


def _is_side_label_bleed(box: TextBox, slide_w: float, slide_h: float) -> bool:
    return (
        box.x < 0
        and box.x >= -0.45
        and box.y >= -OFF_SLIDE_TOLERANCE_IN
        and box.bottom <= slide_h + OFF_SLIDE_TOLERANCE_IN
        and box.h >= slide_h * 0.60
        and box.w <= slide_w * 0.16
    )


def validate_pptx_text_layout(pptx_path: str | Path) -> dict[str, Any]:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    slide_w = _inches(prs.slide_width)
    slide_h = _inches(prs.slide_height)
    boxes = _iter_text_boxes_from_xml(pptx_path)
    text_box_source = "raw_xml"
    if not boxes:
        boxes = _iter_text_boxes(prs)
        text_box_source = "python-pptx_fallback"
    issues: list[dict[str, Any]] = _negative_extent_issues_from_xml(pptx_path)

    for box in boxes:
        width_pt = box.usable_w * PT_PER_INCH
        height_in = max(box.usable_h, box.h) if box.wrap == "none" else box.usable_h
        height_pt = height_in * PT_PER_INCH
        lines_needed = _measure_lines_for_box(box, width_pt)
        line_height = box.font_size_pt * box.line_spacing
        lines_available = max(1, math.floor(height_pt / max(line_height, 1.0) + 1e-6))
        overflow_ratio = lines_needed / max(lines_available, 1)

        if (
            overflow_ratio > OVERFLOW_BLOCK_RATIO
            and not _is_short_label_like(box)
            and not _is_vertical_stack_label(box)
        ):
            placeholder = _is_template_placeholder_text(box.text)
            issues.append(
                _issue(
                    "TEXT-OVERFLOW",
                    f"Text needs {lines_needed} line(s), but the box fits {lines_available}.",
                    box,
                    severity="warning" if placeholder else "blocking",
                    lines_needed=lines_needed,
                    lines_available=lines_available,
                    overflow_ratio=round(overflow_ratio, 2),
                    measured_with_pil=HAS_PIL,
                )
            )

        if box.font_size_pt < MIN_READABLE_FONT_PT and not _is_page_number_like(box):
            issues.append(
                _issue(
                    "TEXT-FONT-TOO-SMALL",
                    f"Font size {box.font_size_pt:.1f}pt is below the readable floor.",
                    box,
                    min_readable_font_pt=MIN_READABLE_FONT_PT,
                )
            )

        if (
            box.x < -OFF_SLIDE_TOLERANCE_IN
            or box.y < -OFF_SLIDE_TOLERANCE_IN
            or box.right > slide_w + OFF_SLIDE_TOLERANCE_IN
            or box.bottom > slide_h + OFF_SLIDE_TOLERANCE_IN
        ) and not _is_side_label_bleed(box, slide_w, slide_h):
            placeholder = _is_template_placeholder_text(box.text)
            issues.append(
                _issue(
                    "TEXT-OFF-SLIDE",
                    "Text box extends outside the slide bounds.",
                    box,
                    severity="warning" if placeholder else "blocking",
                    slide_width=round(slide_w, 3),
                    slide_height=round(slide_h, 3),
                )
            )

        if len(box.text) > 24 and (box.h < 0.35 or box.usable_w < 1.0):
            issues.append(
                _issue(
                    "TEXT-LABEL-TOO-LONG",
                    "A small label-like slot contains sentence-length text.",
                    box,
                    severity="warning",
                    text_length=len(box.text),
                )
            )

    for i, left in enumerate(boxes):
        for right in boxes[i + 1 :]:
            if left.slide_number != right.slide_number:
                continue
            area = _overlap_area(left, right)
            if area < OVERLAP_MIN_AREA_IN2:
                continue
            ratio = area / max(min(left.area, right.area), 0.01)
            if ratio > OVERLAP_BLOCK_RATIO:
                placeholder_overlap = _is_template_placeholder_text(left.text) and _is_template_placeholder_text(right.text)
                issues.append(
                    {
                        "code": "TEXT-OVERLAP",
                        "severity": "warning" if placeholder_overlap else "blocking",
                        "slide_number": left.slide_number,
                        "shape_index": left.shape_index,
                        "shape_name": left.name,
                        "message": "Two text boxes strongly overlap; this usually means stacked text.",
                        "details": {
                            "other_shape_index": right.shape_index,
                            "other_shape_name": right.name,
                            "overlap_ratio": round(ratio, 2),
                            "overlap_area_in2": round(area, 3),
                            "text_preview": left.text[:80],
                            "other_text_preview": right.text[:80],
                        },
                    }
                )

    blocking_count = sum(1 for issue in issues if issue["severity"] == "blocking")
    return {
        "schema_version": SCHEMA_VERSION,
        "pptx_path": str(pptx_path),
        "status": "pass" if blocking_count == 0 else "fail",
        "slide_count": len(prs.slides),
        "text_box_count": len(boxes),
        "text_box_source": text_box_source,
        "blocking_count": blocking_count,
        "warning_count": sum(1 for issue in issues if issue["severity"] == "warning"),
        "issues": issues,
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", help="PPTX file to validate.")
    parser.add_argument("--report", help="Optional JSON report path.")
    args = parser.parse_args(argv)

    report = validate_pptx_text_layout(args.pptx)
    payload = json.dumps(report, ensure_ascii=False, indent=2)
    if args.report:
        Path(args.report).write_text(payload + "\n", encoding="utf-8")
    print(payload)
    return 0 if report["status"] == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
