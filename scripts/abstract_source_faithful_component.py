#!/usr/bin/env python3
"""Create editable placeholder candidates from source-faithful component assets.

Only declared payloads may change. Every fixed visual attribute is verified
against the original source-faithful SVG before the derived asset is written.
"""

from __future__ import annotations

import argparse
import base64
import copy
import hashlib
import json
from pathlib import Path
import sys
import xml.etree.ElementTree as ET
from typing import Any


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

KIT_ROOT = ROOT / "templates" / "components" / "source_templates" / "nsfc_defense_distilled_kit"
SOURCE_COMPONENT = KIT_ROOT / "components" / "source_faithful" / "image_footer_card.svg"
OUTPUT_COMPONENT = KIT_ROOT / "components" / "abstracted" / "image_footer_card_placeholder.svg"
OUTPUT_CATALOG = KIT_ROOT / "source_derived_component_catalog.json"
DEFAULT_OUTPUT = ROOT / "projects" / "nsfc_source_component_extraction_20260728" / "image_footer_card_placeholder_review.pptx"
DEFAULT_PREVIEW = ROOT / "projects" / "nsfc_source_component_extraction_20260728" / "image_footer_card_placeholder.png"

SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"
ET.register_namespace("", SVG_NS)
ET.register_namespace("xlink", XLINK_NS)


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _sha256(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


def _visual_signature(root: ET.Element) -> str:
    """Hash fixed SVG structure while excluding declared payload values."""
    rows: list[dict[str, Any]] = []
    for element in root.iter():
        excluded = {"href", f"{{{XLINK_NS}}}href"}
        if _local_name(element.tag) in {"text", "tspan"}:
            excluded.update({"x", "text-anchor"})
        attributes = {
            key: value
            for key, value in element.attrib.items()
            if key not in excluded and not key.startswith("data-")
        }
        rows.append({"tag": _local_name(element.tag), "attributes": sorted(attributes.items())})
    return _sha256(json.dumps(rows, ensure_ascii=False, separators=(",", ":")).encode("utf-8"))


def _slot_format_signature(root: ET.Element, tag_name: str) -> str:
    rows: list[list[tuple[str, str]]] = []
    for element in root.iter():
        if _local_name(element.tag) != tag_name:
            continue
        rows.append(
            sorted(
                (key, value)
                for key, value in element.attrib.items()
                if key not in {"href", f"{{{XLINK_NS}}}href"} and not key.startswith("data-slot")
            )
        )
    return _sha256(json.dumps(rows, ensure_ascii=False, separators=(",", ":")).encode("utf-8"))


def _caption_style_signature(root: ET.Element) -> str:
    """Preserve caption typography while allowing its generic alignment contract."""
    rows: list[list[tuple[str, str]]] = []
    for element in root.iter():
        if _local_name(element.tag) not in {"text", "tspan"}:
            continue
        rows.append(
            sorted(
                (key, value)
                for key, value in element.attrib.items()
                if key not in {"x", "text-anchor"} and not key.startswith("data-")
            )
        )
    return _sha256(json.dumps(rows, ensure_ascii=False, separators=(",", ":")).encode("utf-8"))


def _placeholder_image_data_uri() -> str:
    placeholder = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100" preserveAspectRatio="none">
  <rect width="100" height="100" fill="#F4F1F5"/>
  <text x="50" y="52" text-anchor="middle" font-family="Arial, sans-serif" font-size="9" fill="#8A8190">IMAGE</text>
</svg>'''
    encoded = base64.b64encode(placeholder.encode("utf-8")).decode("ascii")
    return f"data:image/svg+xml;base64,{encoded}"


def _replace_caption(root: ET.Element) -> float:
    text_nodes = [node for node in root.iter() if _local_name(node.tag) == "text"]
    if len(text_nodes) != 1:
        raise ValueError(f"image_footer_card must have exactly one caption text node; found {len(text_nodes)}")
    text_node = text_nodes[0]
    text_node.attrib["data-slot-id"] = "CAPTION"
    text_node.attrib["data-slot-kind"] = "text"
    caption_box = next(
        (
            node
            for node in root.iter()
            if _local_name(node.tag) == "rect" and node.attrib.get("fill") == "none"
        ),
        None,
    )
    if caption_box is None:
        raise ValueError("image_footer_card caption geometry box is missing")
    center_x = float(caption_box.attrib["x"]) + float(caption_box.attrib["width"]) / 2
    text_node.attrib["x"] = f"{center_x:.2f}"
    text_node.attrib["text-anchor"] = "middle"
    text_node.text = None
    leaves = [node for node in text_node.iter() if _local_name(node.tag) == "tspan"]
    if not leaves:
        text_node.text = "CAPTION"
        return center_x
    for leaf in leaves:
        leaf.text = ""
    leaves[0].text = "CAPTION"
    return center_x


def abstract_image_footer_card(*, output_component: Path = OUTPUT_COMPONENT, catalog_path: Path = OUTPUT_CATALOG) -> dict[str, Any]:
    source_root = ET.fromstring(SOURCE_COMPONENT.read_text(encoding="utf-8"))
    source_visual_signature = _visual_signature(source_root)
    source_image_format = _slot_format_signature(source_root, "image")
    source_caption_style = _caption_style_signature(source_root)

    root = copy.deepcopy(source_root)
    root.attrib["data-component-id"] = "image_footer_card_placeholder"
    root.attrib["data-easyslides-source-fidelity"] = "source_derived_declared_payloads_only"
    images = [node for node in root.iter() if _local_name(node.tag) == "image"]
    if len(images) != 1:
        raise ValueError(f"image_footer_card must have exactly one image node; found {len(images)}")
    image = images[0]
    image.attrib["href"] = _placeholder_image_data_uri()
    image.attrib["data-slot-id"] = "IMAGE"
    image.attrib["data-slot-kind"] = "image"
    caption_center_x = _replace_caption(root)

    if _visual_signature(root) != source_visual_signature:
        raise ValueError("fixed visual structure changed while abstracting image_footer_card")
    if _slot_format_signature(root, "image") != source_image_format:
        raise ValueError("image slot geometry or formatting changed while abstracting image_footer_card")
    if _caption_style_signature(root) != source_caption_style:
        raise ValueError("caption typography changed while abstracting image_footer_card")

    output_component.parent.mkdir(parents=True, exist_ok=True)
    output_component.write_text(ET.tostring(root, encoding="unicode"), encoding="utf-8")
    catalog = {
        "schema_version": "easyslides.source_derived_component_catalog.v1",
        "template_id": "nsfc_defense_distilled",
        "status": "pass",
        "component_count": 1,
        "components": [
            {
                "asset_id": "source_derived/nsfc_defense_distilled/image_footer_card_placeholder",
                "component_id": "image_footer_card_placeholder",
                "display_name": "图像脚注卡（占位版）",
                "description": "保留原图卡的版式、底栏、配色与文字格式；图片和底栏文字改为声明式槽位。",
                "asset_path": output_component.relative_to(ROOT).as_posix(),
                "source_component": SOURCE_COMPONENT.relative_to(ROOT).as_posix(),
                "asset_status": "source_derived_editable_candidate",
                "slots": [
                    {"slot_id": "IMAGE", "kind": "image", "required": True, "source_node": "image"},
                    {"slot_id": "CAPTION", "kind": "text", "required": True, "source_node": "text"},
                ],
                "allowed_mutations": ["replace IMAGE source only", "replace CAPTION text only"],
                "forbidden_mutations": ["change color", "change font", "change font size", "change geometry", "change border", "change crop", "change layer order"],
                "fidelity": {
                    "fixed_visual_signature": source_visual_signature,
                    "image_format_signature": source_image_format,
                    "caption_style_signature": source_caption_style,
                    "caption_alignment": {"horizontal": "center", "center_x": caption_center_x},
                    "fixed_visual_mutation_count": 0,
                },
            }
        ],
    }
    catalog_path.write_text(json.dumps(catalog, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return catalog


def _add_text(slide: Any, *, x: float, y: float, width: float, height: float, text: str, size: float, color: str, bold: bool = False, align: Any = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def render_review_pptx(*, component_path: Path = OUTPUT_COMPONENT, output_path: Path = DEFAULT_OUTPUT, preview_path: Path = DEFAULT_PREVIEW) -> Path:
    import cairosvg
    from PIL import Image
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    preview_path.parent.mkdir(parents=True, exist_ok=True)
    cairosvg.svg2png(url=str(component_path), write_to=str(preview_path), output_width=1512)
    with Image.open(preview_path) as image:
        crop_width, crop_height = image.size
    max_width, max_height = 10.8, 5.55
    ratio = crop_width / crop_height
    fitted_width = min(max_width, max_height * ratio)
    fitted_height = fitted_width / ratio

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333333), Inches(0.62))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor.from_string("751497")
    header.line.width = Pt(0)
    _add_text(slide, x=0.42, y=0.11, width=9.7, height=0.30, text="图像脚注卡（占位版）", size=18, color="FFFFFF", bold=True)
    _add_text(slide, x=0.55, y=6.76, width=11.5, height=0.26, text="可替换：IMAGE / CAPTION", size=10, color="645B69")
    slide.shapes.add_picture(str(preview_path), Inches((13.333333 - fitted_width) / 2), Inches(0.88 + (5.55 - fitted_height) / 2), width=Inches(fitted_width), height=Inches(fitted_height))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Abstract the source-faithful image footer card into placeholders.")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--no-review-pptx", action="store_true")
    args = parser.parse_args(argv)
    catalog = abstract_image_footer_card()
    if not args.no_review_pptx:
        render_review_pptx(output_path=args.out)
    print(json.dumps(catalog, ensure_ascii=False, indent=2))
    if not args.no_review_pptx:
        print(args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
