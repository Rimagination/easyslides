#!/usr/bin/env python3
"""Native PPTX template fill helper.

Analyze a source deck as a slide library, scaffold/validate a fill plan, then
clone, reorder, reuse, and patch native PPTX slide parts without going through
the SVG pipeline.
"""

from __future__ import annotations

import argparse
import json
import posixpath
import re
import sys
from io import BytesIO
from datetime import datetime
from pathlib import Path
from typing import Any
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation


SCHEMA = "template_fill_pptx_plan.v1"
PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CONTENT_TYPE_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
SLIDE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
CHART_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"
PACKAGE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
NOTES_SLIDE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
SLIDE_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
CHART_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.drawingml.chart+xml"
XLSX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
REL_TYPE_BASE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
SHARED_REL_TYPES = {
    REL_TYPE_BASE + name
    for name in (
        "slideLayout",
        "slideMaster",
        "notesMaster",
        "theme",
        "presProps",
        "viewProps",
        "tableStyles",
    )
}
STAGE_OWNED_REL_TYPES = {CHART_REL_TYPE, NOTES_SLIDE_REL_TYPE, SLIDE_REL_TYPE}

ET.register_namespace("", PACKAGE_REL_NS)
ET.register_namespace("p", PRESENTATION_NS)
ET.register_namespace("r", REL_NS)
ET.register_namespace("c", C_NS)


def _shape_text(shape: object) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(
        paragraph.text
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text.strip()
    ).strip()


def _shape_geometry(shape: object) -> dict[str, int]:
    return {
        "left": int(getattr(shape, "left", 0) or 0),
        "top": int(getattr(shape, "top", 0) or 0),
        "width": int(getattr(shape, "width", 0) or 0),
        "height": int(getattr(shape, "height", 0) or 0),
    }


def _slot_role(index: int, text: str, shape: object) -> str:
    name = str(getattr(shape, "name", "") or "").lower()
    if index == 0 or "title" in name:
        return "title"
    if len(text) <= 40:
        return "label"
    return "body"


def _page_type(slide_index: int, slide_count: int, slots: list[dict[str, Any]]) -> str:
    if slide_index == 1:
        return "cover"
    if slide_index == slide_count:
        return "ending"
    if len(slots) <= 2:
        return "chapter"
    return "content"


def _qn(namespace: str, tag: str) -> str:
    return f"{{{namespace}}}{tag}"


def _zip_entries(pptx_path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(pptx_path, "r") as zf:
        return {info.filename: zf.read(info.filename) for info in zf.infolist() if not info.is_dir()}


def _normalize_part(target: str, base: str = "ppt/presentation.xml") -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(base), target)).lstrip("/")


def _relative_target(owner_part: str, target_part: str) -> str:
    return posixpath.relpath(target_part, posixpath.dirname(owner_part))


def _rels_name_for_part(part_name: str) -> str:
    parent = posixpath.dirname(part_name)
    basename = posixpath.basename(part_name)
    return posixpath.join(parent, "_rels", f"{basename}.rels")


def _relationships(entries: dict[str, bytes], rels_name: str) -> dict[str, dict[str, str]]:
    if rels_name not in entries:
        return {}
    root = ET.fromstring(entries[rels_name])
    rels: dict[str, dict[str, str]] = {}
    for rel in root.findall(_qn(PACKAGE_REL_NS, "Relationship")):
        rel_id = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        rel_type = rel.attrib.get("Type")
        if rel_id and target and rel_type:
            rels[rel_id] = {"target": target, "type": rel_type}
    return rels


def _slide_refs(entries: dict[str, bytes]) -> dict[int, tuple[str, str]]:
    pres_root = ET.fromstring(entries["ppt/presentation.xml"])
    rels = _relationships(entries, "ppt/_rels/presentation.xml.rels")
    refs: dict[int, tuple[str, str]] = {}
    for index, slide_id in enumerate(
        pres_root.findall(f".//{_qn(PRESENTATION_NS, 'sldId')}"),
        start=1,
    ):
        rel_id = slide_id.attrib.get(_qn(REL_NS, "id"))
        if not rel_id or rel_id not in rels:
            continue
        rel = rels[rel_id]
        if rel["type"] != SLIDE_REL_TYPE:
            continue
        part_name = _normalize_part(rel["target"])
        refs[index] = (part_name, _rels_name_for_part(part_name))
    return refs


def _chart_shape_refs(slide_xml: bytes, *, source_slide: int) -> list[dict[str, Any]]:
    root = ET.fromstring(slide_xml)
    refs: list[dict[str, Any]] = []
    for order, frame in enumerate(root.findall(f".//{_qn(PRESENTATION_NS, 'graphicFrame')}"), start=1):
        chart = frame.find(f".//{_qn(C_NS, 'chart')}")
        if chart is None:
            continue
        rel_id = chart.attrib.get(_qn(REL_NS, "id"))
        c_nv_pr = frame.find(f".//{_qn(PRESENTATION_NS, 'cNvPr')}")
        shape_id = int(c_nv_pr.attrib.get("id", order)) if c_nv_pr is not None else order
        name = c_nv_pr.attrib.get("name", "") if c_nv_pr is not None else ""
        refs.append(
            {
                "chart_id": f"s{source_slide:02d}_ch{order}",
                "shape_id": shape_id,
                "name": name,
                "rel_id": rel_id,
                "order": order,
            }
        )
    return refs


def _cache_values(cache: ET.Element | None) -> list[str | float]:
    if cache is None:
        return []
    values: list[str | float] = []
    for pt in sorted(cache.findall(_qn(C_NS, "pt")), key=lambda item: int(item.attrib.get("idx", "0"))):
        value = pt.find(_qn(C_NS, "v"))
        text = value.text if value is not None else ""
        if text is None:
            text = ""
        if cache.tag == _qn(C_NS, "numCache"):
            try:
                values.append(float(text))
            except ValueError:
                values.append(text)
        else:
            values.append(text)
    return values


def _series_name(ser: ET.Element) -> str:
    tx = ser.find(_qn(C_NS, "tx"))
    if tx is None:
        return ""
    v = tx.find(_qn(C_NS, "v"))
    if v is not None and v.text:
        return v.text
    cache = tx.find(f".//{_qn(C_NS, 'strCache')}")
    values = _cache_values(cache)
    return str(values[0]) if values else ""


def _read_chart_data(entries: dict[str, bytes], chart_part: str | None) -> dict[str, Any]:
    if not chart_part or chart_part not in entries:
        return {"categories": [], "series": []}
    root = ET.fromstring(entries[chart_part])
    series_payload: list[dict[str, Any]] = []
    categories: list[str | float] = []
    for ser in root.findall(f".//{_qn(C_NS, 'ser')}"):
        cat = ser.find(_qn(C_NS, "cat"))
        cat_cache = cat.find(f".//{_qn(C_NS, 'strCache')}") if cat is not None else None
        if cat_cache is None and cat is not None:
            cat_cache = cat.find(f".//{_qn(C_NS, 'numCache')}")
        ser_categories = _cache_values(cat_cache)
        if ser_categories and not categories:
            categories = ser_categories
        val = ser.find(_qn(C_NS, "val"))
        val_cache = val.find(f".//{_qn(C_NS, 'numCache')}") if val is not None else None
        series_payload.append(
            {
                "name": _series_name(ser),
                "values": _cache_values(val_cache),
            }
        )
    return {"categories": categories, "series": series_payload}


def _analyze_chart_inventory(source: Path) -> dict[int, dict[int, dict[str, Any]]]:
    entries = _zip_entries(source)
    inventory: dict[int, dict[int, dict[str, Any]]] = {}
    for slide_index, (slide_part, rels_name) in _slide_refs(entries).items():
        rels = _relationships(entries, rels_name)
        charts: dict[int, dict[str, Any]] = {}
        for chart_ref in _chart_shape_refs(entries[slide_part], source_slide=slide_index):
            rel_id = chart_ref.get("rel_id")
            rel = rels.get(str(rel_id)) if rel_id else None
            chart_part = None
            if rel and rel.get("type") == CHART_REL_TYPE:
                chart_part = _normalize_part(rel["target"], slide_part)
            chart_ref["chart_part"] = chart_part or ""
            chart_ref.update(_read_chart_data(entries, chart_part))
            charts[int(chart_ref["shape_id"])] = chart_ref
        inventory[slide_index] = charts
    return inventory


def analyze_pptx(pptx_path: str | Path) -> dict[str, Any]:
    source = Path(pptx_path)
    prs = Presentation(str(source))
    chart_inventory = _analyze_chart_inventory(source)
    slides = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slots: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []
        charts: list[dict[str, Any]] = []
        text_parts: list[str] = []
        slot_index = 0
        table_index = 0
        chart_index = 0
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text:
                slot_index += 1
                text_parts.append(text)
                slots.append(
                    {
                        "slot_id": f"s{slide_index:02d}_sh{shape.shape_id}",
                        "shape_id": int(shape.shape_id),
                        "name": getattr(shape, "name", ""),
                        "role": _slot_role(slot_index - 1, text, shape),
                        "geometry": _shape_geometry(shape),
                        "paragraph_count": len(shape.text_frame.paragraphs),
                        "text": text,
                    }
                )
            if getattr(shape, "has_table", False):
                table_index += 1
                tables.append(
                    {
                        "table_id": f"s{slide_index:02d}_tbl{table_index}",
                        "shape_id": int(shape.shape_id),
                        "rows": len(shape.table.rows),
                        "cols": len(shape.table.columns),
                        "geometry": _shape_geometry(shape),
                    }
                )
            if getattr(shape, "has_chart", False):
                chart_index += 1
                chart_info = chart_inventory.get(slide_index, {}).get(int(shape.shape_id), {})
                charts.append(
                    {
                        "chart_id": f"s{slide_index:02d}_ch{chart_index}",
                        "shape_id": int(shape.shape_id),
                        "name": getattr(shape, "name", ""),
                        "geometry": _shape_geometry(shape),
                        "chart_part": chart_info.get("chart_part", ""),
                        "categories": chart_info.get("categories", []),
                        "series": chart_info.get("series", []),
                    }
                )
        slides.append(
            {
                "slide": slide_index,
                "page_type": _page_type(slide_index, len(prs.slides), slots),
                "text_summary": " ".join(text_parts)[:500],
                "slots": slots,
                "tables": tables,
                "charts": charts,
            }
        )
    return {
        "schema": "template_fill_pptx_library.v1",
        "source_pptx": str(source),
        "slide_count": len(prs.slides),
        "slides": slides,
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }


def parse_slide_selection(value: str | None, slide_count: int) -> list[int]:
    if not value:
        return list(range(1, slide_count + 1))
    selected: list[int] = []
    for part in value.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            selected.extend(range(int(start_s), int(end_s) + 1))
        else:
            selected.append(int(part))
    invalid = [idx for idx in selected if idx < 1 or idx > slide_count]
    if invalid:
        raise ValueError(f"Slide selection out of range: {invalid}")
    return selected


def scaffold_plan(library: dict[str, Any], slides: str | None = None) -> dict[str, Any]:
    selected = parse_slide_selection(slides, int(library["slide_count"]))
    by_slide = {item["slide"]: item for item in library["slides"]}
    plan_slides = []
    for slide_num in selected:
        slide = by_slide[slide_num]
        plan_slides.append(
            {
                "source_slide": slide_num,
                "purpose": slide["page_type"],
                "layout_rationale": {
                    "layout_pattern": slide["page_type"],
                    "why_fit": "Scaffolded from source slide inventory; revise before apply.",
                    "risk": "Text capacity and semantic fit require review.",
                },
                "notes": "",
                "transition": None,
                "replacements": [
                    {"slot_id": slot["slot_id"], "text": slot["text"]}
                    for slot in slide.get("slots", [])
                ],
                "table_edits": [],
                "chart_edits": [],
            }
        )
    return {
        "schema": SCHEMA,
        "status": "draft",
        "source_pptx": library["source_pptx"],
        "accepted_warnings": [],
        "slides": plan_slides,
    }


def validate_plan(plan: dict[str, Any], library: dict[str, Any] | None = None) -> tuple[list[str], list[str]]:
    errors: list[str] = []
    warnings: list[str] = []
    if plan.get("schema") != SCHEMA:
        errors.append(f"schema must be {SCHEMA}")
    source = Path(str(plan.get("source_pptx", "")))
    if not source.is_file():
        errors.append(f"source_pptx not found: {source}")
    slides = plan.get("slides")
    if not isinstance(slides, list) or not slides:
        errors.append("slides must be a non-empty list")
        return errors, warnings
    slide_count = None
    known_slots: dict[int, set[str]] = {}
    known_charts: dict[int, set[str]] = {}
    if library:
        slide_count = int(library.get("slide_count", 0) or 0)
        known_slots = {
            int(item["slide"]): {slot["slot_id"] for slot in item.get("slots", [])}
            for item in library.get("slides", [])
        }
        known_charts = {
            int(item["slide"]): {chart["chart_id"] for chart in item.get("charts", [])}
            for item in library.get("slides", [])
        }
    for idx, slide in enumerate(slides, start=1):
        source_slide = slide.get("source_slide")
        if not isinstance(source_slide, int):
            errors.append(f"slides[{idx}].source_slide must be an integer")
            continue
        if slide_count and not (1 <= source_slide <= slide_count):
            errors.append(f"slides[{idx}].source_slide out of range: {source_slide}")
        if not slide.get("layout_rationale"):
            warnings.append(f"slides[{idx}] has no layout_rationale")
        for repl in slide.get("replacements", []):
            slot_id = repl.get("slot_id")
            if not slot_id or "text" not in repl:
                errors.append(f"slides[{idx}] replacement must include slot_id and text")
            elif known_slots and slot_id not in known_slots.get(source_slide, set()):
                errors.append(f"slides[{idx}] unknown slot_id for source slide {source_slide}: {slot_id}")
        for chart_edit in slide.get("chart_edits", []):
            chart_id = chart_edit.get("chart_id")
            if not chart_id and "shape_id" not in chart_edit:
                errors.append(f"slides[{idx}] chart_edit must include chart_id or shape_id")
            elif chart_id and known_charts and chart_id not in known_charts.get(source_slide, set()):
                errors.append(f"slides[{idx}] unknown chart_id for source slide {source_slide}: {chart_id}")
            categories = chart_edit.get("categories")
            series = chart_edit.get("series")
            if categories is not None and not isinstance(categories, list):
                errors.append(f"slides[{idx}] chart_edit categories must be a list")
            if series is not None and not isinstance(series, list):
                errors.append(f"slides[{idx}] chart_edit series must be a list")
            if isinstance(categories, list) and isinstance(series, list):
                for series_index, item in enumerate(series, start=1):
                    if not isinstance(item, dict) or not isinstance(item.get("values"), list):
                        errors.append(f"slides[{idx}] chart_edit series[{series_index}] must include values list")
                    elif len(item["values"]) != len(categories):
                        errors.append(
                            f"slides[{idx}] chart_edit series[{series_index}] values length must match categories"
                        )
    return errors, warnings


def _load_json(path: str | Path) -> dict[str, Any]:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def _write_json(path: str | Path, data: dict[str, Any]) -> None:
    target = Path(path)
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _write_pptx_entries(entries: dict[str, bytes], output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for name in sorted(entries):
            zf.writestr(name, entries[name])


def _max_slide_part_number(entries: dict[str, bytes]) -> int:
    numbers = []
    for name in entries:
        match = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", name)
        if match:
            numbers.append(int(match.group(1)))
    return max(numbers, default=0)


def _max_chart_part_number(entries: dict[str, bytes]) -> int:
    numbers = []
    for name in entries:
        match = re.fullmatch(r"ppt/charts/chart(\d+)\.xml", name)
        if match:
            numbers.append(int(match.group(1)))
    return max(numbers, default=0)


def _max_embedding_part_number(entries: dict[str, bytes]) -> int:
    numbers = []
    for name in entries:
        match = re.search(r"ppt/embeddings/.*?(\d+)\.xlsx$", name)
        if match:
            numbers.append(int(match.group(1)))
    return max(numbers, default=0)


def _max_numeric_rid(rels_root: ET.Element) -> int:
    numbers = []
    for rel in rels_root.findall(_qn(PACKAGE_REL_NS, "Relationship")):
        rel_id = rel.attrib.get("Id", "")
        if rel_id.startswith("rId") and rel_id[3:].isdigit():
            numbers.append(int(rel_id[3:]))
    return max(numbers, default=0)


def _max_slide_id(slide_id_list: ET.Element) -> int:
    ids = []
    for slide_id in slide_id_list.findall(_qn(PRESENTATION_NS, "sldId")):
        value = slide_id.attrib.get("id", "")
        if value.isdigit():
            ids.append(int(value))
    return max(ids, default=255)


def _add_slide_override(content_root: ET.Element, part_name: str) -> None:
    part_name = "/" + part_name.lstrip("/")
    for override in content_root.findall(_qn(CONTENT_TYPE_NS, "Override")):
        if override.attrib.get("PartName") == part_name:
            return
    ET.SubElement(
        content_root,
        _qn(CONTENT_TYPE_NS, "Override"),
        {"PartName": part_name, "ContentType": SLIDE_CONTENT_TYPE},
    )


def _add_chart_override(content_root: ET.Element, part_name: str) -> None:
    part_name = "/" + part_name.lstrip("/")
    for override in content_root.findall(_qn(CONTENT_TYPE_NS, "Override")):
        if override.attrib.get("PartName") == part_name:
            return
    ET.SubElement(
        content_root,
        _qn(CONTENT_TYPE_NS, "Override"),
        {"PartName": part_name, "ContentType": CHART_CONTENT_TYPE},
    )


def _add_default_content_type(content_root: ET.Element, extension: str, content_type: str) -> None:
    extension = extension.lstrip(".")
    for default in content_root.findall(_qn(CONTENT_TYPE_NS, "Default")):
        if default.attrib.get("Extension") == extension:
            return
    ET.SubElement(
        content_root,
        _qn(CONTENT_TYPE_NS, "Default"),
        {"Extension": extension, "ContentType": content_type},
    )


def _override_content_type(content_root: ET.Element, part_name: str) -> str | None:
    part_name = "/" + part_name.lstrip("/")
    for override in content_root.findall(_qn(CONTENT_TYPE_NS, "Override")):
        if override.attrib.get("PartName") == part_name:
            return override.attrib.get("ContentType")
    return None


def _add_part_override(content_root: ET.Element, part_name: str, content_type: str) -> None:
    part_name = "/" + part_name.lstrip("/")
    for override in content_root.findall(_qn(CONTENT_TYPE_NS, "Override")):
        if override.attrib.get("PartName") == part_name:
            return
    ET.SubElement(
        content_root,
        _qn(CONTENT_TYPE_NS, "Override"),
        {"PartName": part_name, "ContentType": content_type},
    )


def _allocate_sibling_part(entries: dict[str, bytes], source_part: str) -> str:
    directory = posixpath.dirname(source_part)
    stem, ext = posixpath.splitext(posixpath.basename(source_part))
    index = 1
    while True:
        candidate = posixpath.join(directory, f"{stem}_tf{index}{ext}")
        if candidate not in entries:
            return candidate
        index += 1


def _clone_private_dependency(
    entries: dict[str, bytes],
    content_root: ET.Element,
    *,
    owner_part: str,
    source_part: str,
    cloned: dict[str, str],
) -> str | None:
    if source_part not in entries:
        return None
    content_type = _override_content_type(content_root, source_part)
    if content_type is None:
        return None
    if source_part in cloned:
        return cloned[source_part]

    new_part = _allocate_sibling_part(entries, source_part)
    entries[new_part] = entries[source_part]
    cloned[source_part] = new_part
    _add_part_override(content_root, new_part, content_type)

    source_rels = _rels_name_for_part(source_part)
    if source_rels in entries:
        rels_root = ET.fromstring(entries[source_rels])
        _rewrite_private_dependency_rels(
            entries,
            content_root,
            rels_root,
            owner_part=new_part,
            cloned=cloned,
        )
        entries[_rels_name_for_part(new_part)] = ET.tostring(
            rels_root,
            encoding="utf-8",
            xml_declaration=True,
        )
    return new_part


def _rewrite_private_dependency_rels(
    entries: dict[str, bytes],
    content_root: ET.Element,
    rels_root: ET.Element,
    *,
    owner_part: str,
    cloned: dict[str, str],
) -> None:
    for rel in rels_root.findall(_qn(PACKAGE_REL_NS, "Relationship")):
        if rel.attrib.get("TargetMode") == "External":
            continue
        rel_type = rel.attrib.get("Type")
        if rel_type in SHARED_REL_TYPES or rel_type in STAGE_OWNED_REL_TYPES:
            continue
        target = rel.attrib.get("Target")
        if not target:
            continue
        source_part = _normalize_part(target, owner_part)
        new_part = _clone_private_dependency(
            entries,
            content_root,
            owner_part=owner_part,
            source_part=source_part,
            cloned=cloned,
        )
        if new_part:
            rel.set("Target", _relative_target(owner_part, new_part))


def _clone_plan_slides(source: Path, output: Path, plan_slides: list[dict[str, Any]]) -> None:
    entries = _zip_entries(source)
    refs = _slide_refs(entries)
    pres_root = ET.fromstring(entries["ppt/presentation.xml"])
    rels_root = ET.fromstring(entries["ppt/_rels/presentation.xml.rels"])
    content_root = ET.fromstring(entries["[Content_Types].xml"])
    slide_id_list = pres_root.find(_qn(PRESENTATION_NS, "sldIdLst"))
    if slide_id_list is None:
        slide_id_list = ET.SubElement(pres_root, _qn(PRESENTATION_NS, "sldIdLst"))

    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for rel in list(rels_root.findall(_qn(PACKAGE_REL_NS, "Relationship"))):
        if rel.attrib.get("Type") == SLIDE_REL_TYPE:
            rels_root.remove(rel)

    next_slide_number = _max_slide_part_number(entries) + 1
    next_slide_id = _max_slide_id(slide_id_list) + 1
    next_rid = _max_numeric_rid(rels_root) + 1

    for offset, plan_slide in enumerate(plan_slides):
        source_slide = int(plan_slide["source_slide"])
        if source_slide not in refs:
            raise ValueError(f"Plan references missing source slide: {source_slide}")
        source_part, source_rels = refs[source_slide]
        new_slide_number = next_slide_number + offset
        new_part = f"ppt/slides/slide{new_slide_number}.xml"
        new_rels = f"ppt/slides/_rels/slide{new_slide_number}.xml.rels"
        new_rid = f"rId{next_rid + offset}"

        entries[new_part] = entries[source_part]
        if source_rels in entries:
            slide_rels_root = ET.fromstring(entries[source_rels])
            _rewrite_private_dependency_rels(
                entries,
                content_root,
                slide_rels_root,
                owner_part=new_part,
                cloned={},
            )
            entries[new_rels] = ET.tostring(
                slide_rels_root,
                encoding="utf-8",
                xml_declaration=True,
            )
        else:
            entries[new_rels] = (
                f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
                f'<Relationships xmlns="{PACKAGE_REL_NS}">\n</Relationships>'
            ).encode("utf-8")
        _add_slide_override(content_root, new_part)
        ET.SubElement(
            rels_root,
            _qn(PACKAGE_REL_NS, "Relationship"),
            {
                "Id": new_rid,
                "Type": SLIDE_REL_TYPE,
                "Target": f"slides/slide{new_slide_number}.xml",
            },
        )
        ET.SubElement(
            slide_id_list,
            _qn(PRESENTATION_NS, "sldId"),
            {"id": str(next_slide_id + offset), _qn(REL_NS, "id"): new_rid},
        )

    entries["ppt/presentation.xml"] = ET.tostring(pres_root, encoding="utf-8", xml_declaration=True)
    entries["ppt/_rels/presentation.xml.rels"] = ET.tostring(
        rels_root,
        encoding="utf-8",
        xml_declaration=True,
    )
    entries["[Content_Types].xml"] = ET.tostring(content_root, encoding="utf-8", xml_declaration=True)
    _write_pptx_entries(entries, output)


def _set_cache_points(cache: ET.Element, values: list[Any]) -> None:
    for child in list(cache):
        if child.tag in {_qn(C_NS, "ptCount"), _qn(C_NS, "pt")}:
            cache.remove(child)
    ET.SubElement(cache, _qn(C_NS, "ptCount"), {"val": str(len(values))})
    for index, value in enumerate(values):
        pt = ET.SubElement(cache, _qn(C_NS, "pt"), {"idx": str(index)})
        node = ET.SubElement(pt, _qn(C_NS, "v"))
        node.text = str(value)


def _set_series_name(ser: ET.Element, name: str) -> None:
    tx = ser.find(_qn(C_NS, "tx"))
    if tx is None:
        tx = ET.SubElement(ser, _qn(C_NS, "tx"))
    v = tx.find(_qn(C_NS, "v"))
    if v is not None:
        v.text = name
    cache = tx.find(f".//{_qn(C_NS, 'strCache')}")
    if cache is not None:
        _set_cache_points(cache, [name])


def _set_chart_xml_data(chart_xml: bytes, chart_edit: dict[str, Any]) -> bytes:
    root = ET.fromstring(chart_xml)
    categories = chart_edit.get("categories")
    series_payload = chart_edit.get("series")
    if not isinstance(series_payload, list):
        return chart_xml
    series_nodes = root.findall(f".//{_qn(C_NS, 'ser')}")
    if len(series_payload) > len(series_nodes):
        raise ValueError("chart_edit has more series than the source chart")
    for index, series_edit in enumerate(series_payload):
        if not isinstance(series_edit, dict):
            raise ValueError("chart_edit series entries must be objects")
        ser = series_nodes[index]
        if "name" in series_edit:
            _set_series_name(ser, str(series_edit["name"]))
        values = series_edit.get("values")
        if isinstance(values, list):
            val = ser.find(_qn(C_NS, "val"))
            val_cache = val.find(f".//{_qn(C_NS, 'numCache')}") if val is not None else None
            if val_cache is None:
                raise ValueError("source chart series has no numeric cache to edit")
            _set_cache_points(val_cache, values)
        if isinstance(categories, list):
            cat = ser.find(_qn(C_NS, "cat"))
            cat_cache = cat.find(f".//{_qn(C_NS, 'strCache')}") if cat is not None else None
            if cat_cache is None and cat is not None:
                cat_cache = cat.find(f".//{_qn(C_NS, 'numCache')}")
            if cat_cache is not None:
                _set_cache_points(cat_cache, categories)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _workbook_bytes_with_chart_data(workbook_bytes: bytes, chart_edit: dict[str, Any]) -> bytes:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - requirements include openpyxl
        raise RuntimeError("openpyxl is required to sync embedded chart workbooks") from exc

    categories = chart_edit.get("categories")
    series_payload = chart_edit.get("series")
    if not isinstance(series_payload, list):
        return workbook_bytes
    input_buffer = BytesIO(workbook_bytes)
    workbook = load_workbook(input_buffer)
    sheet = workbook.active

    max_rows = max(sheet.max_row, (len(categories) + 1) if isinstance(categories, list) else 1)
    max_cols = max(sheet.max_column, len(series_payload) + 1)
    for row in sheet.iter_rows(min_row=1, max_row=max_rows, min_col=1, max_col=max_cols):
        for cell in row:
            cell.value = None

    sheet.cell(row=1, column=1).value = None
    if isinstance(categories, list):
        for row_index, category in enumerate(categories, start=2):
            sheet.cell(row=row_index, column=1).value = category
    for column_index, series_edit in enumerate(series_payload, start=2):
        if not isinstance(series_edit, dict):
            raise ValueError("chart_edit series entries must be objects")
        sheet.cell(row=1, column=column_index).value = series_edit.get("name", f"Series {column_index - 1}")
        values = series_edit.get("values")
        if not isinstance(values, list):
            continue
        for row_index, value in enumerate(values, start=2):
            sheet.cell(row=row_index, column=column_index).value = value

    output_buffer = BytesIO()
    workbook.save(output_buffer)
    return output_buffer.getvalue()


def _chart_edit_matches(ref: dict[str, Any], edit: dict[str, Any]) -> bool:
    chart_id = edit.get("chart_id")
    if chart_id and chart_id == ref.get("chart_id"):
        return True
    if "shape_id" in edit and int(edit["shape_id"]) == int(ref.get("shape_id", -1)):
        return True
    if "order" in edit and int(edit["order"]) == int(ref.get("order", -1)):
        return True
    return False


def _clone_chart_for_slide(
    entries: dict[str, bytes],
    content_root: ET.Element,
    rel: ET.Element,
    *,
    owner_slide_part: str,
    next_chart_number: int,
    next_embedding_number: int,
) -> tuple[str, int, int]:
    source_part = _normalize_part(str(rel.attrib["Target"]), owner_slide_part)
    if source_part not in entries:
        raise ValueError(f"chart part not found: {source_part}")
    next_chart_number += 1
    new_part = f"ppt/charts/chart{next_chart_number}.xml"
    while new_part in entries:
        next_chart_number += 1
        new_part = f"ppt/charts/chart{next_chart_number}.xml"
    entries[new_part] = entries[source_part]
    source_rels = _rels_name_for_part(source_part)
    if source_rels in entries:
        entries[_rels_name_for_part(new_part)] = entries[source_rels]
    rel.set("Target", _relative_target(owner_slide_part, new_part))
    _add_chart_override(content_root, new_part)
    return new_part, next_chart_number, next_embedding_number


def _clone_and_sync_chart_workbook(
    entries: dict[str, bytes],
    content_root: ET.Element,
    chart_part: str,
    chart_edit: dict[str, Any],
    next_embedding_number: int,
) -> int:
    chart_rels_name = _rels_name_for_part(chart_part)
    if chart_rels_name not in entries:
        return next_embedding_number
    chart_rels_root = ET.fromstring(entries[chart_rels_name])
    updated = False
    for rel in chart_rels_root.findall(_qn(PACKAGE_REL_NS, "Relationship")):
        if rel.attrib.get("Type") != PACKAGE_REL_TYPE or rel.attrib.get("TargetMode") == "External":
            continue
        target = rel.attrib.get("Target")
        if not target:
            continue
        source_workbook = _normalize_part(target, chart_part)
        if source_workbook not in entries:
            continue
        next_embedding_number += 1
        new_workbook = f"ppt/embeddings/Microsoft_Excel_Worksheet{next_embedding_number}.xlsx"
        while new_workbook in entries:
            next_embedding_number += 1
            new_workbook = f"ppt/embeddings/Microsoft_Excel_Worksheet{next_embedding_number}.xlsx"
        entries[new_workbook] = _workbook_bytes_with_chart_data(entries[source_workbook], chart_edit)
        rel.set("Target", _relative_target(chart_part, new_workbook))
        _add_default_content_type(content_root, "xlsx", XLSX_CONTENT_TYPE)
        updated = True
    if updated:
        entries[chart_rels_name] = ET.tostring(chart_rels_root, encoding="utf-8", xml_declaration=True)
    return next_embedding_number


def _apply_chart_edits(output: Path, plan_slides: list[dict[str, Any]]) -> None:
    if not any(slide.get("chart_edits") for slide in plan_slides):
        return
    entries = _zip_entries(output)
    refs = _slide_refs(entries)
    content_root = ET.fromstring(entries["[Content_Types].xml"])
    next_chart_number = _max_chart_part_number(entries)
    next_embedding_number = _max_embedding_part_number(entries)
    for output_index, plan_slide in enumerate(plan_slides, start=1):
        chart_edits = plan_slide.get("chart_edits") or []
        if not chart_edits or output_index not in refs:
            continue
        slide_part, rels_name = refs[output_index]
        if rels_name not in entries:
            continue
        rels_root = ET.fromstring(entries[rels_name])
        rels_by_id = {
            rel.attrib.get("Id"): rel
            for rel in rels_root.findall(_qn(PACKAGE_REL_NS, "Relationship"))
        }
        chart_refs = _chart_shape_refs(entries[slide_part], source_slide=int(plan_slide["source_slide"]))
        for edit in chart_edits:
            matches = [ref for ref in chart_refs if _chart_edit_matches(ref, edit)]
            if not matches:
                raise ValueError(f"chart_edit did not match a chart on slide {output_index}: {edit}")
            ref = matches[0]
            rel = rels_by_id.get(ref.get("rel_id"))
            if rel is None or rel.attrib.get("Type") != CHART_REL_TYPE:
                raise ValueError(f"chart relationship not found for {ref.get('chart_id')}")
            chart_part, next_chart_number, next_embedding_number = _clone_chart_for_slide(
                entries,
                content_root,
                rel,
                owner_slide_part=slide_part,
                next_chart_number=next_chart_number,
                next_embedding_number=next_embedding_number,
            )
            entries[chart_part] = _set_chart_xml_data(entries[chart_part], edit)
            next_embedding_number = _clone_and_sync_chart_workbook(
                entries,
                content_root,
                chart_part,
                edit,
                next_embedding_number,
            )
        entries[rels_name] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)
    entries["[Content_Types].xml"] = ET.tostring(content_root, encoding="utf-8", xml_declaration=True)
    _write_pptx_entries(entries, output)


def apply_plan(plan: dict[str, Any], output_path: str | Path) -> None:
    source = Path(plan["source_pptx"])
    output = Path(output_path)
    plan_slides = plan.get("slides", [])
    _clone_plan_slides(source, output, plan_slides)
    prs = Presentation(str(output))
    for plan_slide, slide in zip(plan_slides, prs.slides):
        replacements = {
            repl["slot_id"]: repl["text"]
            for repl in plan_slide.get("replacements", [])
        }
        for shape in slide.shapes:
            slot_id = f"s{plan_slide['source_slide']:02d}_sh{shape.shape_id}"
            if slot_id in replacements and getattr(shape, "has_text_frame", False):
                shape.text = str(replacements[slot_id])
            if getattr(shape, "has_table", False):
                for table_edit in plan_slide.get("table_edits", []):
                    if int(table_edit.get("shape_id", -1)) != int(shape.shape_id):
                        continue
                    for cell in table_edit.get("cells", []):
                        row = int(cell["row"])
                        col = int(cell["col"])
                        shape.table.cell(row, col).text = str(cell["text"])
    prs.save(str(output))
    _apply_chart_edits(output, plan_slides)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Analyze and fill native PPTX templates")
    subparsers = parser.add_subparsers(dest="command", required=True)

    analyze = subparsers.add_parser("analyze", help="Analyze a PPTX slide library")
    analyze.add_argument("pptx")
    analyze.add_argument("-o", "--output", required=True)

    scaffold = subparsers.add_parser("scaffold", help="Create a fill_plan.json")
    scaffold.add_argument("library_json")
    scaffold.add_argument("-o", "--output", required=True)
    scaffold.add_argument("--slides", help="Slide selection, e.g. 1,3-5")

    validate = subparsers.add_parser("validate", help="Validate fill_plan.json")
    validate.add_argument("fill_plan")
    validate.add_argument("--library")
    validate.add_argument("--json", action="store_true")

    apply = subparsers.add_parser("apply", help="Apply a native fill plan")
    apply.add_argument("fill_plan")
    apply.add_argument("-o", "--output", required=True)
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        if args.command == "analyze":
            library = analyze_pptx(args.pptx)
            _write_json(args.output, library)
            print(f"[OK] Wrote slide library: {args.output}")
            return 0
        if args.command == "scaffold":
            plan = scaffold_plan(_load_json(args.library_json), slides=args.slides)
            _write_json(args.output, plan)
            print(f"[OK] Wrote fill plan: {args.output}")
            return 0
        if args.command == "validate":
            library = _load_json(args.library) if args.library else None
            errors, warnings = validate_plan(_load_json(args.fill_plan), library)
            if args.json:
                print(json.dumps({"errors": errors, "warnings": warnings}, ensure_ascii=False, indent=2))
            else:
                for warning in warnings:
                    print(f"[WARN] {warning}")
                for error in errors:
                    print(f"[ERROR] {error}")
                if not errors:
                    print("[OK] Fill plan is valid")
            return 0 if not errors else 1
        if args.command == "apply":
            apply_plan(_load_json(args.fill_plan), args.output)
            print(f"[OK] Wrote filled PPTX: {args.output}")
            return 0
    except Exception as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 1
    parser.error(f"Unknown command: {args.command}")
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
