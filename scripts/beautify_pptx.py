#!/usr/bin/env python3
"""Conservative 1:1 PPTX beautification MVP."""

from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation

SCRIPTS_DIR = Path(__file__).resolve().parent
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_DIR))

from workflow_manifest import update_manifest  # noqa: E402


SCHEMA_VERSION = "easyslides.beautify_report.v1"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
HEX_COLOR = re.compile(r"^#[0-9A-Fa-f]{6}$")
MEDIA_RISK_SUFFIXES = {".mp4", ".mov", ".avi", ".wmv", ".mp3", ".m4a", ".wav"}

ET.register_namespace("a", DRAWING_NS)


def _validate_hex(value: str, field: str) -> str:
    if not HEX_COLOR.match(value):
        raise ValueError(f"{field} must be a #RRGGBB color")
    return value.upper()


def _slide_texts(pptx_path: Path) -> list[list[str]]:
    prs = Presentation(str(pptx_path))
    slides: list[list[str]] = []
    for slide in prs.slides:
        texts = [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        ]
        slides.append(texts)
    return slides


def _package_risks(pptx_path: Path) -> dict[str, int]:
    with zipfile.ZipFile(pptx_path) as zf:
        names = [item.filename for item in zf.infolist() if not item.is_dir()]
    return {
        "charts": sum(name.startswith("ppt/charts/") and name.endswith(".xml") for name in names),
        "smartart_parts": sum(name.startswith("ppt/diagrams/") for name in names),
        "embedded_objects": sum(name.startswith("ppt/embeddings/") for name in names),
        "media_files": sum(Path(name).suffix.lower() in MEDIA_RISK_SUFFIXES for name in names),
    }


def inspect_pptx(pptx_path: str | Path, *, status: str = "inspected", actions: list[str] | None = None) -> dict[str, Any]:
    source = Path(pptx_path)
    if not source.is_file():
        raise FileNotFoundError(f"PPTX not found: {source}")
    texts = _slide_texts(source)
    return {
        "schema_version": SCHEMA_VERSION,
        "route": "beautify-pptx",
        "status": status,
        "source_pptx": str(source),
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "slide_count": len(texts),
        "text_by_slide": texts,
        "package_risks": _package_risks(source),
        "actions": actions or [],
        "notes": [
            "MVP preserves slide count, slide order, visible text, and object structure.",
            "Theme color patch affects only theme-bound colors; direct RGB object colors are left unchanged.",
        ],
    }


def _write_report(report: dict[str, Any], report_dir: Path, *, stage: str, output_pptx: Path | None = None) -> Path:
    report_dir.mkdir(parents=True, exist_ok=True)
    report_path = report_dir / "beautify_report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    artifacts = {
        "source_pptx": report["source_pptx"],
        "beautify_report": report_path.name,
    }
    if output_pptx:
        artifacts["output_pptx"] = str(output_pptx)
    update_manifest(
        report_dir,
        route="beautify-pptx",
        stage=stage,
        status=report["status"],
        artifacts=artifacts,
    )
    return report_path


def _set_scheme_color(root: ET.Element, color_name: str, value: str) -> bool:
    scheme = root.find(f".//{{{DRAWING_NS}}}clrScheme")
    if scheme is None:
        return False
    target = scheme.find(f"{{{DRAWING_NS}}}{color_name}")
    if target is None:
        return False
    for child in list(target):
        target.remove(child)
    ET.SubElement(target, f"{{{DRAWING_NS}}}srgbClr", {"val": value.lstrip("#")})
    return True


def patch_theme_colors(source_pptx: Path, output_pptx: Path, *, primary: str, accent: str) -> list[str]:
    primary = _validate_hex(primary, "primary")
    accent = _validate_hex(accent, "accent")
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(source_pptx, "r") as zf:
        entries = {info.filename: zf.read(info.filename) for info in zf.infolist() if not info.is_dir()}

    patched = False
    for name, data in list(entries.items()):
        if not (name.startswith("ppt/theme/") and name.endswith(".xml")):
            continue
        root = ET.fromstring(data)
        changed = False
        changed = _set_scheme_color(root, "accent1", primary) or changed
        changed = _set_scheme_color(root, "accent2", accent) or changed
        changed = _set_scheme_color(root, "hlink", primary) or changed
        changed = _set_scheme_color(root, "folHlink", accent) or changed
        if changed:
            entries[name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            patched = True

    if not patched:
        raise RuntimeError("No patchable PPT theme color scheme found")

    with zipfile.ZipFile(output_pptx, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in sorted(entries.items()):
            zf.writestr(name, data)
    return ["theme_color_patch"]


def apply_theme_patch(
    source_pptx: str | Path,
    output_pptx: str | Path,
    *,
    primary: str,
    accent: str,
) -> dict[str, Any]:
    source = Path(source_pptx)
    output = Path(output_pptx)
    before_texts = _slide_texts(source)
    actions = patch_theme_colors(source, output, primary=primary, accent=accent)
    after_texts = _slide_texts(output)
    if after_texts != before_texts:
        output.unlink(missing_ok=True)
        raise RuntimeError("Theme patch changed visible text; output was discarded")
    report = inspect_pptx(source, status="applied", actions=actions)
    report["output_pptx"] = str(output)
    report["after_slide_count"] = len(after_texts)
    report["preserved_visible_text"] = True
    return report


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Conservative 1:1 PPTX beautification.")
    subparsers = parser.add_subparsers(dest="command", required=True)

    inspect_cmd = subparsers.add_parser("inspect", help="Inspect a PPTX and write a beautify report.")
    inspect_cmd.add_argument("pptx")
    inspect_cmd.add_argument("--out", required=True, help="Report directory.")

    apply_cmd = subparsers.add_parser("apply", help="Apply a safe theme color patch.")
    apply_cmd.add_argument("pptx")
    apply_cmd.add_argument("-o", "--output", required=True)
    apply_cmd.add_argument("--report-dir", required=True)
    apply_cmd.add_argument("--primary", default="#2454A6")
    apply_cmd.add_argument("--accent", default="#E9B44C")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        if args.command == "inspect":
            report = inspect_pptx(args.pptx)
            path = _write_report(report, Path(args.out), stage="inspect")
            print(json.dumps({"status": "ok", "report": str(path)}, ensure_ascii=False, indent=2))
            return 0
        if args.command == "apply":
            report = apply_theme_patch(
                args.pptx,
                args.output,
                primary=args.primary,
                accent=args.accent,
            )
            path = _write_report(report, Path(args.report_dir), stage="apply", output_pptx=Path(args.output))
            print(json.dumps({"status": "ok", "output": args.output, "report": str(path)}, ensure_ascii=False, indent=2))
            return 0
    except Exception as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 1
    parser.error(f"Unknown command: {args.command}")
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
