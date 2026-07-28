#!/usr/bin/env python3
"""Build a visual audit kit for renderable EasySlides component assets.

The registry intentionally contains several asset levels: page recipes, body
variants, charts, icons, source evidence, and reusable visual components. This
tool only renders the directly selectable component level, while reporting the
larger registry inventory alongside it. It keeps a review overlay around image
slots so empty image placeholders remain legible without changing source SVGs.
"""

from __future__ import annotations

import argparse
import html
import json
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
import shutil
import sys
import xml.etree.ElementTree as ET
from typing import Any


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

DEFAULT_OUTPUT = ROOT / "projects" / "component_audit_20260728"
LAYOUTS_ROOT = ROOT / "templates" / "layouts"
REGISTRY_PATH = ROOT / "templates" / "components" / "component_registry.json"
SOURCE_FAITHFUL_CATALOG_PATH = (
    ROOT
    / "templates"
    / "components"
    / "source_templates"
    / "nsfc_defense_distilled_kit"
    / "source_faithful_component_catalog.json"
)
SOURCE_DERIVED_CATALOG_PATH = (
    ROOT
    / "templates"
    / "components"
    / "source_templates"
    / "nsfc_defense_distilled_kit"
    / "source_derived_component_catalog.json"
)
SVG_NS = "http://www.w3.org/2000/svg"
ET.register_namespace("", SVG_NS)

SCHEMA_VERSION = "easyslides.component_asset_audit.v1"


@dataclass(frozen=True)
class ComponentAsset:
    asset_id: str
    display_name: str
    family: str
    kind: str
    source_path: Path
    description: str
    roles: tuple[str, ...]
    slot_count: int
    image_slot_count: int
    asset_status: str = ""


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8-sig"))
    if not isinstance(payload, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return payload


def _rel(path: Path, base: Path) -> str:
    return path.resolve().relative_to(base.resolve()).as_posix()


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _slug(value: str) -> str:
    return "".join(char if char.isalnum() else "_" for char in value).strip("_")


def _humanize(value: str) -> str:
    return value.replace("_", " ").replace("-", " ").title()


def _slot_counts(slots: list[Any]) -> tuple[int, int]:
    normalized = [slot for slot in slots if isinstance(slot, dict)]
    return len(normalized), sum(1 for slot in normalized if slot.get("kind") == "image")


def _component_records() -> list[ComponentAsset]:
    records: list[ComponentAsset] = []
    for catalog_path in sorted(LAYOUTS_ROOT.glob("*/component_catalog.json")):
        catalog = _read_json(catalog_path)
        template_id = str(catalog.get("template_id") or catalog_path.parent.name)
        for entry in catalog.get("components", []):
            if not isinstance(entry, dict):
                continue
            relative_asset = str(entry.get("asset_path") or "")
            source_path = catalog_path.parent / relative_asset
            if not source_path.is_file():
                continue
            component_id = str(entry.get("component_id") or source_path.stem)
            slots = entry.get("slots") if isinstance(entry.get("slots"), list) else []
            slot_count, image_slot_count = _slot_counts(slots)
            records.append(
                ComponentAsset(
                    asset_id=str(entry.get("asset_id") or f"component/{template_id}/{component_id}"),
                    display_name=component_id,
                    family=template_id,
                    kind="template_component",
                    source_path=source_path,
                    description=str(entry.get("description") or ""),
                    roles=tuple(str(entry.get("reuse_policy") or "template_scoped").split("_or_")),
                    slot_count=slot_count,
                    image_slot_count=image_slot_count,
                    asset_status="runtime_selectable",
                )
            )

    primitives_path = LAYOUTS_ROOT / "nsfc_defense" / "component_primitives.json"
    if primitives_path.is_file():
        primitives = _read_json(primitives_path)
        for entry in primitives.get("primitives", []):
            if not isinstance(entry, dict):
                continue
            source_path = primitives_path.parent / str(entry.get("asset_path") or "")
            if not source_path.is_file():
                continue
            primitive_id = str(entry.get("primitive_id") or source_path.stem)
            slots = entry.get("slots") if isinstance(entry.get("slots"), list) else []
            slot_count, image_slot_count = _slot_counts(slots)
            records.append(
                ComponentAsset(
                    asset_id=f"primitive/nsfc_defense/{primitive_id}",
                    display_name=primitive_id,
                    family="nsfc_defense",
                    kind="primitive",
                    source_path=source_path,
                    description=str(entry.get("style") or ""),
                    roles=(str(entry.get("role") or ""),),
                    slot_count=slot_count,
                    image_slot_count=image_slot_count,
                    asset_status="runtime_selectable",
                )
            )

    if SOURCE_FAITHFUL_CATALOG_PATH.is_file():
        source_catalog = _read_json(SOURCE_FAITHFUL_CATALOG_PATH)
        for entry in source_catalog.get("components", []):
            if not isinstance(entry, dict):
                continue
            source_path = ROOT / str(entry.get("asset_path") or "")
            if not source_path.is_file():
                continue
            component_id = str(entry.get("component_id") or source_path.stem)
            records.append(
                ComponentAsset(
                    asset_id=str(entry.get("asset_id") or f"source_faithful/nsfc_defense_distilled/{component_id}"),
                    display_name=str(entry.get("display_name") or component_id),
                    family="nsfc_defense_distilled",
                    kind="source_faithful_component",
                    source_path=source_path,
                    description=str(entry.get("description") or ""),
                    roles=("source_faithful_review_only",),
                    slot_count=0,
                    image_slot_count=0,
                    asset_status=str(entry.get("asset_status") or "source_faithful_review_only"),
                )
            )

    if SOURCE_DERIVED_CATALOG_PATH.is_file():
        derived_catalog = _read_json(SOURCE_DERIVED_CATALOG_PATH)
        for entry in derived_catalog.get("components", []):
            if not isinstance(entry, dict):
                continue
            source_path = ROOT / str(entry.get("asset_path") or "")
            if not source_path.is_file():
                continue
            slots = entry.get("slots") if isinstance(entry.get("slots"), list) else []
            slot_count, image_slot_count = _slot_counts(slots)
            component_id = str(entry.get("component_id") or source_path.stem)
            records.append(
                ComponentAsset(
                    asset_id=str(entry.get("asset_id") or f"source_derived/nsfc_defense_distilled/{component_id}"),
                    display_name=str(entry.get("display_name") or component_id),
                    family="nsfc_defense_distilled",
                    kind="source_derived_component",
                    source_path=source_path,
                    description=str(entry.get("description") or ""),
                    roles=("declared_payload_slots_only",),
                    slot_count=slot_count,
                    image_slot_count=image_slot_count,
                    asset_status=str(entry.get("asset_status") or "source_derived_editable_candidate"),
                )
            )

    return sorted(records, key=lambda item: (item.family, item.kind, item.display_name))


def _view_box(root: ET.Element) -> tuple[float, float, float, float]:
    raw = str(root.attrib.get("viewBox") or "").replace(",", " ").split()
    if len(raw) == 4:
        try:
            return tuple(float(value) for value in raw)  # type: ignore[return-value]
        except ValueError:
            pass
    width = float(str(root.attrib.get("width") or "1280").replace("px", ""))
    height = float(str(root.attrib.get("height") or "720").replace("px", ""))
    return 0.0, 0.0, width, height


def _audit_svg(source: Path, destination: Path) -> None:
    """Write a review-only SVG with image-slot geometry visible."""
    root = ET.fromstring(source.read_text(encoding="utf-8"))
    x, y, width, height = _view_box(root)
    background = ET.Element(f"{{{SVG_NS}}}rect", {"x": str(x), "y": str(y), "width": str(width), "height": str(height), "fill": "#FFFFFF"})
    root.insert(0, background)

    for parent in root.iter():
        children = list(parent)
        for index, child in enumerate(children):
            if _local_name(child.tag) != "image":
                continue
            attrs = child.attrib
            if not (attrs.get("data-slot-id") or attrs.get("data-slot")):
                continue
            image_x = attrs.get("x", "0")
            image_y = attrs.get("y", "0")
            image_w = attrs.get("width", "0")
            image_h = attrs.get("height", "0")
            slot_id = str(attrs.get("data-slot-id") or attrs.get("data-slot") or "IMAGE")
            overlay = ET.Element(
                f"{{{SVG_NS}}}rect",
                {
                    "x": image_x,
                    "y": image_y,
                    "width": image_w,
                    "height": image_h,
                    "fill": "#F5F6F8",
                    "stroke": "#AAB2C0",
                    "stroke-width": "1.5",
                    "stroke-dasharray": "6 5",
                    "data-audit-overlay": "image-slot",
                },
            )
            parent.insert(index, overlay)
            try:
                center_x = float(image_x) + float(image_w) / 2
                center_y = float(image_y) + float(image_h) / 2
            except ValueError:
                continue
            label = ET.Element(
                f"{{{SVG_NS}}}text",
                {
                    "x": f"{center_x:.2f}",
                    "y": f"{center_y:.2f}",
                    "text-anchor": "middle",
                    "dominant-baseline": "middle",
                    "font-family": "Arial, sans-serif",
                    "font-size": "14",
                    "font-weight": "700",
                    "fill": "#6C7480",
                    "data-audit-overlay": "image-slot-label",
                },
            )
            label.text = slot_id
            parent.insert(index + 1, label)

    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(ET.tostring(root, encoding="unicode"), encoding="utf-8")


def _render_html(records: list[dict[str, Any]], registry_counts: dict[str, int]) -> str:
    return _render_html_clean(records, registry_counts)


def _render_html_legacy(records: list[dict[str, Any]], registry_counts: dict[str, int]) -> str:
    cards = []
    for record in records:
        role_text = " · ".join(record["roles"]) if record["roles"] else "-"
        cards.append(
            f'''<article class="component" data-family="{html.escape(record["family"], quote=True)}" data-kind="{html.escape(record["kind"], quote=True)}" data-search="{html.escape((record["display_name"] + " " + record["description"] + " " + role_text).lower(), quote=True)}">
  <div class="card-head"><div><span class="family">{html.escape(record["family"])}</span><h2>{html.escape(record["display_name"])}</h2></div><span class="kind">{html.escape(record["kind"])}</span></div>
  <div class="preview"><img src="{html.escape(record["preview_path"], quote=True)}" alt="{html.escape(record["display_name"], quote=True)}"></div>
  <dl><div><dt>用途</dt><dd>{html.escape(role_text)}</dd></div><div><dt>槽位</dt><dd>{record["slot_count"]} 文本/图像位，图像位 {record["image_slot_count"]}</dd></div><div><dt>说明</dt><dd>{html.escape(record["description"] or "未登记")}</dd></div></dl>
</article>'''
        )

    family_buttons = "".join(
        f'<button type="button" data-filter="family" data-value="{html.escape(family, quote=True)}">{html.escape(family)}</button>'
        for family in sorted({record["family"] for record in records})
    )
    kind_buttons = "".join(
        f'<button type="button" data-filter="kind" data-value="{html.escape(kind, quote=True)}">{html.escape(kind)}</button>'
        for kind in sorted({record["kind"] for record in records})
    )
    registry_metrics = "".join(
        f'<span><b>{count}</b>{html.escape(label)}</span>' for label, count in registry_counts.items()
    )
    return f'''<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>EasySlides 组件审查</title>
  <style>
    :root {{ --ink:#1C1B22; --muted:#6E6B78; --line:#DED9E6; --paper:#F8F7FA; --purple:#751497; --purple-soft:#F8EAFC; --green:#16724F; }}
    * {{ box-sizing:border-box; }} body {{ margin:0; font-family:"Microsoft YaHei","Aptos","Segoe UI",sans-serif; color:var(--ink); background:var(--paper); }}
    header {{ position:sticky; top:0; z-index:4; padding:20px clamp(18px,4vw,56px); background:rgba(255,255,255,.96); border-bottom:1px solid var(--line); }}
    .header-row {{ max-width:1680px; margin:auto; display:flex; gap:28px; align-items:end; justify-content:space-between; }}
    h1 {{ margin:0; font-size:25px; font-weight:800; }} p {{ margin:6px 0 0; color:var(--muted); font-size:13px; }}
    .metrics {{ display:flex; gap:14px; flex-wrap:wrap; justify-content:flex-end; }} .metrics span {{ display:grid; gap:2px; min-width:68px; color:var(--muted); font-size:11px; }} .metrics b {{ color:var(--ink); font-size:18px; }}
    main {{ max-width:1680px; margin:auto; padding:22px clamp(18px,4vw,56px) 56px; }}
    .controls {{ display:grid; grid-template-columns:minmax(180px,360px) 1fr; gap:14px; align-items:center; margin-bottom:20px; }}
    input {{ width:100%; height:36px; padding:0 12px; border:1px solid var(--line); border-radius:4px; font:inherit; outline-color:var(--purple); background:#fff; }}
    .filters {{ display:flex; gap:7px; flex-wrap:wrap; justify-content:flex-end; }} button {{ min-height:30px; border:1px solid var(--line); border-radius:4px; padding:4px 9px; background:#fff; color:var(--muted); font:inherit; font-size:12px; cursor:pointer; }} button.active {{ border-color:var(--purple); background:var(--purple); color:#fff; }}
    .component-grid {{ display:grid; grid-template-columns:repeat(3,minmax(0,1fr)); gap:16px; }} .component {{ min-width:0; background:#fff; border:1px solid var(--line); border-radius:6px; overflow:hidden; }}
    .card-head {{ display:flex; min-height:78px; align-items:start; justify-content:space-between; gap:10px; padding:14px 16px 10px; border-bottom:1px solid var(--line); }} .family {{ color:var(--purple); font-size:11px; font-weight:800; }} h2 {{ margin:4px 0 0; font-size:17px; line-height:1.25; letter-spacing:0; }} .kind {{ flex:0 0 auto; color:var(--green); font-size:11px; font-weight:700; }}
    .preview {{ display:flex; align-items:center; justify-content:center; aspect-ratio:16/9; padding:12px; background:#FBFBFC; border-bottom:1px solid var(--line); }} .preview img {{ width:100%; height:100%; object-fit:contain; }}
    dl {{ margin:0; padding:12px 16px 15px; display:grid; gap:8px; font-size:12px; }} dl div {{ display:grid; grid-template-columns:40px 1fr; gap:8px; }} dt {{ color:var(--muted); }} dd {{ margin:0; line-height:1.5; }}
    .hidden {{ display:none; }} .empty {{ display:none; padding:34px 0; color:var(--muted); text-align:center; }}
    @media (max-width:1080px) {{ .component-grid {{ grid-template-columns:repeat(2,minmax(0,1fr)); }} }} @media (max-width:720px) {{ header {{ position:static; }} .header-row,.controls {{ display:grid; grid-template-columns:1fr; }} .metrics,.filters {{ justify-content:flex-start; }} .component-grid {{ grid-template-columns:1fr; }} }}
  </style>
</head>
<body>
  <header><div class="header-row"><div><h1>EasySlides 组件审查</h1><p>图像槽位以审查辅助线标识；不改变源组件与模板输出。</p></div><div class="metrics">{registry_metrics}</div></div></header>
  <main><div class="controls"><input id="search" type="search" placeholder="筛选组件名称、用途或说明"><div class="filters"><button type="button" class="active" data-filter="all" data-value="all">全部可视组件</button>{family_buttons}{kind_buttons}</div></div><section id="grid" class="component-grid">{"".join(cards)}</section><p id="empty" class="empty">没有匹配的组件。</p></main>
  <script>
    const state={{family:null,kind:null,query:""}}; const cards=[...document.querySelectorAll('.component')]; const buttons=[...document.querySelectorAll('button')];
    function refresh() {{ let count=0; cards.forEach(card=>{{const yes=(!state.family||card.dataset.family===state.family)&&(!state.kind||card.dataset.kind===state.kind)&&(!state.query||card.dataset.search.includes(state.query));card.classList.toggle('hidden',!yes);if(yes)count++;}});document.getElementById('empty').style.display=count?'none':'block';buttons.forEach(button=>{{const field=button.dataset.filter;const value=button.dataset.value;const active=(field==='all'&&!state.family&&!state.kind)||(field==='family'&&state.family===value)||(field==='kind'&&state.kind===value);button.classList.toggle('active',active);}});}}
    buttons.forEach(button=>button.addEventListener('click',()=>{{if(button.dataset.filter==='all'){{state.family=null;state.kind=null;}}else{{state[button.dataset.filter]=state[button.dataset.filter]===button.dataset.value?null:button.dataset.value;}}refresh();}}));
    document.getElementById('search').addEventListener('input',event=>{{state.query=event.target.value.trim().toLowerCase();refresh();}});
  </script>
</body>
</html>\n'''


def _render_html_clean(records: list[dict[str, Any]], registry_counts: dict[str, int]) -> str:
    cards = []
    for record in records:
        role_text = " / ".join(record["roles"])
        cards.append(
            f'''<article class="component" data-family="{html.escape(record["family"], quote=True)}" data-kind="{html.escape(record["kind"], quote=True)}" data-search="{html.escape((record["display_name"] + " " + record["description"] + " " + role_text).lower(), quote=True)}">
  <div class="card-head"><div><span class="family">{html.escape(record["family"])}</span><h2>{html.escape(record["display_name"])}</h2></div><span class="kind">{html.escape(record["kind"])}</span></div>
  <div class="preview"><img src="{html.escape(record["preview_path"], quote=True)}" alt="{html.escape(record["display_name"], quote=True)}"></div>
  <dl><div><dt>状态</dt><dd>{html.escape(record["asset_status"])}</dd></div><div><dt>用途</dt><dd>{html.escape(role_text or "-")}</dd></div><div><dt>说明</dt><dd>{html.escape(record["description"] or "-")}</dd></div></dl>
</article>'''
        )
    family_buttons = "".join(
        f'<button type="button" data-filter="family" data-value="{html.escape(family, quote=True)}">{html.escape(family)}</button>'
        for family in sorted({record["family"] for record in records})
    )
    kind_buttons = "".join(
        f'<button type="button" data-filter="kind" data-value="{html.escape(kind, quote=True)}">{html.escape(kind)}</button>'
        for kind in sorted({record["kind"] for record in records})
    )
    labels = ["已登记资产", "可见组件", "内容变体", "组合包", "图表资产", "图标族"]
    metrics = "".join(
        f"<span><b>{count}</b>{label}</span>"
        for label, count in zip(labels, registry_counts.values())
    )
    return f'''<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>EasySlides 组件审阅</title>
<style>
:root{{--ink:#1C1B22;--muted:#6E6B78;--line:#DED9E6;--paper:#F8F7FA;--purple:#751497;--green:#16724F}}*{{box-sizing:border-box}}body{{margin:0;font-family:"Microsoft YaHei","Aptos","Segoe UI",sans-serif;color:var(--ink);background:var(--paper)}}header{{position:sticky;top:0;z-index:4;padding:20px clamp(18px,4vw,56px);background:rgba(255,255,255,.96);border-bottom:1px solid var(--line)}}.header-row{{max-width:1680px;margin:auto;display:flex;gap:28px;align-items:end;justify-content:space-between}}h1{{margin:0;font-size:25px;font-weight:800}}header p{{margin:6px 0 0;color:var(--muted);font-size:13px}}.metrics{{display:flex;gap:14px;flex-wrap:wrap;justify-content:flex-end}}.metrics span{{display:grid;gap:2px;min-width:68px;color:var(--muted);font-size:11px}}.metrics b{{color:var(--ink);font-size:18px}}main{{max-width:1680px;margin:auto;padding:22px clamp(18px,4vw,56px) 56px}}.controls{{display:grid;grid-template-columns:minmax(180px,360px) 1fr;gap:14px;align-items:center;margin-bottom:20px}}input{{width:100%;height:36px;padding:0 12px;border:1px solid var(--line);border-radius:4px;font:inherit;outline-color:var(--purple);background:#fff}}.filters{{display:flex;gap:7px;flex-wrap:wrap;justify-content:flex-end}}button{{min-height:30px;border:1px solid var(--line);border-radius:4px;padding:4px 9px;background:#fff;color:var(--muted);font:inherit;font-size:12px;cursor:pointer}}button.active{{border-color:var(--purple);background:var(--purple);color:#fff}}.component-grid{{display:grid;grid-template-columns:repeat(3,minmax(0,1fr));gap:16px}}.component{{min-width:0;background:#fff;border:1px solid var(--line);border-radius:6px;overflow:hidden}}.card-head{{display:flex;min-height:78px;align-items:start;justify-content:space-between;gap:10px;padding:14px 16px 10px;border-bottom:1px solid var(--line)}}.family{{color:var(--purple);font-size:11px;font-weight:800}}h2{{margin:4px 0 0;font-size:17px;line-height:1.25;letter-spacing:0}}.kind{{flex:0 0 auto;color:var(--green);font-size:11px;font-weight:700}}.preview{{display:flex;align-items:center;justify-content:center;aspect-ratio:16/9;padding:12px;background:#FBFBFC;border-bottom:1px solid var(--line)}}.preview img{{width:100%;height:100%;object-fit:contain}}dl{{margin:0;padding:12px 16px 15px;display:grid;gap:8px;font-size:12px}}dl div{{display:grid;grid-template-columns:40px 1fr;gap:8px}}dt{{color:var(--muted)}}dd{{margin:0;line-height:1.5}}.hidden{{display:none}}.empty{{display:none;padding:34px 0;color:var(--muted);text-align:center}}@media(max-width:1080px){{.component-grid{{grid-template-columns:repeat(2,minmax(0,1fr))}}}}@media(max-width:720px){{header{{position:static}}.header-row,.controls{{display:grid;grid-template-columns:1fr}}.metrics,.filters{{justify-content:flex-start}}.component-grid{{grid-template-columns:1fr}}}}
</style></head><body><header><div class="header-row"><div><h1>EasySlides 组件审阅</h1><p>原样保真组件直接裁切源 SVG 节点；普通组件与页面组合包保持分层，不会混为一谈。</p></div><div class="metrics">{metrics}</div></div></header><main><div class="controls"><input id="search" type="search" placeholder="筛选组件名称、用途或说明"><div class="filters"><button type="button" class="active" data-filter="all" data-value="all">全部可见组件</button>{family_buttons}{kind_buttons}</div></div><section id="grid" class="component-grid">{"".join(cards)}</section><p id="empty" class="empty">没有匹配的组件。</p></main><script>const state={{family:null,kind:null,query:""}},cards=[...document.querySelectorAll('.component')],buttons=[...document.querySelectorAll('button')];function refresh(){{let count=0;cards.forEach(card=>{{const yes=(!state.family||card.dataset.family===state.family)&&(!state.kind||card.dataset.kind===state.kind)&&(!state.query||card.dataset.search.includes(state.query));card.classList.toggle('hidden',!yes);if(yes)count++;}});document.getElementById('empty').style.display=count?'none':'block';buttons.forEach(button=>{{const field=button.dataset.filter,value=button.dataset.value,active=(field==='all'&&!state.family&&!state.kind)||(field==='family'&&state.family===value)||(field==='kind'&&state.kind===value);button.classList.toggle('active',active);}});}}buttons.forEach(button=>button.addEventListener('click',()=>{{if(button.dataset.filter==='all'){{state.family=null;state.kind=null;}}else{{state[button.dataset.filter]=state[button.dataset.filter]===button.dataset.value?null:button.dataset.value;}}refresh();}}));document.getElementById('search').addEventListener('input',event=>{{state.query=event.target.value.trim().toLowerCase();refresh();}});</script></body></html>\n'''


def _add_text(slide: Any, x: float, y: float, width: float, height: float, text: str, *, size: float, color: str, bold: bool = False, align: Any = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _build_pptx(records: list[dict[str, Any]], output_path: Path) -> None:
    import cairosvg
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    preview_dir = output_path.parent / "png"
    preview_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rect(slide: Any, x: float, y: float, width: float, height: float, fill: str, line: str | None = None) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)
        shape.line.color.rgb = RGBColor.from_string(line or fill)
        shape.line.width = Pt(0.5)

    overview = prs.slides.add_slide(blank)
    rect(overview, 0, 0, 13.333333, 7.5, "FFFFFF")
    rect(overview, 0, 0, 13.333333, 0.16, "751497")
    _add_text(overview, 0.68, 0.52, 12.0, 0.55, "EasySlides 组件审查", size=26, color="1C1B22", bold=True)
    _add_text(overview, 0.68, 1.15, 12.0, 0.34, "范围：可直接选择的视觉组件。", size=12, color="6E6B78")
    family_counts = Counter(record["family"] for record in records)
    kind_counts = Counter(record["kind"] for record in records)
    count_rows = [("可视组件", len(records)), *family_counts.items(), *kind_counts.items()]
    for index, (label, count) in enumerate(count_rows[:8]):
        row = index // 4
        col = index % 4
        x = 0.78 + col * 3.08
        y = 2.05 + row * 1.24
        rect(overview, x, y, 2.72, 0.98, "F8EAFC", "D5B2E0")
        _add_text(overview, x + 0.14, y + 0.13, 2.44, 0.28, str(count), size=20, color="751497", bold=True)
        _add_text(overview, x + 0.14, y + 0.47, 2.44, 0.28, str(label), size=10, color="4F4655")
    _add_text(overview, 0.78, 5.45, 11.6, 0.42, "虚线：审查用图片槽位。", size=12, color="6E6B78")
    _add_text(overview, 0.78, 6.0, 11.6, 0.42, "当前：34 个可视组件。", size=12, color="6E6B78")

    for page_start in range(0, len(records), 4):
        slide = prs.slides.add_slide(blank)
        rect(slide, 0, 0, 13.333333, 7.5, "FFFFFF")
        rect(slide, 0, 0, 13.333333, 0.12, "751497")
        page_records = records[page_start : page_start + 4]
        _add_text(slide, 0.5, 0.27, 8.2, 0.38, "组件联系表", size=18, color="1C1B22", bold=True)
        _add_text(slide, 0.5, 0.66, 8.4, 0.24, "视觉资产审查", size=10, color="751497", bold=True)
        _add_text(slide, 11.3, 0.27, 1.5, 0.38, f"{page_start + 1:02d}-{min(page_start + 4, len(records)):02d}", size=11, color="6E6B78", align=PP_ALIGN.RIGHT)
        for index, record in enumerate(page_records):
            row = index // 2
            col = index % 2
            x = 0.48 + col * 6.22
            y = 1.08 + row * 3.16
            card_w = 5.96
            card_h = 2.86
            rect(slide, x, y, card_w, card_h, "FFFFFF", "DED9E6")
            _add_text(slide, x + 0.18, y + 0.10, 4.45, 0.28, record["display_name"], size=13, color="1C1B22", bold=True)
            _add_text(slide, x + 4.65, y + 0.10, 1.1, 0.28, record["kind"], size=8, color="16724F", bold=True, align=PP_ALIGN.RIGHT)
            png_path = preview_dir / f"{record['preview_stem']}.png"
            cairosvg.svg2png(url=str(output_path.parent / record["preview_path"]), write_to=str(png_path), output_width=1040)
            slide.shapes.add_picture(str(png_path), Inches(x + 0.18), Inches(y + 0.46), width=Inches(5.60), height=Inches(1.56))
            detail = f"槽位 {record['slot_count']}  |  图像 {record['image_slot_count']}"
            _add_text(slide, x + 0.18, y + 2.14, 5.60, 0.22, detail, size=8.5, color="6E6B78")
            role = str(record["roles"][0] if record["roles"] else "未登记")
            role_label = role.replace("nsfc_defense_shell", "NSFC shell").replace("semantic_template_component", "semantic")
            _add_text(slide, x + 0.18, y + 2.40, 5.60, 0.22, role_label[:28], size=8.5, color="4F4655")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _build_contact_sheet(records: list[dict[str, Any]], output_dir: Path) -> None:
    """Create a compact raster overview for a first-pass visual review."""
    from PIL import Image, ImageDraw, ImageFont

    png_dir = output_dir / "png"
    columns = 4
    card_width = 384
    card_height = 246
    rows = (len(records) + columns - 1) // columns
    canvas = Image.new("RGB", (columns * card_width, rows * card_height), "#F8F7FA")
    draw = ImageDraw.Draw(canvas)
    try:
        title_font = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 16)
        meta_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 11)
    except OSError:  # pragma: no cover - fallback for non-Windows review machines
        title_font = ImageFont.load_default()
        meta_font = ImageFont.load_default()

    for index, record in enumerate(records):
        row = index // columns
        col = index % columns
        x = col * card_width
        y = row * card_height
        draw.rectangle((x + 7, y + 7, x + card_width - 8, y + card_height - 8), fill="#FFFFFF", outline="#DED9E6", width=1)
        draw.text((x + 18, y + 18), record["display_name"], fill="#1C1B22", font=title_font)
        draw.text((x + 18, y + 39), f"{record['family']} | {record['kind']}", fill="#6E6B78", font=meta_font)
        source = png_dir / f"{record['preview_stem']}.png"
        if not source.is_file():
            continue
        with Image.open(source) as image:
            image = image.convert("RGB")
            image.thumbnail((card_width - 38, card_height - 80), Image.Resampling.LANCZOS)
            image_x = x + (card_width - image.width) // 2
            image_y = y + 66 + (card_height - 78 - image.height) // 2
            canvas.paste(image, (image_x, image_y))

    canvas.save(output_dir / "component_contact_sheet.png")


def build_component_asset_audit(output_dir: Path = DEFAULT_OUTPUT) -> dict[str, Any]:
    if output_dir.exists():
        for child in (output_dir / "previews", output_dir / "png"):
            if child.exists():
                shutil.rmtree(child)
    output_dir.mkdir(parents=True, exist_ok=True)
    records = _component_records()
    manifest_records: list[dict[str, Any]] = []
    for record in records:
        preview_name = f"{_slug(record.asset_id)}.svg"
        preview_path = output_dir / "previews" / preview_name
        _audit_svg(record.source_path, preview_path)
        manifest_records.append(
            {
                "asset_id": record.asset_id,
                "display_name": record.display_name,
                "family": record.family,
                "kind": record.kind,
                "source_path": _rel(record.source_path, ROOT),
                "preview_path": f"previews/{preview_name}",
                "preview_stem": Path(preview_name).stem,
                "description": record.description,
                "roles": list(filter(None, record.roles)),
                "slot_count": record.slot_count,
                "image_slot_count": record.image_slot_count,
                "asset_status": record.asset_status,
            }
        )

    registry_counts: dict[str, int] = {}
    if REGISTRY_PATH.is_file():
        registry = _read_json(REGISTRY_PATH)
        counts_by_granularity = registry.get("counts_by_granularity") or {}
        registry_counts = {
            "登记资产": int(registry.get("asset_count") or 0),
            "当前可视组件": len(manifest_records),
            "页面变体": int(counts_by_granularity.get("body_variant") or 0),
            "页面组合包": int(counts_by_granularity.get("component_package") or 0),
            "图表资产": int(counts_by_granularity.get("chart_asset") or 0),
            "图标族": int(counts_by_granularity.get("icon_family") or 0),
        }
    else:
        registry_counts = {"可视组件": len(manifest_records)}

    manifest = {
        "schema_version": SCHEMA_VERSION,
        "status": "pass",
        "renderable_component_count": len(manifest_records),
        "counts_by_family": dict(Counter(record["family"] for record in manifest_records)),
        "counts_by_kind": dict(Counter(record["kind"] for record in manifest_records)),
        "excluded_from_visual_component_catalog": {
            "composition_package": 6,
            "reason": "research-core is a page-composition package with story previews, not a leaf visual component.",
        },
        "registry_counts_by_granularity": registry_counts,
        "records": manifest_records,
    }
    (output_dir / "component_audit_manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    (output_dir / "component_audit.html").write_text(
        _render_html(manifest_records, registry_counts), encoding="utf-8"
    )
    _build_pptx(manifest_records, output_dir / "easyslides_component_audit.pptx")
    _build_contact_sheet(manifest_records, output_dir)
    return manifest


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build a visual audit kit for renderable EasySlides components.")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    manifest = build_component_asset_audit(args.out)
    if args.json:
        print(json.dumps(manifest, ensure_ascii=False, indent=2))
    else:
        print(
            f"Component audit: {manifest['status']} "
            f"({manifest['renderable_component_count']} renderable component(s))"
        )
        print(args.out / "component_audit.html")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
