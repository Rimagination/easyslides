#!/usr/bin/env python3
"""Validate structural editability of a PPTX rebuilt from slide images."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from scripts import layout_metrics
except ImportError:  # pragma: no cover - direct script execution
    import layout_metrics


SCHEMA_VERSION = "easyslides.image_reconstruction_pptx_report.v1"
FULL_SLIDE_PICTURE_THRESHOLD = 0.85


def _inches(value: int | float | None) -> float:
    return layout_metrics.emu_to_in(value)


def _issue(
    code: str,
    severity: str,
    message: str,
    *,
    slide_number: int,
    shape_name: str | None = None,
    suggestion: str | None = None,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "code": code,
        "severity": severity,
        "slide_number": slide_number,
        "message": message,
    }
    if shape_name:
        payload["shape_name"] = shape_name
    if suggestion:
        payload["suggestion"] = suggestion
    return payload


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _shape_area_fraction(shape: Any, slide_area: float) -> float:
    if slide_area <= 0:
        return 0.0
    return max(0.0, _inches(shape.width)) * max(0.0, _inches(shape.height)) / slide_area


def _has_visible_text(shape: Any) -> bool:
    return bool(getattr(shape, "has_text_frame", False) and shape.text_frame and shape.text_frame.text.strip())


def validate_image_reconstruction_pptx(
    pptx_path: str | Path,
    *,
    full_slide_picture_threshold: float = FULL_SLIDE_PICTURE_THRESHOLD,
) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_w = _inches(prs.slide_width)
    slide_h = _inches(prs.slide_height)
    slide_area = slide_w * slide_h
    issues: list[dict[str, Any]] = []
    slides: list[dict[str, Any]] = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        text_frame_count = 0
        native_shape_count = 0
        picture_count = 0
        max_picture_area_fraction = 0.0

        for shape in _iter_shapes(slide.shapes):
            shape_type = getattr(shape, "shape_type", None)
            if _has_visible_text(shape):
                text_frame_count += 1
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                area_fraction = _shape_area_fraction(shape, slide_area)
                max_picture_area_fraction = max(max_picture_area_fraction, area_fraction)
                if area_fraction >= full_slide_picture_threshold:
                    issues.append(
                        _issue(
                            "PPTX-FULL-SLIDE-PICTURE",
                            "blocking",
                            f"Picture covers {area_fraction:.1%} of the slide, which likely means a full-slide screenshot was used.",
                            slide_number=slide_number,
                            shape_name=getattr(shape, "name", None),
                            suggestion="Split the source into visual assets, native structure, and editable text instead of placing one large image.",
                        )
                    )
            elif shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE}:
                native_shape_count += 1

        if text_frame_count == 0:
            issues.append(
                _issue(
                    "PPTX-NO-EDITABLE-TEXT",
                    "warning",
                    "Slide has no editable text frames.",
                    slide_number=slide_number,
                    suggestion="If the source image contains readable text, extract it into native PowerPoint text boxes.",
                )
            )
        if native_shape_count == 0:
            issues.append(
                _issue(
                    "PPTX-NO-NATIVE-STRUCTURE",
                    "warning",
                    "Slide has no native structure shapes, tables, charts, or lines.",
                    slide_number=slide_number,
                    suggestion="Simple panels, arrows, dividers, and badges should be rebuilt as native DrawingML/PPT shapes.",
                )
            )
        if picture_count == 1 and text_frame_count == 0 and native_shape_count == 0:
            issues.append(
                _issue(
                    "PPTX-SINGLE-PICTURE-ONLY",
                    "blocking",
                    "Slide appears to contain only one picture and no editable text or native structure.",
                    slide_number=slide_number,
                )
            )

        slides.append(
            {
                "slide_number": slide_number,
                "text_frame_count": text_frame_count,
                "native_shape_count": native_shape_count,
                "picture_count": picture_count,
                "max_picture_area_fraction": round(max_picture_area_fraction, 4),
            }
        )

    blocking_count = sum(1 for issue in issues if issue["severity"] == "blocking")
    warning_count = sum(1 for issue in issues if issue["severity"] == "warning")
    return {
        "schema_version": SCHEMA_VERSION,
        "status": "fail" if blocking_count else "pass",
        "pptx_path": str(Path(pptx_path).resolve()),
        "slide_count": len(prs.slides),
        "blocking_count": blocking_count,
        "warning_count": warning_count,
        "thresholds": {"full_slide_picture_area_fraction": full_slide_picture_threshold},
        "slides": slides,
        "issues": issues,
    }


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Validate PPTX structural editability after slide-image reconstruction.")
    parser.add_argument("pptx", help="PPTX file to validate.")
    parser.add_argument("--report", help="Optional JSON report path.")
    parser.add_argument("--full-slide-picture-threshold", type=float, default=FULL_SLIDE_PICTURE_THRESHOLD)
    parser.add_argument("--quiet", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    report = validate_image_reconstruction_pptx(
        args.pptx,
        full_slide_picture_threshold=args.full_slide_picture_threshold,
    )
    if args.report:
        _write_json(Path(args.report), report)
    if not args.quiet:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["status"] == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
