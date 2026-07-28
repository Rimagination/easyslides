#!/usr/bin/env python3
"""Render EasySlides component package stories as native PPTX previews."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import sys
import textwrap
from typing import Any


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

try:
    from scripts.component_package import (
        INSTALLED_PACKAGES_ROOT,
        PACKAGES_ROOT,
        STORY_SCHEMA_VERSION,
        is_public_component_package,
        load_component_packages_from_roots,
        validate_component_story_payload,
    )
except ModuleNotFoundError:  # pragma: no cover
    from component_package import (
        INSTALLED_PACKAGES_ROOT,
        PACKAGES_ROOT,
        STORY_SCHEMA_VERSION,
        is_public_component_package,
        load_component_packages_from_roots,
        validate_component_story_payload,
    )

try:
    from scripts.component_renderer_registry import register_renderer_handler, render_registered
except ModuleNotFoundError:  # pragma: no cover
    from component_renderer_registry import register_renderer_handler, render_registered


DEFAULT_OUTPUT = ROOT / "templates" / "components" / "gallery" / "component_gallery.pptx"
SCHEMA_VERSION = "easyslides.component_pptx_renderer_report.v1"
CANVAS_W = 1280
CANVAS_H = 720


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8-sig"))
    if not isinstance(payload, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return payload


def _story_payload(package_dir: Path, story_ref: dict[str, Any]) -> dict[str, Any]:
    story_path = package_dir / str(story_ref["payload"])
    story = _read_json(story_path)
    if story.get("schema_version") != STORY_SCHEMA_VERSION:
        raise ValueError(f"{story_path} must use {STORY_SCHEMA_VERSION}")
    payload = story.get("payload")
    if not isinstance(payload, dict):
        raise ValueError(f"{story_path} payload must be an object")
    return payload


def _px(value: float) -> float:
    return value / 96.0


def _clip(value: Any, limit: int = 90) -> str:
    text = " ".join(str(value or "").split())
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 3)].rstrip() + "..."


def _wrap(value: Any, width: int = 22, max_lines: int = 4) -> str:
    text = _clip(value, width * max_lines)
    lines = textwrap.wrap(text, width=width, break_long_words=False, replace_whitespace=False)
    if not lines:
        return ""
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = lines[-1].rstrip(". ") + "..."
    return "\n".join(lines)


def _rgb(hex_value: str) -> tuple[int, int, int]:
    value = hex_value.strip().lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def build_component_pptx(
    *,
    packages_root: Path = PACKAGES_ROOT,
    installed_root: Path | None = INSTALLED_PACKAGES_ROOT,
    output_path: Path = DEFAULT_OUTPUT,
    validate_text_layout: bool = False,
) -> dict[str, Any]:
    """Build a native PPTX preview deck from component package stories."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    package_roots = [packages_root]
    if installed_root and Path(installed_root).resolve() != Path(packages_root).resolve():
        package_roots.append(Path(installed_root))
    public_packages = [
        (package_dir, package)
        for package_dir, package in load_component_packages_from_roots(package_roots)
        if is_public_component_package(package_dir)
    ]
    if not public_packages:
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "not_applicable",
            "output": "",
            "package_count": 0,
            "story_count": 0,
            "slide_count": 0,
            "center_anchor_textbox_count": 0,
            "text_layout_status": "not_applicable",
            "text_layout_report": None,
            "packages": [],
            "reason": "no_public_component_packages",
        }

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    output_path.parent.mkdir(parents=True, exist_ok=True)
    stats = {"center_anchor_textbox_count": 0}

    def add_rect(slide, x, y, w, h, fill, line="#D7DDE5", radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(_px(x)),
            Inches(_px(y)),
            Inches(_px(w)),
            Inches(_px(h)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*_rgb(fill))
        shape.line.color.rgb = RGBColor(*_rgb(line if line else fill))
        shape.line.width = Pt(1)
        return shape

    def add_textbox(
        slide,
        x,
        y,
        w,
        h,
        text,
        *,
        size=16,
        bold=False,
        color="#172033",
        align=PP_ALIGN.CENTER,
        wrap_width=22,
        max_lines=4,
    ):
        box = slide.shapes.add_textbox(Inches(_px(x)), Inches(_px(y)), Inches(_px(w)), Inches(_px(h)))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = Inches(0.03)
        frame.margin_right = Inches(0.03)
        frame.margin_top = Inches(0)
        frame.margin_bottom = Inches(0)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = _wrap(text, width=wrap_width, max_lines=max_lines)
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*_rgb(color))
        stats["center_anchor_textbox_count"] += 1
        return box

    def add_line(slide, x1, y1, x2, y2, color="#9CB6C9", width=2):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(_px(x1)),
            Inches(_px(y1)),
            Inches(_px(x2)),
            Inches(_px(y2)),
        )
        line.line.color.rgb = RGBColor(*_rgb(color))
        line.line.width = Pt(width)
        return line

    def add_shell(slide, component_id: str, story_id: str, status: str):
        add_rect(slide, 38, 30, 1204, 660, "#FFFFFF", "#D7DDE5", radius=True)
        add_rect(slide, 38, 30, 1204, 104, "#172033", "#172033", radius=False)
        add_textbox(
            slide,
            70,
            47,
            820,
            34,
            component_id,
            size=20,
            bold=True,
            color="#FFFFFF",
            align=PP_ALIGN.LEFT,
            wrap_width=40,
            max_lines=1,
        )
        add_textbox(
            slide,
            70,
            82,
            780,
            24,
            f"story: {story_id}",
            size=10,
            bold=True,
            color="#B7C3D3",
            align=PP_ALIGN.LEFT,
            wrap_width=90,
            max_lines=1,
        )
        badge_fill = "#0F7B55" if status == "pass" else "#B42318"
        add_rect(slide, 1040, 58, 128, 34, badge_fill, badge_fill, radius=True)
        add_textbox(
            slide,
            1040,
            58,
            128,
            34,
            status.upper(),
            size=11,
            bold=True,
            color="#FFFFFF",
            wrap_width=12,
            max_lines=1,
        )

    def render_three_card(slide, payload: dict[str, Any]):
        items = payload.get("items") if isinstance(payload.get("items"), list) else []
        boxes = [(74, 166, 348, 430), (466, 166, 348, 430), (858, 166, 348, 430)]
        for index, (x, y, w, h) in enumerate(boxes, start=1):
            item = items[index - 1] if index - 1 < len(items) and isinstance(items[index - 1], dict) else {}
            add_rect(slide, x, y, w, h, "#FDFEFE", "#CDD6E0", radius=True)
            add_rect(slide, x, y, w, 10, "#1C75BC", "#1C75BC", radius=False)
            add_rect(slide, x + 20, y + 28, 44, 44, "#EAF3FA", "#C9DFF0", radius=True)
            add_textbox(slide, x + 20, y + 28, 44, 44, f"{index:02d}", size=12, bold=True, color="#1C75BC", max_lines=1)
            add_textbox(
                slide,
                x + 34,
                y + 104,
                w - 68,
                70,
                item.get("title", ""),
                size=18,
                bold=True,
                wrap_width=14,
                max_lines=2,
            )
            add_line(slide, x + 34, y + 196, x + w - 34, y + 196, "#D7DDE5", width=1)
            add_textbox(
                slide,
                x + 34,
                y + 220,
                w - 68,
                138,
                item.get("body", ""),
                size=13,
                color="#4B5B6D",
                wrap_width=25,
                max_lines=5,
            )

    def render_process(slide, payload: dict[str, Any]):
        items = payload.get("items") if isinstance(payload.get("items"), list) else []
        count = max(3, min(5, len(items) or 4))
        gap = 24
        w = (1088 - gap * (count - 1)) / count
        y = 238
        add_line(slide, 104, 318, 1176, 318, "#9CB6C9", width=3)
        for index in range(count):
            x = 96 + index * (w + gap)
            item = items[index] if index < len(items) and isinstance(items[index], dict) else {}
            add_rect(slide, x, y, w, 250, "#FFFFFF", "#CDD6E0", radius=True)
            add_rect(slide, x + 18, y - 30, 58, 58, "#145F8F", "#145F8F", radius=True)
            add_textbox(slide, x + 18, y - 30, 58, 58, str(index + 1), size=15, bold=True, color="#FFFFFF", max_lines=1)
            add_textbox(
                slide,
                x + 20,
                y + 54,
                w - 40,
                56,
                item.get("title", ""),
                size=14,
                bold=True,
                wrap_width=12,
                max_lines=2,
            )
            add_textbox(
                slide,
                x + 20,
                y + 134,
                w - 40,
                84,
                item.get("body", ""),
                size=11,
                color="#4B5B6D",
                wrap_width=18,
                max_lines=4,
            )

    def render_figure(slide, payload: dict[str, Any]):
        add_rect(slide, 82, 156, 690, 452, "#E9EEF3", "#C4CDD8", radius=True)
        add_rect(slide, 122, 198, 610, 322, "#F8FAFC", "#D7DDE5", radius=False)
        add_textbox(slide, 260, 330, 334, 50, "FIGURE", size=22, bold=True, color="#8290A2", max_lines=1)
        add_textbox(
            slide,
            230,
            384,
            394,
            32,
            "preserve aspect ratio | source-linked",
            size=10,
            bold=True,
            color="#8290A2",
            max_lines=1,
        )
        add_rect(slide, 824, 156, 374, 452, "#FFFFFF", "#CDD6E0", radius=True)
        add_rect(slide, 824, 156, 374, 10, "#1C75BC", "#1C75BC", radius=False)
        add_textbox(
            slide,
            858,
            206,
            306,
            84,
            payload.get("takeaway", ""),
            size=18,
            bold=True,
            wrap_width=16,
            max_lines=2,
        )
        add_line(slide, 858, 318, 1164, 318, "#D7DDE5", width=1)
        add_textbox(
            slide,
            858,
            350,
            306,
            148,
            payload.get("bullets", ""),
            size=12,
            color="#4B5B6D",
            align=PP_ALIGN.LEFT,
            wrap_width=26,
            max_lines=5,
        )
        add_textbox(
            slide,
            858,
            542,
            306,
            34,
            payload.get("caption") or payload.get("source") or payload.get("image") or "",
            size=9,
            bold=True,
            color="#69788A",
            align=PP_ALIGN.LEFT,
            wrap_width=42,
            max_lines=2,
        )

    def render_kpi_row(slide, payload: dict[str, Any]):
        items = payload.get("items") if isinstance(payload.get("items"), list) else []
        boxes = [(94, 246, 340, 190), (470, 246, 340, 190), (846, 246, 340, 190)]
        for index, (x, y, w, h) in enumerate(boxes):
            item = items[index] if index < len(items) and isinstance(items[index], dict) else {}
            add_rect(slide, x, y, w, h, "#FFFFFF", "#CDD6E0", radius=True)
            add_rect(slide, x, y, w, 9, "#0F766E", "#0F766E", radius=False)
            add_textbox(
                slide,
                x + 34,
                y + 28,
                w - 68,
                58,
                item.get("metric", ""),
                size=25,
                bold=True,
                color="#0F766E",
                wrap_width=8,
                max_lines=1,
            )
            add_textbox(
                slide,
                x + 34,
                y + 92,
                w - 68,
                34,
                item.get("label", ""),
                size=14,
                bold=True,
                wrap_width=14,
                max_lines=1,
            )
            add_textbox(
                slide,
                x + 34,
                y + 134,
                w - 68,
                42,
                item.get("note", ""),
                size=11,
                color="#617083",
                wrap_width=18,
                max_lines=2,
            )

    def render_comparison(slide, payload: dict[str, Any]):
        items = payload.get("items") if isinstance(payload.get("items"), list) else []
        boxes = [(82, 184, 520, 372), (678, 184, 520, 372)]
        accents = ["#1C75BC", "#0F766E"]
        for index, (x, y, w, h) in enumerate(boxes):
            item = items[index] if index < len(items) and isinstance(items[index], dict) else {}
            add_rect(slide, x, y, w, h, "#FFFFFF", "#CDD6E0", radius=True)
            add_rect(slide, x, y, w, 11, accents[index], accents[index], radius=False)
            add_textbox(
                slide,
                x + 38,
                y + 58,
                w - 76,
                52,
                item.get("title", ""),
                size=18,
                bold=True,
                wrap_width=18,
                max_lines=1,
            )
            add_line(slide, x + 38, y + 136, x + w - 38, y + 136, "#D7DDE5", width=1)
            add_textbox(
                slide,
                x + 38,
                y + 170,
                w - 76,
                128,
                item.get("body", ""),
                size=12,
                color="#4B5B6D",
                wrap_width=34,
                max_lines=5,
            )
        add_rect(slide, 202, 584, 876, 54, "#EFF6FB", "#C9DFF0", radius=True)
        add_textbox(
            slide,
            242,
            592,
            796,
            38,
            payload.get("synthesis", ""),
            size=11,
            bold=True,
            color="#145F8F",
            wrap_width=34,
            max_lines=2,
        )

    def render_evidence_stack(slide, payload: dict[str, Any]):
        items = payload.get("items") if isinstance(payload.get("items"), list) else []
        count = max(3, min(5, len(items) or 3))
        row_gap = 14
        row_h = (306 - row_gap * (count - 1)) / count
        add_rect(slide, 116, 162, 1048, 86, "#172033", "#172033", radius=True)
        add_textbox(
            slide,
            156,
            176,
            968,
            58,
            payload.get("claim", ""),
            size=18,
            bold=True,
            color="#FFFFFF",
            wrap_width=42,
            max_lines=2,
        )
        for index in range(count):
            y = 286 + index * (row_h + row_gap)
            item = items[index] if index < len(items) and isinstance(items[index], dict) else {}
            add_rect(slide, 116, y, 1048, row_h, "#FFFFFF", "#CDD6E0", radius=True)
            add_rect(slide, 138, y + row_h / 2 - 18, 36, 36, "#EAF3FA", "#C9DFF0", radius=True)
            add_textbox(
                slide,
                138,
                y + row_h / 2 - 18,
                36,
                36,
                str(index + 1),
                size=11,
                bold=True,
                color="#1C75BC",
                max_lines=1,
            )
            add_textbox(
                slide,
                198,
                y + 2,
                924,
                row_h - 4,
                item.get("evidence", ""),
                size=11,
                bold=True,
                color="#4B5B6D",
                wrap_width=30,
                max_lines=2,
            )

    register_renderer_handler("three_card_summary", "native_pptx", render_three_card)
    register_renderer_handler("process_timeline", "native_pptx", render_process)
    register_renderer_handler("figure_with_notes", "native_pptx", render_figure)
    register_renderer_handler("kpi_row_3", "native_pptx", render_kpi_row)
    register_renderer_handler("comparison_pair", "native_pptx", render_comparison)
    register_renderer_handler("evidence_stack", "native_pptx", render_evidence_stack)

    story_count = 0
    package_rows: list[dict[str, Any]] = []
    for package_dir, package in public_packages:
        component_id = str(package.get("component_id") or package_dir.name)
        renderer_id = str(package.get("renderer_id") or component_id)
        stories = []
        for story_ref in package.get("stories", []):
            if not isinstance(story_ref, dict):
                continue
            story_id = str(story_ref.get("story_id") or "story")
            payload = _story_payload(package_dir, story_ref)
            payload_report = validate_component_story_payload(str(package.get("source_asset_id") or ""), payload)
            status = "pass" if payload_report["passed"] else "fail"
            slide = prs.slides.add_slide(blank)
            add_shell(slide, component_id, story_id, status)
            try:
                render_registered("native_pptx", renderer_id, slide, payload)
            except (KeyError, ValueError):
                add_textbox(slide, 300, 330, 680, 70, f"No native renderer for {component_id}", size=20, bold=True)
            stories.append({"story_id": story_id, "status": status})
            story_count += 1
        package_rows.append({"component_id": component_id, "story_count": len(stories), "stories": stories})

    prs.save(output_path)

    text_layout_report = None
    if validate_text_layout:
        try:
            from scripts.validate_pptx_text_layout import validate_pptx_text_layout
        except ModuleNotFoundError:  # pragma: no cover
            from validate_pptx_text_layout import validate_pptx_text_layout

        text_layout_report = validate_pptx_text_layout(output_path)

    status = "pass"
    if text_layout_report and text_layout_report.get("status") == "fail":
        status = "fail"
    return {
        "schema_version": SCHEMA_VERSION,
        "status": status,
        "output": str(output_path),
        "package_count": len(package_rows),
        "story_count": story_count,
        "slide_count": len(prs.slides),
        "center_anchor_textbox_count": stats["center_anchor_textbox_count"],
        "text_layout_status": text_layout_report.get("status") if isinstance(text_layout_report, dict) else "skipped",
        "text_layout_report": text_layout_report,
        "packages": package_rows,
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Render EasySlides component package stories to a native PPTX preview deck.")
    parser.add_argument("--packages-root", type=Path, default=PACKAGES_ROOT)
    parser.add_argument("--installed-root", type=Path, default=INSTALLED_PACKAGES_ROOT)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--validate-text-layout", action="store_true")
    parser.add_argument("--json", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    report = build_component_pptx(
        packages_root=args.packages_root,
        installed_root=args.installed_root,
        output_path=args.out,
        validate_text_layout=args.validate_text_layout,
    )
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"Component PPTX renderer: {report['status']} ({report['slide_count']} slide(s))")
        print(report["output"])
        text_layout_report = report.get("text_layout_report")
        if isinstance(text_layout_report, dict) and text_layout_report.get("status") == "fail":
            for item in text_layout_report.get("issues", []):
                print(f"- {item.get('code')}: {item.get('message')}")
    return 0 if report["status"] in {"pass", "not_applicable"} else 1


if __name__ == "__main__":
    raise SystemExit(main())
