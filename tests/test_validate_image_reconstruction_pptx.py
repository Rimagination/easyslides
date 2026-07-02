import base64
import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)


def write_png(path: Path) -> None:
    path.write_bytes(PNG_1X1)


class ValidateImageReconstructionPptxTests(unittest.TestCase):
    def test_blocks_single_full_slide_picture(self):
        from scripts.validate_image_reconstruction_pptx import validate_image_reconstruction_pptx

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image = root / "pixel.png"
            pptx = root / "single_picture.pptx"
            write_png(image)

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(pptx)

            report = validate_image_reconstruction_pptx(pptx)

        self.assertEqual(report["status"], "fail")
        codes = {issue["code"] for issue in report["issues"]}
        self.assertIn("PPTX-FULL-SLIDE-PICTURE", codes)
        self.assertIn("PPTX-SINGLE-PICTURE-ONLY", codes)

    def test_text_and_native_shape_pass_without_blocking_issues(self):
        from scripts.validate_image_reconstruction_pptx import validate_image_reconstruction_pptx

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image = root / "pixel.png"
            pptx = root / "editable.pptx"
            write_png(image)

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
            run = textbox.text_frame.paragraphs[0].add_run()
            run.text = "Editable title"
            run.font.size = Pt(24)
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(3), Inches(2))
            slide.shapes.add_picture(str(image), Inches(4), Inches(1), width=Inches(1), height=Inches(1))
            prs.save(pptx)

            report = validate_image_reconstruction_pptx(pptx)

        self.assertEqual(report["status"], "pass")
        self.assertEqual(report["blocking_count"], 0)
        self.assertEqual(report["slides"][0]["text_frame_count"], 1)
        self.assertEqual(report["slides"][0]["native_shape_count"], 1)

    def test_cli_help_is_printable(self):
        result = subprocess.run(
            [sys.executable, "scripts/validate_image_reconstruction_pptx.py", "--help"],
            cwd=ROOT,
            text=True,
            capture_output=True,
            encoding="utf-8",
        )

        self.assertEqual(result.returncode, 0, result.stdout + result.stderr)
        self.assertIn("Validate PPTX structural editability", result.stdout)

    def test_cli_writes_report(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image = root / "pixel.png"
            pptx = root / "single_picture.pptx"
            report_path = root / "report.json"
            write_png(image)

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(pptx)

            result = subprocess.run(
                [
                    sys.executable,
                    "scripts/validate_image_reconstruction_pptx.py",
                    str(pptx),
                    "--report",
                    str(report_path),
                    "--quiet",
                ],
                cwd=ROOT,
                text=True,
                capture_output=True,
                encoding="utf-8",
            )

            self.assertEqual(result.returncode, 1, result.stdout + result.stderr)
            report = json.loads(report_path.read_text(encoding="utf-8"))

        self.assertEqual(report["schema_version"], "easyslides.image_reconstruction_pptx_report.v1")


if __name__ == "__main__":
    unittest.main()
