import tempfile
import unittest
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def make_template_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(0.7))
    title.name = "Title Slot"
    title.text = "Original Title"

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(7), Inches(2))
    body.name = "Body Slot"
    body.text = "Original body text with enough words to be treated as body content."

    prs.save(path)


def make_two_slide_template_pptx(path: Path) -> None:
    prs = Presentation()
    for idx, title_text in enumerate(("Alpha", "Beta"), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(0.7))
        title.name = f"Title Slot {idx}"
        title.text = title_text
    prs.save(path)


def make_chart_template_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.5))
    title.text = "Chart Slide"

    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Old Series", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.2),
        Inches(7),
        Inches(4),
        chart_data,
    )
    prs.save(path)


def inject_private_tag_part(path: Path) -> None:
    from scripts.template_fill_pptx import PACKAGE_REL_NS

    with zipfile.ZipFile(path, "r") as zf:
        entries = {info.filename: zf.read(info.filename) for info in zf.infolist() if not info.is_dir()}
    rels_name = "ppt/slides/_rels/slide1.xml.rels"
    rels = ET.fromstring(entries[rels_name])
    ET.SubElement(
        rels,
        f"{{{PACKAGE_REL_NS}}}Relationship",
        {
            "Id": "rId99",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags",
            "Target": "../tags/tag1.xml",
        },
    )
    entries[rels_name] = ET.tostring(rels, encoding="utf-8", xml_declaration=True)
    entries["ppt/tags/tag1.xml"] = b'<?xml version="1.0" encoding="UTF-8"?><tags><tag>one</tag></tags>'
    content_types = ET.fromstring(entries["[Content_Types].xml"])
    ET.SubElement(
        content_types,
        "{http://schemas.openxmlformats.org/package/2006/content-types}Override",
        {
            "PartName": "/ppt/tags/tag1.xml",
            "ContentType": "application/vnd.easyslides.test.tags+xml",
        },
    )
    entries["[Content_Types].xml"] = ET.tostring(content_types, encoding="utf-8", xml_declaration=True)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in sorted(entries.items()):
            zf.writestr(name, data)


def chart_parts_by_slide(pptx_path: Path) -> list[str]:
    from scripts.template_fill_pptx import _normalize_part, _qn
    from scripts.template_fill_pptx import CHART_REL_TYPE, PACKAGE_REL_NS, PRESENTATION_NS, REL_NS, SLIDE_REL_TYPE

    with zipfile.ZipFile(pptx_path) as zf:
        pres = ET.fromstring(zf.read("ppt/presentation.xml"))
        pres_rels = ET.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
        rels = {
            rel.attrib["Id"]: rel
            for rel in pres_rels.findall(_qn(PACKAGE_REL_NS, "Relationship"))
        }
        chart_parts: list[str] = []
        for slide_id in pres.findall(f".//{_qn(PRESENTATION_NS, 'sldId')}"):
            rel = rels[slide_id.attrib[_qn(REL_NS, "id")]]
            if rel.attrib["Type"] != SLIDE_REL_TYPE:
                continue
            slide_part = _normalize_part(rel.attrib["Target"])
            slide_rels_name = str(Path(slide_part).parent / "_rels" / f"{Path(slide_part).name}.rels").replace("\\", "/")
            slide_rels = ET.fromstring(zf.read(slide_rels_name))
            for item in slide_rels.findall(_qn(PACKAGE_REL_NS, "Relationship")):
                if item.attrib.get("Type") == CHART_REL_TYPE:
                    chart_parts.append(_normalize_part(item.attrib["Target"], slide_part))
                    break
    return chart_parts


def workbook_parts_by_chart(pptx_path: Path, chart_parts: list[str]) -> list[str]:
    from scripts.template_fill_pptx import _normalize_part, _qn
    from scripts.template_fill_pptx import PACKAGE_REL_NS, PACKAGE_REL_TYPE

    workbook_parts: list[str] = []
    with zipfile.ZipFile(pptx_path) as zf:
        for chart_part in chart_parts:
            rels_name = str(Path(chart_part).parent / "_rels" / f"{Path(chart_part).name}.rels").replace("\\", "/")
            rels = ET.fromstring(zf.read(rels_name))
            for rel in rels.findall(_qn(PACKAGE_REL_NS, "Relationship")):
                if rel.attrib.get("Type") == PACKAGE_REL_TYPE:
                    workbook_parts.append(_normalize_part(rel.attrib["Target"], chart_part))
                    break
    return workbook_parts


def workbook_matrix(pptx_path: Path, workbook_part: str) -> list[list[object]]:
    from openpyxl import load_workbook

    with zipfile.ZipFile(pptx_path) as zf:
        data = zf.read(workbook_part)
    workbook = load_workbook(BytesIO(data), data_only=True)
    sheet = workbook.active
    return [
        [sheet.cell(row=row, column=col).value for col in range(1, 3)]
        for row in range(1, 4)
    ]


def private_tag_targets_by_slide(pptx_path: Path) -> list[str]:
    from scripts.template_fill_pptx import _normalize_part, _qn
    from scripts.template_fill_pptx import PACKAGE_REL_NS, PRESENTATION_NS, REL_NS, SLIDE_REL_TYPE

    targets: list[str] = []
    tag_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags"
    with zipfile.ZipFile(pptx_path) as zf:
        pres = ET.fromstring(zf.read("ppt/presentation.xml"))
        pres_rels = ET.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
        rels = {
            rel.attrib["Id"]: rel
            for rel in pres_rels.findall(_qn(PACKAGE_REL_NS, "Relationship"))
        }
        for slide_id in pres.findall(f".//{_qn(PRESENTATION_NS, 'sldId')}"):
            rel = rels[slide_id.attrib[_qn(REL_NS, "id")]]
            if rel.attrib["Type"] != SLIDE_REL_TYPE:
                continue
            slide_part = _normalize_part(rel.attrib["Target"])
            slide_rels_name = str(Path(slide_part).parent / "_rels" / f"{Path(slide_part).name}.rels").replace("\\", "/")
            slide_rels = ET.fromstring(zf.read(slide_rels_name))
            for item in slide_rels.findall(_qn(PACKAGE_REL_NS, "Relationship")):
                if item.attrib.get("Type") == tag_rel_type:
                    targets.append(_normalize_part(item.attrib["Target"], slide_part))
                    break
    return targets


class TemplateFillPptxTests(unittest.TestCase):
    def test_analyze_scaffold_validate_and_apply_text_replacements(self):
        from scripts.template_fill_pptx import analyze_pptx, apply_plan, scaffold_plan, validate_plan

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "template.pptx"
            output = tmp_path / "filled.pptx"
            make_template_pptx(source)

            library = analyze_pptx(source)
            self.assertEqual(library["schema"], "template_fill_pptx_library.v1")
            self.assertEqual(library["slide_count"], 1)
            self.assertEqual(len(library["slides"][0]["slots"]), 2)

            plan = scaffold_plan(library)
            plan["slides"][0]["replacements"][0]["text"] = "New Title"
            errors, warnings = validate_plan(plan, library)
            self.assertEqual(errors, [])
            self.assertIsInstance(warnings, list)

            apply_plan(plan, output)

            filled = Presentation(str(output))
            texts = [
                shape.text
                for shape in filled.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
            ]
            self.assertIn("New Title", texts)
            self.assertIn(
                "Original body text with enough words to be treated as body content.",
                texts,
            )

    def test_validate_rejects_unknown_slot_when_library_is_available(self):
        from scripts.template_fill_pptx import analyze_pptx, scaffold_plan, validate_plan

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "template.pptx"
            make_template_pptx(source)

            library = analyze_pptx(source)
            plan = scaffold_plan(library)
            plan["slides"][0]["replacements"].append({"slot_id": "missing_slot", "text": "Nope"})

            errors, _warnings = validate_plan(plan, library)

            self.assertTrue(any("unknown slot_id" in error for error in errors))

    def test_apply_supports_reordered_and_reused_source_slides(self):
        from scripts.template_fill_pptx import analyze_pptx, apply_plan

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "template.pptx"
            output = tmp_path / "filled.pptx"
            make_two_slide_template_pptx(source)

            library = analyze_pptx(source)
            alpha_slot = library["slides"][0]["slots"][0]["slot_id"]
            beta_slot = library["slides"][1]["slots"][0]["slot_id"]
            plan = {
                "schema": "template_fill_pptx_plan.v1",
                "status": "draft",
                "source_pptx": str(source),
                "accepted_warnings": [],
                "slides": [
                    {
                        "source_slide": 2,
                        "layout_rationale": {"layout_pattern": "reuse-beta"},
                        "replacements": [{"slot_id": beta_slot, "text": "First output"}],
                        "table_edits": [],
                    },
                    {
                        "source_slide": 1,
                        "layout_rationale": {"layout_pattern": "alpha"},
                        "replacements": [{"slot_id": alpha_slot, "text": "Second output"}],
                        "table_edits": [],
                    },
                    {
                        "source_slide": 2,
                        "layout_rationale": {"layout_pattern": "reuse-beta-again"},
                        "replacements": [{"slot_id": beta_slot, "text": "Third output"}],
                        "table_edits": [],
                    },
                ],
            }

            apply_plan(plan, output)

            filled = Presentation(str(output))
            self.assertEqual(len(filled.slides), 3)
            ordered_titles = [
                next(
                    shape.text
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False) and shape.text
                )
                for slide in filled.slides
            ]
            self.assertEqual(ordered_titles, ["First output", "Second output", "Third output"])

    def test_analyze_and_apply_chart_edits_clone_chart_parts_per_output_slide(self):
        from scripts.template_fill_pptx import analyze_pptx, apply_plan

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "chart_template.pptx"
            output = tmp_path / "filled_chart.pptx"
            make_chart_template_pptx(source)

            library = analyze_pptx(source)
            chart = library["slides"][0]["charts"][0]
            self.assertEqual(chart["categories"], ["A", "B"])
            self.assertEqual(chart["series"][0]["name"], "Old Series")
            self.assertEqual(chart["series"][0]["values"], [1.0, 2.0])

            plan = {
                "schema": "template_fill_pptx_plan.v1",
                "status": "draft",
                "source_pptx": str(source),
                "accepted_warnings": [],
                "slides": [
                    {
                        "source_slide": 1,
                        "layout_rationale": {"layout_pattern": "chart"},
                        "replacements": [],
                        "table_edits": [],
                        "chart_edits": [
                            {
                                "chart_id": chart["chart_id"],
                                "categories": ["North", "South"],
                                "series": [{"name": "First Series", "values": [12, 18]}],
                            }
                        ],
                    },
                    {
                        "source_slide": 1,
                        "layout_rationale": {"layout_pattern": "chart-reuse"},
                        "replacements": [],
                        "table_edits": [],
                        "chart_edits": [
                            {
                                "chart_id": chart["chart_id"],
                                "categories": ["East", "West"],
                                "series": [{"name": "Second Series", "values": [7, 9]}],
                            }
                        ],
                    },
                ],
            }

            apply_plan(plan, output)

            parts = chart_parts_by_slide(output)
            self.assertEqual(len(parts), 2)
            self.assertNotEqual(parts[0], parts[1])
            workbook_parts = workbook_parts_by_chart(output, parts)
            self.assertEqual(len(workbook_parts), 2)
            self.assertNotEqual(workbook_parts[0], workbook_parts[1])
            with zipfile.ZipFile(output) as zf:
                first_chart = zf.read(parts[0]).decode("utf-8")
                second_chart = zf.read(parts[1]).decode("utf-8")
            self.assertIn("First Series", first_chart)
            self.assertIn("North", first_chart)
            self.assertIn(">12<", first_chart)
            self.assertIn("Second Series", second_chart)
            self.assertIn("West", second_chart)
            self.assertIn(">9<", second_chart)
            self.assertEqual(
                workbook_matrix(output, workbook_parts[0]),
                [[None, "First Series"], ["North", 12], ["South", 18]],
            )
            self.assertEqual(
                workbook_matrix(output, workbook_parts[1]),
                [[None, "Second Series"], ["East", 7], ["West", 9]],
            )

    def test_apply_deep_clones_structured_private_slide_dependencies(self):
        from scripts.template_fill_pptx import analyze_pptx, apply_plan

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "template.pptx"
            output = tmp_path / "filled_private_parts.pptx"
            make_template_pptx(source)
            inject_private_tag_part(source)

            library = analyze_pptx(source)
            slot = library["slides"][0]["slots"][0]["slot_id"]
            plan = {
                "schema": "template_fill_pptx_plan.v1",
                "status": "draft",
                "source_pptx": str(source),
                "accepted_warnings": [],
                "slides": [
                    {
                        "source_slide": 1,
                        "layout_rationale": {"layout_pattern": "private-a"},
                        "replacements": [{"slot_id": slot, "text": "Clone A"}],
                        "table_edits": [],
                        "chart_edits": [],
                    },
                    {
                        "source_slide": 1,
                        "layout_rationale": {"layout_pattern": "private-b"},
                        "replacements": [{"slot_id": slot, "text": "Clone B"}],
                        "table_edits": [],
                        "chart_edits": [],
                    },
                ],
            }

            apply_plan(plan, output)

            targets = private_tag_targets_by_slide(output)
            self.assertEqual(len(targets), 2)
            self.assertNotEqual(targets[0], targets[1])
            self.assertTrue(targets[0].startswith("ppt/tags/tag1_tf"))
            self.assertTrue(targets[1].startswith("ppt/tags/tag1_tf"))
            with zipfile.ZipFile(output) as zf:
                self.assertIn(targets[0], zf.namelist())
                self.assertIn(targets[1], zf.namelist())


if __name__ == "__main__":
    unittest.main()
