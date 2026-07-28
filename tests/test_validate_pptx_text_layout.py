import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


def save_deck(build):
    tmp = tempfile.TemporaryDirectory()
    path = Path(tmp.name) / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build(slide)
    prs.save(path)
    return tmp, path


class ValidatePptxTextLayoutTests(unittest.TestCase):
    def test_clean_text_layout_passes(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        tmp, path = save_deck(
            lambda slide: add_textbox(slide, 1.0, 1.0, 5.0, 1.0, "Short readable text", 24)
        )
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertEqual(report["status"], "pass", report["issues"])
        self.assertEqual(report["blocking_count"], 0)

    def test_long_text_in_small_box_is_blocking_overflow(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        long_text = " ".join(["capacity"] * 80)
        tmp, path = save_deck(
            lambda slide: add_textbox(slide, 0.7, 0.7, 1.4, 0.45, long_text, 22, word_wrap=True)
        )
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertEqual(report["status"], "fail")
        self.assertTrue(
            any(issue["code"] == "TEXT-OVERFLOW" for issue in report["issues"]),
            report["issues"],
        )

    def test_sentence_in_tall_centered_box_is_not_treated_as_small_label(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            box = add_textbox(
                slide,
                1.0,
                1.0,
                10.0,
                3.0,
                "Monitoring, mechanism, and population evidence support the conclusion.",
                18,
                word_wrap=True,
            )
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertNotIn(
            "TEXT-LABEL-TOO-LONG",
            {issue["code"] for issue in report["issues"]},
        )

    def test_no_wrap_text_uses_explicit_lines_for_overflow(self):
        from scripts.validate_pptx_text_layout import PT_PER_INCH, TextBox, _measure_lines_for_box

        text = "LongNativeAxisLabel\nSecondLongAxisLabel"
        nowrap = TextBox(
            1, 1, "nowrap", text, 1, 1, 0.35, 0.6, 0.3, 0.5, "Arial", 18, 1.2, wrap="none"
        )
        wrapped = TextBox(
            1, 2, "wrapped", text, 1, 1, 0.35, 0.6, 0.3, 0.5, "Arial", 18, 1.2, wrap="square"
        )
        width_pt = nowrap.usable_w * PT_PER_INCH

        self.assertEqual(_measure_lines_for_box(nowrap, width_pt), 2)
        self.assertGreater(_measure_lines_for_box(wrapped, width_pt), 2)

    def test_explicit_line_breaks_use_natural_line_height(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            add_textbox(slide, 1.0, 1.0, 2.0, 1.3, "第一行\n第二行\n第三行\n第四行", 18)

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertNotIn("TEXT-OVERFLOW", {issue["code"] for issue in report["issues"]})

    def test_highly_overlapping_text_boxes_are_blocking(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            add_textbox(slide, 1.0, 1.0, 3.0, 1.0, "First card body", 20)
            add_textbox(slide, 1.1, 1.05, 3.0, 1.0, "Second card body", 20)

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertEqual(report["status"], "fail")
        self.assertTrue(
            any(issue["code"] == "TEXT-OVERLAP" for issue in report["issues"]),
            report["issues"],
        )

    def test_raw_xml_fallback_extracts_native_text_shapes(self):
        from scripts.validate_pptx_text_layout import _iter_text_boxes_from_xml

        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "raw_xml_text.pptx"
            with zipfile.ZipFile(path, "w") as zf:
                zf.writestr(
                    "ppt/presentation.xml",
                    """<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>
</p:presentation>""",
                )
                zf.writestr(
                    "ppt/_rels/presentation.xml.rels",
                    """<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>""",
                )
                zf.writestr(
                    "ppt/slides/slide1.xml",
                    """<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld><p:spTree>
    <p:sp>
      <p:nvSpPr><p:cNvPr id="2" name="Native Text"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
      <p:spPr><a:xfrm><a:off x="914400" y="1828800"/><a:ext cx="2743200" cy="914400"/></a:xfrm></p:spPr>
      <p:txBody>
        <a:bodyPr lIns="0" rIns="0" tIns="0" bIns="0" wrap="none"/>
        <a:p><a:r><a:rPr sz="2400"><a:latin typeface="Arial"/></a:rPr><a:t>Native fallback text</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
  </p:spTree></p:cSld>
</p:sld>""",
                )

            boxes = _iter_text_boxes_from_xml(path)

        self.assertEqual(len(boxes), 1)
        self.assertEqual(boxes[0].text, "Native fallback text")
        self.assertEqual(boxes[0].name, "Native Text")
        self.assertAlmostEqual(boxes[0].x, 1.0)
        self.assertAlmostEqual(boxes[0].y, 2.0)
        self.assertLess(boxes[0].w, 3.0)
        self.assertGreater(boxes[0].w, 1.0)
        self.assertAlmostEqual(boxes[0].usable_w, 3.0)
        self.assertEqual(boxes[0].font_size_pt, 24.0)
        self.assertEqual(boxes[0].wrap, "none")

    def test_negative_drawingml_extent_is_blocking(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        tmp, path = save_deck(
            lambda slide: add_textbox(slide, 1.0, 1.0, 2.0, 0.5, "01-", 24)
        )
        with tmp:
            rewritten = path.with_suffix(".rewritten.pptx")
            with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(rewritten, "w") as dst:
                for info in src.infolist():
                    data = src.read(info.filename)
                    if info.filename == "ppt/slides/slide1.xml":
                        data = data.replace(b'<a:ext cx="1828800"', b'<a:ext cx="-1828800"', 1)
                    dst.writestr(info, data)
            rewritten.replace(path)

            report = validate_pptx_text_layout(path)

        self.assertEqual(report["status"], "fail")
        self.assertTrue(
            any(issue["code"] == "PPTX-INVALID-NEGATIVE-EXTENT" for issue in report["issues"]),
            report["issues"],
        )

    def test_control_text_must_be_vertically_centered_in_container(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            add_control_container(slide, 1.0, 1.0, 4.0, 0.7)
            add_textbox(slide, 1.2, 1.03, 2.2, 0.25, "Badge", 24)

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertEqual(report["status"], "fail")
        self.assertTrue(
            any(issue["code"] == "CONTROL-TEXT-VERTICAL-MISALIGN" for issue in report["issues"]),
            report["issues"],
        )

    def test_control_text_center_lock_passes(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            add_control_container(slide, 1.0, 1.0, 4.0, 0.7)
            label = add_textbox(slide, 1.0, 1.0, 4.0, 0.7, "Badge", 24)
            label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertEqual(report["status"], "pass", report["issues"])
        self.assertFalse(
            any(issue["code"] == "CONTROL-TEXT-VERTICAL-MISALIGN" for issue in report["issues"]),
            report["issues"],
        )

    def test_control_container_prefers_full_text_coverage_over_partial_background(self):
        from scripts.validate_pptx_text_layout import (
            ShapeBox,
            TextBox,
            _control_container_for_text,
        )

        label = TextBox(1, 1, "label", "Focus", 2.0, 2.0, 1.8, 0.3, 1.8, 0.3, "Arial", 24, 1.2)
        partial_background = ShapeBox(1, 2, "legacy", 1.9, 2.15, 2.2, 0.6, "#751497", "rect")
        component_control = ShapeBox(1, 3, "live", 1.7, 1.8, 2.4, 0.7, "#C00000", "rect")

        selected = _control_container_for_text(
            label,
            [partial_background, component_control],
            [label],
        )

        self.assertIsNotNone(selected)
        self.assertEqual(selected.shape_index, component_control.shape_index)

    def test_bottom_caption_is_not_treated_as_control_text(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            add_control_container(slide, 1.0, 1.0, 4.0, 2.2)
            add_textbox(slide, 1.2, 2.7, 3.6, 0.25, "Figure caption", 12)

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertNotIn("CONTROL-TEXT-VERTICAL-MISALIGN", {issue["code"] for issue in report["issues"]})

    def test_compound_card_text_is_not_treated_as_single_control(self):
        from scripts.validate_pptx_text_layout import validate_pptx_text_layout

        def build(slide):
            add_control_container(slide, 1.0, 1.0, 4.0, 0.95)
            add_textbox(slide, 1.2, 1.12, 2.0, 0.22, "Card title", 16)
            add_textbox(slide, 1.2, 1.45, 3.4, 0.25, "Card body copy", 12)

        tmp, path = save_deck(build)
        with tmp:
            report = validate_pptx_text_layout(path)

        self.assertNotIn("CONTROL-TEXT-VERTICAL-MISALIGN", {issue["code"] for issue in report["issues"]})

    def test_short_and_vertical_labels_are_not_treated_as_body_overflow(self):
        from scripts.validate_pptx_text_layout import TextBox, _is_short_label_like, _is_vertical_stack_label

        short = TextBox(1, 1, "label", "TaOX", 1, 1, 0.4, 0.3, 0.35, 0.28, "Arial", 20, 1.2)
        vertical = TextBox(1, 2, "vertical", "B\ni\nn", 1, 1, 0.14, 0.7, 1.0, 0.5, "Arial", 14, 1.2)
        body = TextBox(1, 3, "body", "Long label text", 1, 1, 0.4, 0.3, 0.35, 0.28, "Arial", 20, 1.2)

        self.assertTrue(_is_short_label_like(short))
        self.assertTrue(_is_vertical_stack_label(vertical))
        self.assertFalse(_is_short_label_like(body))

    def test_final_pptx_qa_docs_include_text_layout_gate(self):
        docs = [
            ROOT / "SKILL.md",
            ROOT / "references" / "workflow-create.md",
            ROOT / "scripts" / "docs" / "svg-pipeline.md",
        ]

        for path in docs:
            text = path.read_text(encoding="utf-8")
            self.assertIn("scripts/validate_pptx_text_layout.py", text, str(path))
            self.assertIn("text_layout_report.json", text, str(path))


def add_textbox(slide, x, y, w, h, text, size_pt, word_wrap=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if word_wrap is not None:
        box.text_frame.word_wrap = word_wrap
    box.text_frame.clear()
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = "Arial"
    return box


def add_control_container(slide, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(25, 85, 145)
    shape.line.fill.background()
    return shape


if __name__ == "__main__":
    unittest.main()
