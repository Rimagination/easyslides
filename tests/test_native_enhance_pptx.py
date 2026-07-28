import json
import tempfile
import unittest
import wave
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def make_source_deck(path: Path, slide_count: int = 1) -> None:
    prs = Presentation()
    for idx in range(slide_count):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(0.7))
        box.text = f"Slide {idx + 1}"
    prs.save(path)


def make_silent_wav(path: Path, duration_seconds: float = 0.15) -> None:
    frame_rate = 8000
    frame_count = int(frame_rate * duration_seconds)
    with wave.open(str(path), "wb") as audio:
        audio.setnchannels(1)
        audio.setsampwidth(2)
        audio.setframerate(frame_rate)
        audio.writeframes(b"\x00\x00" * frame_count)


class NativeEnhancePptxTests(unittest.TestCase):
    def test_init_project_writes_contract_artifacts(self):
        from scripts.native_enhance_pptx import init_project, plan, validate

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "source.pptx"
            make_source_deck(source)

            project = init_project(source, name="demo deck", base_dir=tmp_path / "projects")

            self.assertEqual(project.name, "demo_deck")
            for rel in ("sources", "analysis", "notes", "audio", "exports", "validation"):
                self.assertTrue((project / rel).is_dir(), rel)

            project_json = json.loads((project / "project.json").read_text(encoding="utf-8"))
            self.assertEqual(project_json["schema"], "native_pptx_enhancement_project.v1")
            self.assertEqual(project_json["slide_count"], 1)
            self.assertTrue((project / project_json["source_pptx"]).is_file())

            slide_index = json.loads((project / "analysis" / "slide_index.json").read_text(encoding="utf-8"))
            self.assertEqual(slide_index["slides"][0]["slide"], 1)
            self.assertIsInstance(slide_index["slides"][0]["partname"], str)

            refreshed = plan(project)
            self.assertEqual(refreshed["schema"], "native_pptx_enhancement_plan.v1")

            errors, warnings = validate(project)
            self.assertEqual(errors, [])
            self.assertEqual(warnings, [])

    def test_validate_warns_when_note_count_does_not_match_slide_count(self):
        from scripts.native_enhance_pptx import init_project, validate

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "source.pptx"
            make_source_deck(source, slide_count=2)
            project = init_project(source, name="notes_demo", base_dir=tmp_path / "projects")
            (project / "notes" / "slide_001.md").write_text("Speaker note", encoding="utf-8")

            errors, warnings = validate(project)

            self.assertEqual(errors, [])
            self.assertTrue(any("note count 1 != slide count 2" in warning for warning in warnings))

    def test_apply_adds_notes_and_transition_without_changing_visible_text(self):
        from scripts.native_enhance_pptx import apply, init_project

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "source.pptx"
            make_source_deck(source, slide_count=2)
            project = init_project(source, name="apply_demo", base_dir=tmp_path / "projects")
            (project / "notes" / "001.md").write_text("# Narration\n\nKeep this note.", encoding="utf-8")

            plan_path = project / "analysis" / "enhancement_plan.json"
            plan = json.loads(plan_path.read_text(encoding="utf-8"))
            plan["status"] = "confirmed"
            plan["modules"]["transitions"]["enabled"] = True
            plan["modules"]["transitions"]["apply_without_audio"] = True
            plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

            output = apply(project)

            prs = Presentation(str(output))
            visible_text = [
                shape.text
                for shape in prs.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
            ]
            self.assertIn("Slide 1", visible_text)

            with zipfile.ZipFile(output) as zf:
                names = set(zf.namelist())
                self.assertIn("ppt/notesSlides/notesSlide1.xml", names)
                self.assertIn("ppt/notesSlides/_rels/notesSlide1.xml.rels", names)
                self.assertIn("ppt/notesMasters/notesMaster1.xml", names)
                notes_xml = zf.read("ppt/notesSlides/notesSlide1.xml").decode("utf-8")
                slide_rels = zf.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
                slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
                content_types = zf.read("[Content_Types].xml").decode("utf-8")

            self.assertIn("Keep this note.", notes_xml)
            self.assertIn("relationships/notesSlide", slide_rels)
            self.assertIn("<p:transition", slide_xml)
            self.assertIn("/ppt/notesSlides/notesSlide1.xml", content_types)

            report = json.loads((project / "validation" / "apply_report.json").read_text(encoding="utf-8"))
            self.assertEqual(report["notes_applied"], 1)
            self.assertEqual(report["transition_only_slides"], 2)

    def test_apply_embeds_audio_and_sets_recorded_timing(self):
        from scripts.native_enhance_pptx import apply, init_project

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "source.pptx"
            make_source_deck(source)
            project = init_project(source, name="audio_demo", base_dir=tmp_path / "projects")
            make_silent_wav(project / "audio" / "001.wav")

            plan_path = project / "analysis" / "enhancement_plan.json"
            plan = json.loads(plan_path.read_text(encoding="utf-8"))
            plan["status"] = "confirmed"
            plan["modules"]["audio"]["enabled"] = True
            plan["modules"]["timings"]["enabled"] = True
            plan["modules"]["transitions"]["enabled"] = True
            plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

            output = apply(project)

            with zipfile.ZipFile(output) as zf:
                names = set(zf.namelist())
                self.assertIn("ppt/media/native_enhance_audio_001.wav", names)
                self.assertIn("ppt/media/native_enhance_audio_poster.png", names)
                slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
                slide_rels = zf.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
                content_types = zf.read("[Content_Types].xml").decode("utf-8")

            self.assertIn("<a:audioFile", slide_xml)
            self.assertIn("advTm=", slide_xml)
            self.assertIn("relationships/audio", slide_rels)
            self.assertIn("relationships/media", slide_rels)
            self.assertIn('Extension="wav"', content_types)

            report = json.loads((project / "validation" / "apply_report.json").read_text(encoding="utf-8"))
            self.assertEqual(report["audio_embedded"], 1)


if __name__ == "__main__":
    unittest.main()
