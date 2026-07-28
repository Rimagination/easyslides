import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]


def run_cli(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, *args],
        cwd=ROOT,
        text=True,
        capture_output=True,
        encoding="utf-8",
        errors="replace",
    )


def make_simple_pptx(path: Path) -> None:
    prs = Presentation()
    for title_text in ("First Slide", "Second Slide"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(0.6))
        title.text = title_text
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(7), Inches(1.0))
        body.text = f"Body for {title_text}"
    prs.save(path)


def visible_texts(path: Path) -> list[list[str]]:
    prs = Presentation(str(path))
    return [
        [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        ]
        for slide in prs.slides
    ]


class BeautifyPptxTests(unittest.TestCase):
    def test_inspect_writes_gap_aware_report_and_manifest(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "source.pptx"
            report_dir = tmp_path / "beautify"
            make_simple_pptx(source)

            result = run_cli("scripts/beautify_pptx.py", "inspect", str(source), "--out", str(report_dir))

            self.assertEqual(result.returncode, 0, result.stdout + result.stderr)
            report = json.loads((report_dir / "beautify_report.json").read_text(encoding="utf-8"))
            manifest = json.loads((report_dir / "workflow_manifest.json").read_text(encoding="utf-8"))
            self.assertEqual(report["schema_version"], "easyslides.beautify_report.v1")
            self.assertEqual(report["route"], "beautify-pptx")
            self.assertEqual(report["slide_count"], 2)
            self.assertEqual(report["status"], "inspected")
            self.assertEqual(report["text_by_slide"][0], ["First Slide", "Body for First Slide"])
            self.assertEqual(manifest["current"]["route"], "beautify-pptx")
            self.assertEqual(manifest["current"]["stage"], "inspect")

    def test_apply_theme_patch_preserves_slide_count_and_visible_text(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "source.pptx"
            output = tmp_path / "beautified.pptx"
            report_dir = tmp_path / "beautify"
            make_simple_pptx(source)

            before_text = visible_texts(source)
            result = run_cli(
                "scripts/beautify_pptx.py",
                "apply",
                str(source),
                "-o",
                str(output),
                "--report-dir",
                str(report_dir),
                "--primary",
                "#123456",
                "--accent",
                "#ABCDEF",
            )

            self.assertEqual(result.returncode, 0, result.stdout + result.stderr)
            self.assertEqual(visible_texts(output), before_text)
            self.assertEqual(len(Presentation(str(output)).slides), 2)
            with zipfile.ZipFile(output) as zf:
                theme_xml = zf.read("ppt/theme/theme1.xml").decode("utf-8")
            self.assertIn("123456", theme_xml)
            self.assertIn("ABCDEF", theme_xml)
            report = json.loads((report_dir / "beautify_report.json").read_text(encoding="utf-8"))
            self.assertEqual(report["status"], "applied")
            self.assertIn("theme_color_patch", report["actions"])


if __name__ == "__main__":
    unittest.main()
