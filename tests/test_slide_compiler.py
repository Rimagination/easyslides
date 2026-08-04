from __future__ import annotations

import json
import tempfile
import unittest
from copy import deepcopy
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
NSFC = ROOT / "templates" / "layouts" / "nsfc_defense"
ACADEMIC_GENERAL = ROOT / "templates" / "layouts" / "academic_general"
FIXTURE = ROOT / "tests" / "fixtures" / "nsfc_composable_deck_plan.json"


class SlideCompilerTests(unittest.TestCase):
    def _compile_fixture(self) -> dict:
        from scripts.slide_compiler import compile_slides
        from scripts.template_compiler import compile_template

        deck_plan = json.loads(FIXTURE.read_text(encoding="utf-8"))
        template_ir = compile_template(NSFC)["template_ir"]
        return compile_slides(deck_plan, template_ir)

    def test_compiles_source_derived_component_scene_with_semantic_slots(self) -> None:
        slide_ir = self._compile_fixture()
        slide = slide_ir["slides"][0]
        component_layers = [layer for layer in slide["layers"] if layer["layer_type"] == "component"]

        self.assertEqual(slide_ir["schema_version"], "easyslides.slide_ir.v1")
        self.assertEqual(slide["shell_id"], "content")
        self.assertEqual(slide["body_variant_id"], "need_relationship_evidence")
        self.assertEqual(slide["shell_payload"]["PAGE_NUMBER"], "01")
        self.assertIn("\n", slide["shell_payload"]["KEY_MESSAGE"])
        self.assertEqual(
            {layer["instance_id"] for layer in component_layers},
            {"statement", "relationship", "track_one", "track_two", "callout"},
        )
        statement = next(layer for layer in component_layers if layer["instance_id"] == "statement")
        self.assertIn("STATEMENT", statement["payload"])
        self.assertIn("LEFT_LABEL", next(layer for layer in component_layers if layer["instance_id"] == "relationship")["payload"])
        self.assertNotIn("PRIMARY_BODY", component_layers[0]["payload"])
        for layer in component_layers:
            self.assertGreater(layer["frame"]["width"], 0)
            self.assertGreater(layer["frame"]["height"], 0)

    def test_template_owned_page_number_is_generalized_beyond_nsfc(self) -> None:
        from scripts.slide_compiler import compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(ACADEMIC_GENERAL)["template_ir"]
        slide_ir = compile_slides(
            {
                "template_id": "academic_general",
                "slides": [
                    {
                        "page": "P01",
                        "role": "content",
                        "body_variant_id": "comparison_synthesis",
                        "shell_payload": {
                            "SECTION_NUM": "01",
                            "PAGE_TITLE": "Two methods answer different questions",
                            "LOGO": "LAB",
                            "KEY_MESSAGE": "Method choice follows the evidence needed.",
                            "SOURCE": "Source: study protocol",
                            "SECTION_NAME": "Methods",
                        },
                        "slot_payload": {
                            "LEFT_TITLE": "Method A",
                            "LEFT_BODY": "Measures broad patterns across the cohort.",
                            "RIGHT_TITLE": "Method B",
                            "RIGHT_BODY": "Resolves the mechanism within each case.",
                            "SYNTHESIS": "Use both methods when the claim needs scale and mechanism.",
                        },
                    },
                    {
                        "page": "P02",
                        "role": "content",
                        "body_variant_id": "comparison_synthesis",
                        "shell_payload": {
                            "SECTION_NUM": "01",
                            "PAGE_TITLE": "The combined design reduces ambiguity",
                            "LOGO": "LAB",
                            "KEY_MESSAGE": "Complementary evidence makes the conclusion testable.",
                            "SOURCE": "Source: analysis plan",
                            "SECTION_NAME": "Methods",
                        },
                        "slot_payload": {
                            "LEFT_TITLE": "Breadth",
                            "LEFT_BODY": "Detects reproducible population-level signals.",
                            "RIGHT_TITLE": "Depth",
                            "RIGHT_BODY": "Tests the causal explanation for those signals.",
                            "SYNTHESIS": "The joint design turns association into an auditable argument.",
                        },
                    },
                ],
            },
            template_ir,
        )

        self.assertEqual(
            [slide["shell_payload"]["PAGE_NUM"] for slide in slide_ir["slides"]],
            ["01", "02"],
        )

    def test_academic_general_open_composition_accepts_only_local_registered_components(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(ACADEMIC_GENERAL)["template_ir"]
        plan = {
            "template_id": "academic_general",
            "slides": [
                {
                    "page": "P01",
                    "role": "content",
                    "body_variant_id": "open_component_composition",
                    "shell_payload": {
                        "SECTION_NUM": "02",
                        "PAGE_TITLE": "A selected metric clarifies the decision",
                        "LOGO": "LAB",
                        "KEY_MESSAGE": "Use a local metric component when a reviewed variant does not fit.",
                        "SOURCE": "Source: decision memo",
                        "SECTION_NAME": "Evidence",
                    },
                    "body_components": [
                        {
                            "component_id": "metric_tile",
                            "instance_id": "primary_metric",
                            "frame": {"x": 456, "y": 260, "width": 368, "height": 190},
                            "slot_payload": {"VALUE": "86%", "LABEL": "Evidence coverage"},
                        }
                    ],
                }
            ],
        }

        slide_ir = compile_slides(plan, template_ir)
        layers = slide_ir["slides"][0]["layers"]
        metric = next(layer for layer in layers if layer.get("instance_id") == "primary_metric")

        self.assertEqual(metric["asset_id"], "component/academic_general/metric_tile")
        self.assertEqual(metric["composition_source"], "explicit_body_components")

        plan["slides"][0]["body_components"][0]["component_id"] = "component/nsfc_defense/metric_tile"
        with self.assertRaisesRegex(SlideCompileError, "unregistered component"):
            compile_slides(plan, template_ir)

    def test_cover_uses_shared_title_and_date_payload_names_without_duplicate_dates(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        payload = {
            "TITLE": "Memristive neural computing circuits",
            "PROJECT_TYPE": "National Natural Science Foundation project",
            "SUBTITLE": "Final defense presentation",
            "AFFILIATION": "Research University",
            "PRESENTER": "Liang",
            "DATE": "2026-07-28",
        }
        slide_ir = compile_slides(
            {"template_id": "nsfc_defense", "slides": [{"page": "C01", "role": "cover", "shell_payload": payload}]},
            template_ir,
        )

        self.assertEqual(slide_ir["slides"][0]["shell_id"], "cover")
        self.assertEqual(slide_ir["slides"][0]["shell_payload"], payload)
        with self.assertRaises(SlideCompileError):
            compile_slides(
                {
                    "template_id": "nsfc_defense",
                    "slides": [{"page": "C01", "role": "cover", "shell_payload": {"DATE_02": "duplicate"}}],
                },
                template_ir,
            )

    def test_nsfc_ending_defaults_to_critique_and_rejects_listening_copy(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        base_payload = {
            "AFFILIATION": "东海环健",
            "PRESENTER": "林澈",
            "DATE": "2026.07",
        }
        slide_ir = compile_slides(
            {"template_id": "nsfc_defense", "slides": [{"page": "E01", "role": "ending", "shell_payload": base_payload}]},
            template_ir,
        )
        self.assertEqual(slide_ir["slides"][0]["shell_payload"]["CLOSING_TITLE"], "敬请批评指正")

        with self.assertRaisesRegex(SlideCompileError, "may not contain '聆听'"):
            compile_slides(
                {
                    "template_id": "nsfc_defense",
                    "slides": [
                        {
                            "page": "E02",
                            "role": "ending",
                            "shell_payload": {**base_payload, "CLOSING_TITLE": "感谢聆听"},
                        }
                    ],
                },
                template_ir,
            )

        with self.assertRaisesRegex(SlideCompileError, "one closing line only"):
            compile_slides(
                {
                    "template_id": "nsfc_defense",
                    "slides": [
                        {
                            "page": "E03",
                            "role": "ending",
                            "shell_payload": {**base_payload, "CLOSING_SUBTITLE": "副标题"},
                        }
                    ],
                },
                template_ir,
            )

    def test_nsfc_content_title_requires_one_visual_line(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        newline_plan = json.loads(FIXTURE.read_text(encoding="utf-8"))
        newline_plan["slides"][0]["shell_payload"]["PAGE_TITLE"] = "Risk\nchain"
        with self.assertRaisesRegex(SlideCompileError, "must be one visual line"):
            compile_slides(newline_plan, template_ir)

        oversized_plan = json.loads(FIXTURE.read_text(encoding="utf-8"))
        oversized_plan["slides"][0]["shell_payload"]["PAGE_TITLE"] = "\u4e2d" * 11
        with self.assertRaisesRegex(SlideCompileError, "single-line budget"):
            compile_slides(oversized_plan, template_ir)

    def test_renders_source_derived_scene_to_svg_and_native_pptx(self) -> None:
        from scripts.slide_compiler import render_slide_ir_to_pptx, render_slide_ir_to_svg

        slide_ir = self._compile_fixture()
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            svg_report = render_slide_ir_to_svg(slide_ir, root / "svg")
            svg = Path(svg_report["svg_files"][0])
            svg_root = ET.parse(svg).getroot()
            title = next(node for node in svg_root.iter() if node.attrib.get("data-slot-id") == "PAGE_TITLE")
            key_message = next(
                node for node in svg_root.iter() if node.attrib.get("data-easyslides-generated") == "square_bullets"
            )
            instances = [
                node.attrib["data-easyslides-instance"]
                for node in svg_root.iter()
                if node.attrib.get("data-easyslides-instance")
            ]
            pptx_path = root / "compiled.pptx"
            pptx_report = render_slide_ir_to_pptx(slide_ir, pptx_path, svg_output_dir=root / "pptx-svg")
            presentation = Presentation(pptx_path)
            caption = next(
                layer["payload"]["CAPTION"]
                for layer in slide_ir["slides"][0]["layers"]
                if layer.get("instance_id") == "track_one"
            )

            def iter_shapes(shapes):
                for shape in shapes:
                    yield shape
                    if hasattr(shape, "shapes"):
                        yield from iter_shapes(shape.shapes)

            caption_shape = next(
                shape
                for shape in iter_shapes(presentation.slides[0].shapes)
                if getattr(shape, "has_text_frame", False) and shape.text.strip() == caption
            )

            self.assertEqual(svg_report["status"], "pass")
            self.assertEqual(set(instances), {"statement", "relationship", "track_one", "track_two", "callout"})
            self.assertEqual(title.attrib.get("data-pptx-no-wrap"), "true")
            self.assertEqual(len([node for node in title if node.tag.endswith("tspan")]), 1)
            self.assertEqual(len([node for node in key_message if node.tag.endswith("rect")]), 2)
            self.assertIn('data-easyslides-instance="track_one"', svg.read_text(encoding="utf-8"))
            self.assertFalse(
                any(node.attrib.get("data-easyslides-clear-region") for node in svg_root.iter()),
                "clear_region must remain a non-rendered layout constraint",
            )
            self.assertEqual(pptx_report["status"], "pass")
            self.assertEqual(pptx_report["native_component_bounds"]["status"], "pass")
            self.assertEqual(pptx_report["native_component_bounds"]["checked_component_count"], 5)
            self.assertEqual(len(presentation.slides), 1)
            self.assertGreater(caption_shape.width, Inches(1))
            self.assertGreater(caption_shape.height, Inches(0.2))
            self.assertLess(caption_shape.height, Inches(0.8))
            self.assertEqual(caption_shape.text_frame.vertical_anchor, MSO_ANCHOR.MIDDLE)

    def test_all_nsfc_component_scenes_stay_inside_native_pptx_frames(self) -> None:
        from scripts.render_nsfc_component_first_gallery import build_plan
        from scripts.slide_compiler import compile_slides, render_slide_ir_to_pptx, render_slide_ir_to_svg
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        slide_ir = compile_slides(build_plan(template_ir), template_ir)
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            svg_report = render_slide_ir_to_svg(slide_ir, root / "svg")
            report = render_slide_ir_to_pptx(
                slide_ir,
                root / "all_component_scenes.pptx",
                svg_output_dir=root / "pptx-svg",
            )
            for svg_path in svg_report["svg_files"]:
                svg_root = ET.parse(svg_path).getroot()
                ids = [node.attrib["id"] for node in svg_root.iter() if node.attrib.get("id")]
                self.assertEqual(len(ids), len(set(ids)), svg_path)

            for index, title in ((4, "技术指标看板"), (6, "方法比较边界")):
                svg_root = ET.parse(root / "svg" / f"{index:02d}_content.svg").getroot()
                title_node = next(node for node in svg_root.iter() if node.attrib.get("data-slot-id") == "PAGE_TITLE")
                header = next(
                    node
                    for node in svg_root.iter()
                    if node.tag.rsplit("}", 1)[-1] == "rect"
                    and node.attrib.get("x") == "0"
                    and node.attrib.get("y") == "0"
                    and node.attrib.get("width") == "1280"
                )
                self.assertEqual("".join(title_node.itertext()).strip(), title)
                self.assertEqual(header.attrib.get("fill"), "url(#ggrad2)")
                self.assertEqual(
                    sum(1 for node in svg_root.iter() if node.attrib.get("id") == "ggrad2"),
                    1,
                )

        self.assertEqual(report["status"], "pass")
        self.assertEqual(report["native_component_bounds"]["status"], "pass")
        expected_component_count = sum(
            1
            for slide in slide_ir["slides"]
            for layer in slide["layers"]
            if layer.get("layer_type") == "component"
        )
        self.assertEqual(
            report["native_component_bounds"]["checked_component_count"],
            expected_component_count,
        )

    def test_narrow_nsfc_labels_use_deterministic_balanced_cjk_stack_wrapping(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides, render_slide_ir_to_svg
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        plan = {
            "template_id": "nsfc_defense",
            "slides": [
                {
                    "page": "P01",
                    "role": "content",
                    "section": "01",
                    "story_role": "bottleneck_chain",
                    "body_variant_id": "evidence_chain",
                    "shell_payload": {"PAGE_TITLE": "关键瓶颈链条", "KEY_MESSAGE": "用可验证证据串联关键瓶颈"},
                    "slot_payload": {
                        "STATEMENT": "形成可验证的研究判断",
                        "CHAIN_TAG": "关键证据",
                        **{
                            f"EVIDENCE_{index}_{field}": "关键证据" if field == "LABEL" else str(ROOT / "templates" / "reference" / "template_asset_sources" / "nsfc_defense_distilled" / "assets" / "image10.jpg")
                            for index in range(1, 5)
                            for field in ("IMAGE", "LABEL")
                        },
                    },
                }
            ],
        }
        slide_ir = compile_slides(plan, template_ir)
        tag_layer = next(layer for layer in slide_ir["slides"][0]["layers"] if layer.get("instance_id") == "key_tag")
        self.assertEqual(tag_layer["payload"]["TAG"], "关键\n证据")

        with tempfile.TemporaryDirectory() as tmp:
            report = render_slide_ir_to_svg(slide_ir, Path(tmp) / "svg")
            root = ET.parse(report["svg_files"][0]).getroot()
            tag = next(node for node in root.iter() if node.attrib.get("data-slot-id") == "TAG")
        self.assertEqual(["".join(tspan.itertext()) for tspan in tag], ["关键", "证据"])
        self.assertEqual(tag.attrib.get("data-pptx-no-wrap"), "true")

        overflow = deepcopy(plan)
        overflow["slides"][0]["slot_payload"]["CHAIN_TAG"] = "关键瓶颈链条超长"
        with self.assertRaisesRegex(SlideCompileError, "stacked-label budget"):
            compile_slides(overflow, template_ir)

    def test_nsfc_key_message_is_required_distinct_and_limited_to_two_lines(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        missing = json.loads(FIXTURE.read_text(encoding="utf-8"))
        missing["slides"][0]["shell_payload"].pop("KEY_MESSAGE")
        with self.assertRaisesRegex(SlideCompileError, "requires KEY_MESSAGE"):
            compile_slides(missing, template_ir)

        oversized = json.loads(FIXTURE.read_text(encoding="utf-8"))
        oversized["slides"][0]["shell_payload"]["KEY_MESSAGE"] = "第一条\n第二条\n第三条"
        with self.assertRaisesRegex(SlideCompileError, "one or two"):
            compile_slides(oversized, template_ir)

        repeated = json.loads(FIXTURE.read_text(encoding="utf-8"))
        repeated["slides"][0]["shell_payload"]["KEY_MESSAGE"] = repeated["slides"][0]["slot_payload"]["STATEMENT"]
        with self.assertRaisesRegex(SlideCompileError, "repeats body slot STATEMENT"):
            compile_slides(repeated, template_ir)

    def test_source_guided_content_rejects_direct_body_components(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        with self.assertRaisesRegex(SlideCompileError, "forbids direct body_components"):
            compile_slides(
                {
                    "template_id": "nsfc_defense",
                    "slides": [
                        {
                            "page": "P01",
                            "role": "content",
                            "section": "01",
                            "story_role": "national_need_evidence",
                            "body_variant_id": "need_relationship_evidence",
                            "shell_payload": {"PAGE_TITLE": "National need"},
                            "slot_payload": {},
                            "body_components": [{"component_id": "statement_panel", "frame": {"x": 92, "y": 170, "width": 510, "height": 360}, "slot_payload": {}}],
                        }
                    ],
                },
                template_ir,
            )

    def test_source_guided_variant_component_must_stay_inside_body_canvas(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = deepcopy(compile_template(NSFC)["template_ir"])
        variant = next(
            row
            for row in template_ir["body_variants"]
            if row["variant_id"] == "need_relationship_evidence"
        )
        statement_region = next(
            row for row in variant["regions"] if row["region_id"] == "statement"
        )
        statement_region["frame"]["x"] = 1.0
        deck_plan = json.loads(FIXTURE.read_text(encoding="utf-8"))

        with self.assertRaisesRegex(SlideCompileError, "falls outside the content body_canvas"):
            compile_slides(deck_plan, template_ir)

    def test_source_guided_content_requires_matching_section_and_story_role(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        with self.assertRaisesRegex(SlideCompileError, "does not permit story_role"):
            compile_slides(
                {
                    "template_id": "nsfc_defense",
                    "slides": [
                        {
                            "page": "P01",
                            "role": "content",
                            "section": "02",
                            "story_role": "method_comparison",
                            "body_variant_id": "need_relationship_evidence",
                            "shell_payload": {"PAGE_TITLE": "Method contrast"},
                            "slot_payload": {},
                        }
                    ],
                },
                template_ir,
            )

    def test_ambiguous_content_requires_an_explicit_choice(self) -> None:
        from scripts.slide_compiler import SlideCompileError, compile_slides
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        with self.assertRaisesRegex(SlideCompileError, "ambiguous"):
            compile_slides(
                {
                    "template_id": "nsfc_defense",
                    "slides": [{"page": "P01", "role": "content", "shell_payload": {"PAGE_TITLE": "Undefined content"}, "slot_payload": {}}],
                },
                template_ir,
            )

    def test_nsfc_full_grant_scenario_requires_complete_bound_story(self) -> None:
        from scripts.slide_compiler import SlideCompileError, validate_scenario_contract
        from scripts.template_compiler import compile_template

        template_ir = compile_template(NSFC)["template_ir"]
        profile = template_ir["story_structure"]["grant_cn_profile"]
        bindings = {
            row["grant_role"]: row
            for row in profile["variant_bindings"]
        }
        slides = []
        for index, grant_role in enumerate(profile["full_deck_roles"], start=1):
            role = (
                grant_role
                if grant_role in {"cover", "toc", "ending"}
                else "chapter"
                if grant_role.startswith("chapter_")
                else "content"
            )
            slide = {"page": f"P{index:02d}", "role": role, "grant_role": grant_role}
            binding = bindings.get(grant_role)
            if binding:
                slide.update(
                    {
                        "section": binding["section"],
                        "story_role": binding["story_role"],
                        "body_variant_id": binding["body_variant_id"],
                    }
                )
            slides.append(slide)
        plan = {
            "template_id": "nsfc_defense",
            "scenario_id": "nsfc_grant_cn",
            "scenario_mode": "full",
            "slides": slides,
        }

        audit = validate_scenario_contract(plan, template_ir)
        self.assertEqual(audit["status"], "pass")
        self.assertEqual(len(audit["declared_roles"]), 14)
        self.assertEqual(audit["optional_roles"], ["toc"])

        without_toc = deepcopy(plan)
        without_toc["slides"] = [
            slide for slide in without_toc["slides"] if slide["grant_role"] != "toc"
        ]
        self.assertEqual(validate_scenario_contract(without_toc, template_ir)["status"], "pass")

        missing = deepcopy(plan)
        missing["slides"] = missing["slides"][:-1]
        with self.assertRaisesRegex(SlideCompileError, "missing grant_role"):
            validate_scenario_contract(missing, template_ir)

        wrong_variant = deepcopy(plan)
        wrong_variant["slides"][3]["body_variant_id"] = "comparison_evidence"
        with self.assertRaisesRegex(SlideCompileError, "requires body_variant_id"):
            validate_scenario_contract(wrong_variant, template_ir)


if __name__ == "__main__":
    unittest.main()
